from __future__ import annotations

from datetime import date, datetime
import base64
import html
import hashlib
import io
from pathlib import Path
import re
import requests
import subprocess
import shutil
import tempfile
import uuid
import unicodedata

import pandas as pd
import altair as alt
import streamlit as st
import streamlit.components.v1 as components
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pypdf import PdfReader
import fitz
import pytesseract
from PIL import Image
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image as RLImage, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from data import MUNICIPIOS_JALISCO
from db import access_profile, client_with_token, configured, download_project_images, public_client, public_key, register_access, safe_name, upload_files, valid_official_email
from exports import build_docx, build_pdf

MODULE_PROJECTS = "Programas / Proyectos"
MODULE_BOARD = "Junta de Gobierno"
MODULE_COMMITTEES = "Comités"
MODULE_OFFICIAL_LETTERS = "Oficios Dirección General"
ALL_MODULES = [MODULE_PROJECTS, MODULE_BOARD, MODULE_COMMITTEES, MODULE_OFFICIAL_LETTERS]
USER_DIRECTIONS = ["Dirección General", "Dirección de Administración", "Dirección Jurídica", "Dirección de Operaciones", "Dirección de Planeación", "Órgano Interno de Control"]
PROJECT_DIRECTIONS = ["Dirección de Operaciones", "Dirección de Proyectos"]
MASTER_ADMIN_EMAIL = "yani.limberopulos@jalisco.gob.mx"


def user_can(module: str) -> bool:
    user = st.session_state.get("user", {})
    return user.get("rol") == "administrador" or module in (user.get("modulos") or [])


def is_master_admin() -> bool:
    """La autoridad destructiva se limita a una identidad explícita."""
    return str(st.session_state.get("user", {}).get("email") or "").strip().lower() == MASTER_ADMIN_EMAIL


def _remove_storage_paths(client, paths: list[str]) -> None:
    clean = list(dict.fromkeys(str(path) for path in paths if path))
    for start in range(0, len(clean), 100):
        client.storage.from_("expedientes").remove(clean[start:start + 100])


def delete_project_master(client, project_id: str) -> None:
    documents = client.table("documentos").select("ruta_storage").eq("proyecto_id", project_id).execute().data or []
    _remove_storage_paths(client, [row.get("ruta_storage") for row in documents])
    client.table("proyectos").delete().eq("id", project_id).execute()


def delete_board_session_master(client, session_id: str) -> None:
    documents = client.table("documentos_sesion_junta").select("ruta_storage").eq("sesion_id", session_id).execute().data or []
    agreements = client.table("acuerdos_junta").select("id").eq("sesion_id", session_id).execute().data or []
    agreement_ids = [row["id"] for row in agreements]
    agreement_files = []
    if agreement_ids:
        agreement_files = client.table("archivos_acuerdo").select("ruta_storage").in_("acuerdo_id", agreement_ids).execute().data or []
    _remove_storage_paths(client, [row.get("ruta_storage") for row in documents + agreement_files])
    client.table("sesiones_junta").delete().eq("id", session_id).execute()


def delete_committee_session_master(client, session_id: str) -> None:
    documents = client.table("documentos_sesion_comite").select("ruta_storage").eq("sesion_id", session_id).execute().data or []
    agreements = client.table("acuerdos_comite").select("id").eq("sesion_id", session_id).execute().data or []
    agreement_ids = [row["id"] for row in agreements]
    followup_files = []
    if agreement_ids:
        followup_files = (client.table("archivos_acuerdo_comite").select("ruta_storage")
                          .in_("acuerdo_id", agreement_ids).execute().data or [])
    _remove_storage_paths(client, [row.get("ruta_storage") for row in documents + followup_files])
    client.table("sesiones_comite").delete().eq("id", session_id).execute()


def master_delete_control(label: str, object_id: str, key: str, delete_action) -> None:
    """Control destructivo con confirmación escrita, visible sólo al administrador maestro."""
    if not is_master_admin():
        return
    with st.container(key=f"master_delete_panel_{key}"):
        with st.expander(f"🟧 Administración maestra · Eliminar {label}"):
            st.warning("Esta acción es definitiva y eliminará también acuerdos, seguimiento y documentos relacionados.")
            confirmation = st.text_input(
                "Para confirmar, escribe ELIMINAR", key=f"master_delete_text_{key}",
                placeholder="ELIMINAR",
            )
            if st.button(
                f"Eliminar definitivamente {label}", key=f"master_delete_button_{key}",
                type="primary", use_container_width=True, disabled=confirmation.strip() != "ELIMINAR",
            ):
                try:
                    delete_action()
                    st.success(f"{label.capitalize()} eliminado definitivamente.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"No fue posible eliminar {label}: {exc}")



def _reset_agreement_documents(client, agreement_id: str, kind: str) -> None:
    """Elimina opcionalmente documentos del expediente del acuerdo y sus archivos físicos."""
    if kind == "board":
        table_name = "archivos_acuerdo"
    else:
        table_name = "archivos_acuerdo_comite"

    rows = (
        client.table(table_name)
        .select("id,ruta_storage")
        .eq("acuerdo_id", agreement_id)
        .execute()
        .data
        or []
    )
    _remove_storage_paths(client, [row.get("ruta_storage") for row in rows])
    if rows:
        client.table(table_name).delete().eq("acuerdo_id", agreement_id).execute()


def master_reset_agreement_control(
    client,
    *,
    agreement_id: str,
    agreement_number: str,
    kind: str,
    key: str,
) -> None:
    """Reinicia el seguimiento sin borrar el acuerdo base. Sólo administrador maestro."""
    if not is_master_admin():
        return

    label = "Junta de Gobierno" if kind == "board" else "Comité"
    rpc_name = "reset_acuerdo_junta_master" if kind == "board" else "reset_acuerdo_comite_master"

    with st.container(key=f"master_admin_panel_{key}"):
        with st.expander("🟧 Administración maestra · Reiniciar seguimiento del acuerdo"):
            st.warning(
                "Esta acción NO elimina el acuerdo, su número, título, texto ni la sesión. "
                "Sí pondrá en blanco su seguimiento: áreas y personas responsables, fecha compromiso, "
                "estatus, resultado, cierre/cumplimiento, comentarios/historial y avisos pendientes."
            )
    
            delete_documents = st.checkbox(
                "También eliminar los documentos de seguimiento de este acuerdo",
                value=False,
                key=f"master_reset_docs_{key}",
                help="Si no marcas esta opción, los documentos permanecerán en el expediente.",
            )
    
            confirmation = st.text_input(
                f"Para confirmar el reinicio de {agreement_number or 'este acuerdo'}, escribe LIMPIAR",
                key=f"master_reset_text_{key}",
                placeholder="LIMPIAR",
            )
    
            if st.button(
                f"Reiniciar seguimiento · {label}",
                key=f"master_reset_button_{key}",
                type="primary",
                use_container_width=True,
                disabled=confirmation.strip().upper() != "LIMPIAR",
            ):
                try:
                    if delete_documents:
                        _reset_agreement_documents(client, agreement_id, kind)
    
                    client.rpc(
                        rpc_name,
                        {
                            "p_acuerdo_id": agreement_id,
                        },
                    ).execute()
    
                    # Cambia la generación de llaves de los widgets. En Streamlit, los widgets ya
                    # instanciados pueden conservar valores anteriores aunque la BD ya se haya limpiado.
                    nonce_key = f"agreement_widget_nonce_{agreement_id}"
                    st.session_state[nonce_key] = int(st.session_state.get(nonce_key, 0)) + 1

                    st.success("Seguimiento reiniciado. El acuerdo base se conservó.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"No fue posible reiniciar el seguimiento: {exc}")

def user_can_project_direction(direction: str) -> bool:
    user = st.session_state.get("user", {})
    return user.get("rol") == "administrador" or direction in (user.get("direcciones_proyectos") or [])

st.set_page_config(page_title="COINVIERTE | Gestión Institucional", page_icon="🏛️", layout="wide")

st.markdown("""
<style>
:root { --ink:#35434b; --gray:#858e93; --blue:#0798cf; --green:#009b4c; --teal:#16ad8f; --purple:#a990c7; --orange:#f68b08; --paper:#f6f8f9; }
.stApp { background:linear-gradient(180deg,#fbfcfc 0%,#f1f5f6 100%); color:var(--ink); }
.block-container { max-width:1480px; padding-top:2.1rem; padding-bottom:4rem; }
[data-testid="stSidebar"] { background:linear-gradient(180deg,#535f66 0%,#778187 100%); border-right:0; }
[data-testid="stSidebar"] * { color:white; }
[data-testid="stSidebar"] hr { border-color:rgba(255,255,255,.22); }
[data-testid="stSidebar"] .stButton button { background:rgba(255,255,255,.08); border:1px solid rgba(255,255,255,.22); color:white; text-align:left; padding:.68rem .85rem; border-radius:10px; }
[data-testid="stSidebar"] .stButton button:hover { background:var(--blue); border-color:#fff; }
.brand { display:flex; align-items:center; gap:13px; }
.brand-mark { width:44px; height:44px; flex:0 0 44px; }
.brand-name { font-size:1.32rem; font-weight:800; letter-spacing:.055em; line-height:1; }
.brand-sub { font-size:.67rem; letter-spacing:.08em; opacity:.72; margin-top:5px; text-transform:uppercase; }
.side-logo { background:#fff; border-radius:13px; padding:12px 10px; margin:.3rem 0 1.35rem; box-shadow:0 8px 20px rgba(0,0,0,.12); }
.side-logo img { display:block; width:100%; height:auto; }
.hero { position:relative; overflow:hidden; background:#fff; padding:2.25rem 2.55rem 1.65rem; border-radius:22px; margin-bottom:1.2rem; border:1px solid #e1e7e9; border-bottom:7px solid var(--orange); box-shadow:0 18px 42px rgba(53,67,75,.11); }
.hero:after { content:""; position:absolute; width:260px; height:260px; border:42px solid rgba(7,152,207,.06); border-radius:50%; right:-120px; top:-125px; box-shadow:0 0 0 38px rgba(0,155,76,.035); }
.hero-logo { display:block; width:min(760px,78%); max-height:145px; object-fit:contain; object-position:left center; position:relative; z-index:1; }
.hero-copy { max-width:820px; font-size:1.02rem; color:#667279; margin:1.25rem 0 0; line-height:1.55; position:relative; z-index:1; }
.welcome { background:white; padding:.8rem 1rem; border-radius:11px; border:1px solid #e5ebee; margin:.4rem 0 1.2rem; color:#52616d; font-size:.9rem; }
.card { background:rgba(255,255,255,.97); border:1px solid #dfe7e9; border-top:5px solid var(--accent,var(--blue)); border-radius:17px; padding:1.45rem; min-height:175px; box-shadow:0 8px 24px rgba(53,67,75,.06); transition:.2s ease; }
.card:hover { transform:translateY(-3px); box-shadow:0 13px 30px rgba(53,67,75,.11); border-color:var(--accent,var(--blue)); }
.card-blue { --accent:var(--blue); } .card-green { --accent:var(--green); } .card-purple { --accent:var(--purple); }
.card-icon { display:inline-flex; align-items:center; justify-content:center; width:46px; height:46px; border-radius:13px; background:color-mix(in srgb,var(--accent,var(--blue)) 13%,white); color:var(--accent,var(--blue)); font-size:1.3rem; font-weight:800; margin-bottom:.65rem; }
.card h3 { margin:.15rem 0 .65rem; color:var(--ink); font-size:1.25rem; }
.muted { color:#647580; line-height:1.5; }
.choice-card { background:#fff; border:1px solid #dfe7e9; border-top:7px solid var(--accent,var(--blue)); border-radius:20px; padding:2rem 1.8rem 1.65rem; min-height:205px; box-shadow:0 12px 32px rgba(53,67,75,.08); text-align:center; margin-top:.7rem; }
.choice-card .choice-icon { width:68px; height:68px; border-radius:20px; display:flex; align-items:center; justify-content:center; margin:0 auto 1rem; background:color-mix(in srgb,var(--accent,var(--blue)) 14%,white); color:var(--accent,var(--blue)); font-size:1.6rem; font-weight:800; }
.choice-card h3 { font-size:1.45rem; margin:.25rem 0 .7rem; }
.choice-card p { color:#6d7980; margin:0; line-height:1.45; }
.choice-operations { --accent:var(--blue); } .choice-projects { --accent:var(--green); }
.choice-new { --accent:var(--orange); } .choice-edit { --accent:var(--purple); } .choice-view { --accent:var(--teal); }
.choice-title { text-align:center; margin:.7rem 0 .25rem; }
.choice-subtitle { text-align:center; color:#69767d; margin-bottom:1.2rem; }
.metric-grid { display:grid; grid-template-columns:repeat(4,1fr); gap:12px; margin:1rem 0 1.25rem; }
.metric-box { background:#fff; border:1px solid #dfe7e9; border-radius:14px; padding:1rem; border-top:4px solid var(--metric,var(--blue)); }
.metric-box .metric-label { color:#748087; font-size:.78rem; font-weight:700; text-transform:uppercase; letter-spacing:.03em; }
.metric-box .metric-value { color:var(--ink); font-size:1.55rem; font-weight:800; margin-top:.25rem; }
.metric-blue{--metric:var(--blue)} .metric-green{--metric:var(--green)} .metric-orange{--metric:var(--orange)} .metric-purple{--metric:var(--purple)}
.goal-heading { padding:.75rem 1rem; border-radius:10px; margin:1rem 0 .8rem; font-weight:800; border-left:8px solid var(--status-color); background:color-mix(in srgb,var(--status-color) 10%,white); }
.status-red{--status-color:#b85c62}.status-yellow{--status-color:#c5a44a}.status-green{--status-color:#65a37a}.status-gray{--status-color:#858e93}
[class*="st-key-master_admin_panel_"] details > summary { background:linear-gradient(90deg,#f68b08,#e56f00)!important; color:white!important; border:1px solid #d96500!important; border-radius:10px!important; font-weight:800!important; }
[class*="st-key-master_admin_panel_"] details > summary:hover { background:linear-gradient(90deg,#ff9b20,#f07400)!important; color:white!important; }
[class*="st-key-master_admin_panel_"] details > summary * { color:white!important; }
[class*="st-key-master_delete_panel_"] details > summary { background:linear-gradient(90deg,#f68b08,#e56f00)!important; color:white!important; border:1px solid #d96500!important; border-radius:10px!important; font-weight:800!important; }
[class*="st-key-master_delete_panel_"] details > summary:hover { background:linear-gradient(90deg,#ff9b20,#f07400)!important; color:white!important; }
[class*="st-key-master_delete_panel_"] details > summary * { color:white!important; }
.fin-progress-wrap { margin:.8rem 0 1rem; }
.fin-progress-head { display:flex; justify-content:space-between; gap:12px; align-items:end; margin-bottom:.45rem; }
.fin-progress-label { color:#5f6c74; font-weight:800; }
.fin-progress-value { color:var(--ink); font-size:1.35rem; font-weight:850; }
.fin-progress-track { width:100%; height:18px; background:#e7ecee; border-radius:999px; overflow:hidden; border:1px solid #dbe3e6; }
.fin-progress-bar { height:100%; width:var(--progress,0%); background:var(--progress-color,#f68b08); border-radius:999px; transition:width .25s ease; }
.fin-progress-foot { margin-top:.45rem; color:#647580; font-size:.88rem; display:flex; justify-content:space-between; gap:12px; flex-wrap:wrap; }
.year-card { background:#fff; border:1px solid #dfe7e9; border-top:6px solid var(--accent,var(--blue)); border-radius:18px; padding:1.25rem; text-align:center; box-shadow:0 8px 22px rgba(53,67,75,.07); margin-top:.5rem; }
.year-card h2 { font-size:2rem; margin:.1rem 0; color:var(--accent,var(--blue)); }
.analytics-banner { margin:1.35rem 0 .35rem; padding:1.25rem 1.5rem; border-radius:18px; color:#fff;
 background:linear-gradient(100deg,#173b63 0%,#0798cf 48%,#16ad8f 100%); box-shadow:0 12px 28px rgba(23,59,99,.2); }
.analytics-banner h3 { margin:0 0 .2rem; font-size:1.35rem; }.analytics-banner p { margin:0; opacity:.88; }
.analytics-metrics { display:grid; grid-template-columns:repeat(3,1fr); gap:15px; margin:1rem 0 1.7rem; }
.analytics-metric { --tone:var(--blue); position:relative; overflow:hidden; background:#fff; border:1px solid #dfe7e9;
 border-radius:18px; padding:1.15rem 1.25rem; min-height:125px; box-shadow:0 8px 22px rgba(53,67,75,.07); }
.analytics-metric:before { content:""; position:absolute; left:0; top:0; bottom:0; width:7px; background:var(--tone); }
.analytics-metric:after { content:""; position:absolute; width:88px; height:88px; border-radius:50%; right:-35px; top:-38px; background:color-mix(in srgb,var(--tone) 12%,white); }
.analytics-value { font-size:2.35rem; line-height:1; font-weight:850; color:var(--tone); margin:.1rem 0 .65rem; }
.analytics-label { color:#5f6c74; font-weight:700; line-height:1.3; max-width:90%; }
.session-column { background:rgba(255,255,255,.72); border:1px solid #dfe7e9; border-radius:18px; padding:1rem 1rem .4rem; min-height:220px; }
.session-card { background:#fff; border-left:7px solid var(--accent,var(--blue)); border-radius:13px; padding:1rem 1.1rem; margin:.7rem 0 .35rem; box-shadow:0 5px 16px rgba(53,67,75,.07); }
.session-card h4 { margin:0 0 .25rem; font-size:1.08rem; }.session-card p { margin:0; color:#738087; font-size:.83rem; }
[class*="st-key-session_open_"] div[data-testid="stButton"] button { background:#e7f4ec!important; border-color:#a8d3b7!important; color:#245c3b!important; font-weight:700; }
[class*="st-key-session_open_"] div[data-testid="stButton"] button:hover { background:#d6ecdf!important; border-color:#79b78e!important; }
@media(max-width:900px){.metric-grid,.analytics-metrics{grid-template-columns:repeat(2,1fr)}}
div[data-testid="stForm"] { background:white; padding:1.55rem; border-radius:18px; border:1px solid #dfe7e9; box-shadow:0 8px 24px rgba(20,55,70,.045); }
div[data-testid="stForm"] h3 { color:var(--gray); border-left:5px solid var(--orange); border-bottom:1px solid #e4ebed; padding:.15rem 0 .65rem .75rem; margin-top:1.2rem; }
.stButton button, .stFormSubmitButton button { border-radius:10px; }
.stButton button[kind="primary"], .stFormSubmitButton button[kind="primary"] { background:linear-gradient(90deg,#173b63,#6750a4)!important; border:0!important; color:white!important; }
.stButton button[kind="primary"]:hover, .stFormSubmitButton button[kind="primary"]:hover { background:linear-gradient(90deg,#0798cf,#6750a4)!important; }
[data-baseweb="tag"] { background-color:#6750a4!important; }
[data-testid="stBaseButton-primary"], [data-testid="stFormSubmitButton"] button { background:linear-gradient(90deg,#173b63,#6750a4)!important; border-color:#173b63!important; color:white!important; }
.stTabs [data-baseweb="tab"][aria-selected="true"] { color:#173b63!important; }
.stTabs [data-baseweb="tab-highlight"] { background-color:#6750a4!important; }
.stFormSubmitButton button[kind="primary"] { background:linear-gradient(90deg,var(--green),var(--teal)); border:0; }
.stFormSubmitButton button[kind="primary"]:hover { background:linear-gradient(90deg,var(--blue),var(--teal)); }
div[data-baseweb="radio"] div[aria-checked="true"] { color:var(--orange); }
div[data-baseweb="input"]:focus-within, div[data-baseweb="select"]:focus-within, div[data-baseweb="textarea"]:focus-within { border-color:var(--blue); box-shadow:0 0 0 1px var(--blue); }
[data-testid="stFileUploaderDropzone"] { background:#f5faf9; border-color:var(--teal); }
h1,h2,h3 { letter-spacing:-.018em; color:var(--ink); }
</style>
""", unsafe_allow_html=True)


def logo_data_uri():
    logo = Path("assets/logo_coinvierte.jpeg")
    if not logo.exists():
        return ""
    return "data:image/jpeg;base64," + base64.b64encode(logo.read_bytes()).decode()


def brand_html(sidebar=False):
    src = logo_data_uri()
    if src:
        css_class = "side-logo" if sidebar else ""
        return f'<div class="{css_class}"><img src="{src}" alt="COINVIERTE"></div>' if sidebar else f'<img class="hero-logo" src="{src}" alt="COINVIERTE">'
    return '<div class="brand"><div><div class="brand-name">COINVIERTE</div><div class="brand-sub">Agencia de Coinversión para el Desarrollo Sostenible de Jalisco</div></div></div>'


def logo_header():
    identity = brand_html()
    st.markdown(f'''<div class="hero">{identity}
    <p class="hero-copy">Plataforma institucional para la gestión integral, documentación y seguimiento de programas y proyectos.</p></div>''', unsafe_allow_html=True)


def login():
    logo_header()
    st.subheader("Acceso institucional")
    if not configured():
        st.warning("Modo demostración: falta configurar Supabase. Puedes ingresar con cualquier correo @jalisco.gob.mx.")
    login_tab, activation_tab = st.tabs(["Ingresar", "Activar acceso con código"])
    with login_tab:
        with st.form("login"):
            email = st.text_input("Correo institucional", placeholder="nombre@jalisco.gob.mx")
            password = st.text_input("Contraseña", type="password")
            submitted = st.form_submit_button("Ingresar", type="primary", use_container_width=True)
    with activation_tab:
        st.caption("Utiliza el código temporal entregado por el administrador. El código sólo puede usarse una vez.")
        with st.form("activate_access"):
            activation_email = st.text_input("Correo autorizado", placeholder="nombre@jalisco.gob.mx", key="activation_email")
            code = st.text_input("Código temporal", max_chars=8)
            new_password = st.text_input("Crea una contraseña", type="password", key="new_password")
            confirm_password = st.text_input("Confirma la contraseña", type="password")
            activate = st.form_submit_button("Activar mi acceso", type="primary", use_container_width=True)
    if activate:
        if not configured():
            st.error("Primero debes conectar Supabase.")
        elif not valid_official_email(activation_email):
            st.error("El correo debe pertenecer a @jalisco.gob.mx.")
        elif len(new_password) < 8:
            st.error("La contraseña debe tener al menos 8 caracteres.")
        elif new_password != confirm_password:
            st.error("Las contraseñas no coinciden.")
        else:
            try:
                auth = public_client().auth.sign_up({"email": activation_email.lower().strip(), "password": new_password})
                redeemed = public_client().rpc("canjear_codigo_acceso", {"p_email": activation_email.lower().strip(),
                                                                          "p_codigo": code.strip()}).execute().data
                if not redeemed:
                    st.error("El código es incorrecto, ya fue utilizado o está vencido.")
                else:
                    if auth.session and auth.user:
                        st.success("Acceso activado correctamente. Ya puedes ingresar.")
                    else:
                        st.success("Acceso activado. Revisa tu correo si Supabase solicita confirmar la cuenta.")
            except Exception as exc:
                st.error(f"No fue posible activar el acceso: {exc}")
    if submitted:
        if not valid_official_email(email):
            st.error("El acceso está limitado a cuentas @jalisco.gob.mx.")
        elif not configured():
            st.session_state.user = {"email": email.lower(), "id": "demo"}
            st.rerun()
        else:
            try:
                auth = public_client().auth.sign_in_with_password({"email": email, "password": password})
            except Exception as exc:
                message = str(exc)
                if "Invalid login credentials" in message:
                    st.error("Supabase rechazó el correo o la contraseña: Invalid login credentials.")
                elif "Email not confirmed" in message:
                    st.error("La cuenta existe, pero el correo todavía no está confirmado en Supabase.")
                else:
                    st.error(f"Supabase no pudo autenticar la cuenta: {message}")
                return
            if not auth.user or not valid_official_email(auth.user.email or ""):
                public_client().auth.sign_out()
                st.error("La cuenta no pertenece al dominio autorizado.")
                return
            st.session_state.access_token = auth.session.access_token
            st.session_state.refresh_token = auth.session.refresh_token
            try:
                user_client = client_with_token(auth.session.access_token, auth.session.refresh_token)
                profile = access_profile(user_client, auth.user.email)
            except Exception as exc:
                st.error(f"La contraseña fue aceptada, pero falló la consulta de autorización: {exc}")
                return
            if not profile or not profile.get("activo"):
                public_client().auth.sign_out()
                st.session_state.pop("access_token", None)
                st.session_state.pop("refresh_token", None)
                st.error("La contraseña fue aceptada, pero tu acceso no está autorizado o está suspendido.")
                return
            st.session_state.user = {"email": auth.user.email, "id": str(auth.user.id),
                                     "nombre": profile.get("nombre") or auth.user.email,
                                     "rol": profile.get("rol", "usuario"), "direccion": profile.get("direccion"),
                                     "modulos": profile.get("modulos") or [],
                                     "direcciones_proyectos": profile.get("direcciones_proyectos") or []}
            register_access(user_client)
            st.rerun()


def landing():
    logo_header()
    st.markdown(f'<div class="welcome">Sesión institucional activa · <b>{st.session_state.user["email"]}</b></div>', unsafe_allow_html=True)
    sections = [item for item in [
        ("01", "Programas / Proyectos", "Alta, consulta, edición y seguimiento de expedientes."),
        ("02", "Junta de Gobierno", "Actas, acuerdos y documentación de las sesiones."),
        ("03", "Comités", "Integración, sesiones, actas y dictaminación."),
        ("04", "Oficios Dirección General", "Archivo anual y mensual de los oficios firmados por la Dirección General."),
    ] if user_can(item[1])]
    if not sections:
        st.info("Tu cuenta está activa, pero todavía no tiene módulos asignados. Solicita al administrador que configure tus permisos.")
        return
    cols = st.columns(len(sections))
    card_styles = ["card-blue", "card-green", "card-purple", "card-blue"]
    for col, (icon, title, text), card_style in zip(cols, sections, card_styles):
        with col:
            st.markdown(f'<div class="card {card_style}"><div class="card-icon">{icon}</div><h3>{title}</h3><p class="muted">{text}</p></div>', unsafe_allow_html=True)
            if st.button(f"Abrir {title}", key=title, use_container_width=True):
                st.session_state.page = title
                st.rerun()


def objective_fields(existing=None):
    existing = existing or [""]
    if "objective_count" not in st.session_state:
        st.session_state.objective_count = max(1, len(existing))
    values = []
    for index in range(st.session_state.objective_count):
        default = existing[index] if index < len(existing) else ""
        values.append(st.text_area(f"Objetivo específico {index + 1}", value=default, key=f"obj_{index}"))
    col1, col2 = st.columns([1, 4])
    if col1.form_submit_button("＋ Agregar objetivo"):
        st.session_state.objective_count += 1
        st.rerun()
    if st.session_state.objective_count > 1 and col2.form_submit_button("Quitar último"):
        st.session_state.objective_count -= 1
        st.rerun()
    return values


def execution_goal_fields(existing=None):
    existing = existing or []
    if "execution_goal_count" not in st.session_state:
        st.session_state.execution_goal_count = max(1, len(existing))
    goals, evidence_groups = [], []
    statuses = ["Por iniciar", "En progreso", "Terminada"]
    for index in range(st.session_state.execution_goal_count):
        saved = existing[index] if index < len(existing) else {}
        saved_status = saved.get("estatus", "Por iniciar")
        display_status = st.session_state.get(f"goal_status_{index}", saved_status)
        status_class = {"Por iniciar": "status-red", "En progreso": "status-yellow", "Terminada": "status-green"}.get(display_status, "status-gray")
        st.markdown(f'<div class="goal-heading {status_class}">Meta de ejecución {index + 1} · {display_status}</div>', unsafe_allow_html=True)
        name = st.text_input("Nombre de la meta", value=saved.get("nombre", ""), key=f"goal_name_{index}")
        description = st.text_area("Descripción", value=saved.get("descripcion", ""), key=f"goal_description_{index}", height=90)
        g1, g2 = st.columns(2)
        status = g1.selectbox("Estatus", statuses, index=statuses.index(saved_status) if saved_status in statuses else 0,
                              key=f"goal_status_{index}")
        target_date = g2.text_input("Fecha objetivo", value=saved.get("fecha_objetivo", ""),
                                    placeholder="Ej. 30/09/2026", key=f"goal_date_{index}")
        evidence = st.file_uploader("Evidencia de la meta (fotografías o documentos)", accept_multiple_files=True,
                                    key=f"goal_evidence_{index}")
        goals.append({"nombre": name.strip(), "descripcion": description.strip(), "estatus": status,
                      "fecha_objetivo": target_date.strip(),
                      "evidencias_nombres": [file.name for file in evidence]})
        evidence_groups.append(evidence)
    c1, c2 = st.columns([1, 3])
    add_goal = c1.form_submit_button("＋ Agregar meta", use_container_width=True)
    remove_goal = c2.form_submit_button("Quitar última meta", use_container_width=True) if st.session_state.execution_goal_count > 1 else False
    return goals, evidence_groups, add_goal, remove_goal


RISK_COLUMNS = ["id", "riesgo", "categoria", "descripcion", "causa", "consecuencia",
                "probabilidad", "impacto", "mitigacion", "responsable", "fecha_compromiso",
                "estatus", "observaciones", "puntaje", "nivel", "eliminar"]


def risk_level(score: int) -> str:
    if score <= 4: return "Bajo"
    if score <= 9: return "Moderado"
    if score <= 16: return "Alto"
    return "Crítico"


def normalize_risks(records) -> list[dict]:
    clean = []
    for raw in records or []:
        item = dict(raw)
        if item.get("eliminar") is True:
            continue
        name = str(item.get("riesgo") or "").strip()
        if not name:
            continue
        try:
            probability = max(1, min(5, int(float(item.get("probabilidad") or 1))))
            impact = max(1, min(5, int(float(item.get("impacto") or 1))))
        except (TypeError, ValueError):
            probability, impact = 1, 1
        score = probability * impact
        clean.append({"id": str(item.get("id") or uuid.uuid4()), "riesgo": name,
            "categoria": str(item.get("categoria") or "Otro").strip(),
            "descripcion": str(item.get("descripcion") or "").strip(),
            "causa": str(item.get("causa") or "").strip(),
            "consecuencia": str(item.get("consecuencia") or "").strip(),
            "probabilidad": probability, "impacto": impact,
            "mitigacion": str(item.get("mitigacion") or "").strip(),
            "responsable": str(item.get("responsable") or "").strip(),
            "fecha_compromiso": str(item.get("fecha_compromiso") or "").strip(),
            "estatus": str(item.get("estatus") or "Por iniciar").strip(),
            "observaciones": str(item.get("observaciones") or "").strip(),
            "puntaje": score, "nivel": risk_level(score), "eliminar": False})
    return clean


def risks_from_excel(uploaded_file) -> list[dict]:
    frame = pd.read_excel(uploaded_file, sheet_name="Matriz de riesgos", header=3)
    aliases = {"ID (no modificar)": "id", "Riesgo": "riesgo", "Categoría": "categoria",
        "Descripción": "descripcion", "Causa": "causa", "Consecuencia": "consecuencia",
        "Probabilidad (1-5)": "probabilidad", "Impacto (1-5)": "impacto",
        "Mitigación": "mitigacion", "Responsable": "responsable",
        "Fecha compromiso": "fecha_compromiso", "Estatus": "estatus", "Observaciones": "observaciones"}
    missing = [column for column in aliases if column not in frame.columns]
    if missing:
        raise ValueError("Faltan columnas de la plantilla: " + ", ".join(missing))
    frame = frame.rename(columns=aliases)[list(aliases.values())].where(pd.notna(frame), "")
    return normalize_risks(frame.to_dict("records"))


def risk_summary(risks: list[dict]) -> str:
    if not risks:
        return "Todavía no se han registrado riesgos."
    counts = {level: sum(1 for risk in risks if risk["nivel"] == level)
              for level in ["Bajo", "Moderado", "Alto", "Crítico"]}
    priority = sorted(risks, key=lambda risk: risk["puntaje"], reverse=True)[:3]
    main = ", ".join(f'{risk["riesgo"]} ({risk["puntaje"]}, {risk["nivel"]})' for risk in priority)
    materialized = sum(1 for risk in risks if risk["estatus"] == "Materializado")
    closed = sum(1 for risk in risks if risk["estatus"] == "Mitigado / cerrado")
    return (f'Riesgos: {len(risks)} · Críticos: {counts["Crítico"]} · Altos: {counts["Alto"]} · '
            f'Moderados: {counts["Moderado"]} · Bajos: {counts["Bajo"]}. Materializados: {materialized}; '
            f'mitigados/cerrados: {closed}. Principales: {main}.')


def _financial_progress_color(percent: float) -> tuple[str, str]:
    if percent > 100:
        return "#c23b3f", "Excede monto autorizado"
    if percent >= 100:
        return "#b7f34a", "Dispersión completa"
    if percent >= 60:
        return "#86c98a", "Avance alto"
    if percent >= 30:
        return "#f1d34f", "Avance medio"
    return "#f39a24", "Avance inicial"


def _project_transfers(client, project_id: str) -> list[dict]:
    if not client or not project_id:
        return []
    return (client.table("transferencias_proyecto").select("*")
            .eq("proyecto_id", project_id).order("fecha_transferencia", desc=False)
            .order("created_at", desc=False).execute().data or [])


def _transfer_documents(client, transfer_id: str) -> list[dict]:
    return (client.table("documentos_transferencia_proyecto").select("*")
            .eq("transferencia_id", transfer_id).order("created_at").execute().data or [])


def _financial_totals(client, project: dict) -> tuple[float, float, float]:
    total_project = float(project.get("monto") or 0)
    transfers = _project_transfers(client, str(project.get("id") or ""))
    dispersed = sum(float(row.get("importe") or 0) for row in transfers)
    percent = (dispersed / total_project * 100) if total_project else 0.0
    return total_project, dispersed, percent


def _render_financial_progress(total_project: float, dispersed: float, percent: float):
    pending = total_project - dispersed
    color, label = _financial_progress_color(percent)
    width = max(0.0, min(percent, 100.0))
    st.markdown(
        f"""<div class="fin-progress-wrap">
        <div class="fin-progress-head">
          <div><div class="fin-progress-label">Avance de dispersión</div>
          <div class="fin-progress-value">{percent:.1f}%</div></div>
          <div class="muted">{html.escape(label)}</div>
        </div>
        <div class="fin-progress-track"><div class="fin-progress-bar" style="--progress:{width:.2f}%;--progress-color:{color};"></div></div>
        <div class="fin-progress-foot">
          <span><b>${dispersed:,.2f}</b> dispersados de <b>${total_project:,.2f}</b></span>
          <span>Pendiente: <b>${pending:,.2f}</b></span>
        </div></div>""",
        unsafe_allow_html=True,
    )
    if percent > 100:
        st.error("La dispersión registrada excede el monto autorizado del proyecto.")


def _sync_project_financial_progress(client, project: dict) -> None:
    if not project or not project.get("id"):
        return
    _, dispersed, percent = _financial_totals(client, project)
    current = dict(project.get("avance_proyecto") or {})
    current["presupuesto_dispersado"] = round(dispersed, 2)
    current["porcentaje_financiero"] = round(percent, 2)
    client.table("proyectos").update({
        "avance_proyecto": current,
        "updated_at": datetime.now().isoformat(),
    }).eq("id", project["id"]).execute()
    project["avance_proyecto"] = current


def _upload_transfer_document(client, project_id: str, transfer_id: str, document_type: str, uploaded) -> None:
    folder = safe_name(document_type.lower().replace(" ", "_"))
    path = f"proyectos/{project_id}/transferencias/{transfer_id}/{folder}/{uuid.uuid4().hex}_{safe_name(uploaded.name)}"
    _upload_junta_document(client, path, uploaded)
    client.table("documentos_transferencia_proyecto").insert({
        "transferencia_id": transfer_id,
        "tipo_documento": document_type,
        "nombre_visible": Path(uploaded.name).stem,
        "nombre_archivo": uploaded.name,
        "ruta_storage": path,
        "mime_type": uploaded.type,
        "tamano_bytes": uploaded.size,
        "subido_por": st.session_state.user["id"],
        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"),
    }).execute()


def _delete_transfer(client, transfer: dict, project: dict) -> None:
    documents = _transfer_documents(client, str(transfer["id"]))
    _remove_storage_paths(client, [row.get("ruta_storage") for row in documents])
    client.table("transferencias_proyecto").delete().eq("id", transfer["id"]).execute()
    _sync_project_financial_progress(client, project)


def project_financial_documents(project: dict):
    if not project or project.get("direccion") != "Dirección de Proyectos":
        return
    if not configured():
        st.info("La comprobación financiera estará disponible al conectar Supabase.")
        return

    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
    project_id = str(project["id"])
    total_project, dispersed, percent = _financial_totals(client, project)

    st.markdown("---")
    st.markdown("## Comprobación y dispersión de recursos")
    st.caption("Registra cada transferencia y conserva sus comprobantes, facturas/CFDI y documentos soporte.")
    _render_financial_progress(total_project, dispersed, percent)

    transfers = _project_transfers(client, project_id)
    st.markdown("### Registrar nueva transferencia")
    with st.form(f"new_transfer_{project_id}", clear_on_submit=True):
        c1, c2 = st.columns(2)
        transfer_date = c1.date_input("Fecha de transferencia *", value=date.today())
        transfer_amount = c2.number_input("Importe dispersado (MXN) *", min_value=0.0, step=1000.0, format="%.2f")
        c3, c4 = st.columns(2)
        beneficiary = c3.text_input("Beneficiario / proveedor")
        reference = c4.text_input("Referencia / folio")
        concept = st.text_area("Concepto / descripción *", height=90)
        c5, c6 = st.columns(2)
        proof = c5.file_uploader("Comprobante de transferencia", type=["pdf", "jpg", "jpeg", "png"], key=f"transfer_proof_{project_id}")
        invoice = c6.file_uploader("Factura / CFDI", type=["pdf", "xml", "jpg", "jpeg", "png"], key=f"transfer_invoice_{project_id}")
        support = st.file_uploader("Otros documentos soporte", type=["pdf", "docx", "xlsx", "xls", "jpg", "jpeg", "png"],
                                   accept_multiple_files=True, key=f"transfer_support_{project_id}")
        add_transfer = st.form_submit_button("Guardar transferencia", type="primary", use_container_width=True)

    if add_transfer:
        if transfer_amount <= 0:
            st.error("El importe debe ser mayor a cero.")
        elif not concept.strip():
            st.error("Escribe el concepto de la transferencia.")
        else:
            try:
                result = client.table("transferencias_proyecto").insert({
                    "proyecto_id": project_id,
                    "fecha_transferencia": transfer_date.isoformat(),
                    "importe": float(transfer_amount),
                    "concepto": concept.strip(),
                    "beneficiario": beneficiary.strip() or None,
                    "referencia": reference.strip() or None,
                    "creado_por": st.session_state.user["id"],
                    "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"),
                }).execute()
                transfer_id = str(result.data[0]["id"])
                if proof:
                    _upload_transfer_document(client, project_id, transfer_id, "Comprobante de transferencia", proof)
                if invoice:
                    _upload_transfer_document(client, project_id, transfer_id, "Factura / CFDI", invoice)
                for file in support or []:
                    _upload_transfer_document(client, project_id, transfer_id, "Documento soporte", file)
                _sync_project_financial_progress(client, project)
                st.success("Transferencia y documentos guardados.")
                st.rerun()
            except Exception as exc:
                st.error(f"No fue posible guardar la transferencia: {exc}")

    st.markdown(f"### Historial de transferencias ({len(transfers)})")
    if not transfers:
        st.info("Todavía no hay transferencias registradas para este proyecto.")
        return

    for index, transfer in enumerate(transfers, 1):
        transfer_id = str(transfer["id"])
        amount_value = float(transfer.get("importe") or 0)
        title = f"{index}. {transfer.get('fecha_transferencia') or 'Sin fecha'} · ${amount_value:,.2f}"
        with st.expander(title, expanded=False):
            st.markdown(f"**Concepto:** {transfer.get('concepto') or 'Sin descripción'}")
            st.markdown(f"**Beneficiario / proveedor:** {transfer.get('beneficiario') or 'Sin información'}  \n"
                        f"**Referencia / folio:** {transfer.get('referencia') or 'Sin información'}")

            with st.form(f"edit_transfer_{transfer_id}"):
                ec1, ec2 = st.columns(2)
                saved_date = date.fromisoformat(str(transfer.get("fecha_transferencia"))[:10]) if transfer.get("fecha_transferencia") else date.today()
                new_date = ec1.date_input("Fecha", value=saved_date, key=f"edit_date_{transfer_id}")
                new_amount = ec2.number_input("Importe", min_value=0.0, value=amount_value, step=1000.0, format="%.2f",
                                              key=f"edit_amount_{transfer_id}")
                new_beneficiary = st.text_input("Beneficiario / proveedor", value=transfer.get("beneficiario") or "", key=f"edit_beneficiary_{transfer_id}")
                new_reference = st.text_input("Referencia / folio", value=transfer.get("referencia") or "", key=f"edit_reference_{transfer_id}")
                new_concept = st.text_area("Concepto / descripción", value=transfer.get("concepto") or "", key=f"edit_concept_{transfer_id}")
                save_edit = st.form_submit_button("Guardar cambios", use_container_width=True)
            if save_edit:
                if new_amount <= 0 or not new_concept.strip():
                    st.error("El importe debe ser mayor a cero y el concepto es obligatorio.")
                else:
                    client.table("transferencias_proyecto").update({
                        "fecha_transferencia": new_date.isoformat(),
                        "importe": float(new_amount),
                        "beneficiario": new_beneficiary.strip() or None,
                        "referencia": new_reference.strip() or None,
                        "concepto": new_concept.strip(),
                        "updated_at": datetime.now().isoformat(),
                    }).eq("id", transfer_id).execute()
                    _sync_project_financial_progress(client, project)
                    st.success("Transferencia actualizada.")
                    st.rerun()

            pending_key = f"delete_transfer_pending_{transfer_id}"
            if st.button("Eliminar transferencia", key=f"delete_transfer_{transfer_id}"):
                st.session_state[pending_key] = True
                st.rerun()
            if st.session_state.get(pending_key):
                st.warning("Esta acción eliminará la transferencia y todos sus documentos asociados.")
                y, n, _ = st.columns([1, 1, 4])
                if y.button("Sí, eliminar", key=f"confirm_delete_transfer_{transfer_id}", type="primary"):
                    _delete_transfer(client, transfer, project)
                    st.session_state.pop(pending_key, None)
                    st.success("Transferencia eliminada.")
                    st.rerun()
                if n.button("Cancelar", key=f"cancel_delete_transfer_{transfer_id}"):
                    st.session_state.pop(pending_key, None)
                    st.rerun()

            documents = _transfer_documents(client, transfer_id)
            st.markdown(f"#### Documentos ({len(documents)})")
            upload_nonce = st.session_state.get(f"transfer_doc_nonce_{transfer_id}", 0)
            d1, d2 = st.columns(2)
            new_type = d1.selectbox("Tipo de documento", ["Comprobante de transferencia", "Factura / CFDI", "Documento soporte"],
                                    key=f"transfer_doc_type_{transfer_id}_{upload_nonce}")
            new_doc = d2.file_uploader("Agregar documento", type=["pdf", "xml", "docx", "xlsx", "xls", "jpg", "jpeg", "png"],
                                       key=f"transfer_doc_file_{transfer_id}_{upload_nonce}")
            if st.button("Incorporar documento", key=f"add_transfer_doc_{transfer_id}", disabled=not new_doc, use_container_width=True):
                _upload_transfer_document(client, project_id, transfer_id, new_type, new_doc)
                st.session_state[f"transfer_doc_nonce_{transfer_id}"] = upload_nonce + 1
                st.success("Documento incorporado.")
                st.rerun()

            for document in documents:
                _document_card(client, document, f"transfer_doc_{transfer_id}", "documentos_transferencia_proyecto")


def project_form(direction: str, project=None):
    project = project or {}
    risk_context = str(project.get("id") or "new")
    if st.session_state.get("risk_context") != risk_context:
        st.session_state.risk_context = risk_context
        saved_risks = ((project.get("avance_proyecto") or {}).get("matriz_riesgos") or [])
        st.session_state.risk_rows = normalize_risks(saved_risks)
    st.subheader("Editar proyecto" if project else "Dar de alta nuevo proyecto")
    with st.form("project_form", clear_on_submit=False):
        is_projects = direction == "Dirección de Proyectos"
        if is_projects:
            general_tab, advance_tab, risks_tab = st.tabs(["Ficha general", "Avance", "Matriz de riesgos"])
        else:
            general_tab, advance_tab, risks_tab = st.container(), None, None

        with general_tab:
            st.markdown("### Información General")
            name = st.text_input("Nombre del proyecto *", value=project.get("nombre", ""))
            applicant = st.text_input("Nombre del solicitante *", value=project.get("solicitante", ""))
            municipality = st.selectbox("Municipio de ejecución *", MUNICIPIOS_JALISCO,
                                        index=MUNICIPIOS_JALISCO.index(project["municipio"]) if project.get("municipio") in MUNICIPIOS_JALISCO else 0)
            c1, c2 = st.columns(2)
            year = c1.number_input("Año de inicio *", min_value=2000, max_value=2100,
                                   value=int(project.get("anio_inicio", datetime.now().year)), step=1)
            amount = c2.number_input("Monto (MXN) *", min_value=0.0, value=float(project.get("monto", 0)), step=1000.0, format="%.2f")
            general = st.text_area("Objetivo general *", value=project.get("objetivo_general", ""), height=130)
            st.markdown("#### Objetivos específicos")
            objectives = objective_fields(project.get("objetivos_especificos"))

            st.markdown("### Gestión Documental")
            legal = st.file_uploader("Documentación jurídica", accept_multiple_files=True, key="legal")
            auxiliary = st.file_uploader("Documentación auxiliar", accept_multiple_files=True, key="aux")
            committee = st.file_uploader("Acta del Comité de Dictaminación", accept_multiple_files=False, key="committee")
            board = st.file_uploader("Acta de aprobación de Junta de Gobierno", accept_multiple_files=False, key="board")
            agreement = st.file_uploader("Convenio de colaboración", accept_multiple_files=False, key="agreement")

            st.markdown("### Evidencia fotográfica")
            photos = st.file_uploader("Fotografías generales (máximo 5 MB por archivo)", type=["jpg", "jpeg", "png", "webp"],
                                      accept_multiple_files=True, key="photos")

            if not is_projects:
                st.markdown("### Monitoreo y Seguimiento")
                previous_monitoring = project.get("monitoreo", {}) or {}
                m1, m2 = st.columns(2)
                statuses = ["Sin iniciar", "En planeación", "En ejecución", "Suspendido", "Concluido"]
                current_status = previous_monitoring.get("estatus", "Sin iniciar")
                status = m1.selectbox("Estatus del proyecto", statuses,
                                      index=statuses.index(current_status) if current_status in statuses else 0)
                responsible = m2.text_input("Responsable del seguimiento", value=previous_monitoring.get("responsable", ""))
                m3, m4 = st.columns(2)
                period = m3.text_input("Periodo de seguimiento", value=previous_monitoring.get("periodo", ""))
                progress = m4.slider("Porcentaje de avance", 0, 100, int(previous_monitoring.get("avance", 0)), 5)
                monitoring_progress = st.text_area("Principales avances", value=previous_monitoring.get("avances", ""), height=100)
                pending = st.text_area("Pendientes o riesgos", value=previous_monitoring.get("pendientes", ""), height=90)
                next_actions = st.text_area("Próximas acciones", value=previous_monitoring.get("proximas_acciones", ""), height=90)
                observations = st.text_area("Observaciones de seguimiento", value=previous_monitoring.get("observaciones", ""), height=90)

        execution_goals, goal_evidence_groups, add_goal, remove_goal = [], [], False, False
        budget_dispersed = 0.0
        if is_projects:
            with advance_tab:
                saved_advance = project.get("avance_proyecto", {}) or {}
                st.markdown("### Ejecución financiera")
                if project and configured():
                    financial_client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
                    _, budget_dispersed, budget_pct = _financial_totals(financial_client, {**project, "monto": amount})
                    _render_financial_progress(float(amount), budget_dispersed, budget_pct)
                    st.caption("El avance financiero se calcula automáticamente con las transferencias registradas.")
                else:
                    budget_dispersed = 0.0
                    budget_pct = 0.0
                    _render_financial_progress(float(amount), 0.0, 0.0)
                    st.info("Guarda primero el proyecto para poder registrar transferencias, facturas y comprobantes.")
                saved_goals = saved_advance.get("metas", []) or []
                st.markdown("### Metas de ejecución")
                execution_goals, goal_evidence_groups, add_goal, remove_goal = execution_goal_fields(saved_goals)
                completed = sum(1 for goal in execution_goals if goal["estatus"] == "Terminada")
                in_progress = sum(1 for goal in execution_goals if goal["estatus"] == "En progreso")
                physical_pct = ((completed * 100 + in_progress * 50) / len(execution_goals)) if execution_goals else 0
                traffic = "Verde" if physical_pct >= 80 else "Amarillo" if physical_pct >= 40 else "Rojo"
                traffic_class = "metric-green" if traffic == "Verde" else "metric-orange" if traffic == "Amarillo" else "metric-blue"
                st.markdown(f'''<div class="metric-grid">
                  <div class="metric-box metric-blue"><div class="metric-label">Avance financiero</div><div class="metric-value">{budget_pct:.1f}%</div></div>
                  <div class="metric-box metric-green"><div class="metric-label">Metas terminadas</div><div class="metric-value">{completed}/{len(execution_goals)}</div></div>
                  <div class="metric-box metric-purple"><div class="metric-label">Avance físico</div><div class="metric-value">{physical_pct:.1f}%</div></div>
                  <div class="metric-box {traffic_class}"><div class="metric-label">Semáforo general</div><div class="metric-value">{traffic}</div></div>
                </div>''', unsafe_allow_html=True)
                refresh_metrics = st.form_submit_button("Actualizar indicadores", use_container_width=True)

            with risks_tab:
                st.markdown("### Matriz de riesgos")
                st.caption("Edita la tabla. Para borrar un riesgo, marca Eliminar o usa el control de eliminación de filas.")
                template_path = Path("plantilla_matriz_riesgos_coinvierte.xlsx")
                if template_path.exists():
                    template_b64 = base64.b64encode(template_path.read_bytes()).decode()
                    st.markdown(f'<a download="plantilla_matriz_riesgos_coinvierte.xlsx" href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{template_b64}">⬇️ Descargar plantilla Excel</a>', unsafe_allow_html=True)
                risk_file = st.file_uploader("Cargar plantilla llena (.xlsx)", type=["xlsx"], key=f"risk_excel_{risk_context}")
                import_risks = st.form_submit_button("Ingestar riesgos desde Excel", use_container_width=True)
                risk_frame = pd.DataFrame(st.session_state.get("risk_rows", []), columns=RISK_COLUMNS)
                edited_risks = st.data_editor(risk_frame, key=f"risk_editor_{risk_context}", num_rows="dynamic",
                    use_container_width=True, hide_index=True, disabled=["id", "puntaje", "nivel"],
                    column_config={"id": None,
                        "riesgo": st.column_config.TextColumn("Riesgo", required=True, width="medium"),
                        "categoria": st.column_config.SelectboxColumn("Categoría", options=["Financiero", "Operativo", "Jurídico / normativo", "Técnico", "Ambiental", "Social", "Reputacional", "Cronograma", "Otro"], width="medium"),
                        "descripcion": st.column_config.TextColumn("Descripción", width="large"),
                        "causa": st.column_config.TextColumn("Causa", width="large"),
                        "consecuencia": st.column_config.TextColumn("Consecuencia", width="large"),
                        "probabilidad": st.column_config.NumberColumn("Probabilidad 1–5", min_value=1, max_value=5, step=1, required=True),
                        "impacto": st.column_config.NumberColumn("Impacto 1–5", min_value=1, max_value=5, step=1, required=True),
                        "mitigacion": st.column_config.TextColumn("Mitigación", width="large"),
                        "responsable": st.column_config.TextColumn("Responsable", width="medium"),
                        "fecha_compromiso": st.column_config.TextColumn("Fecha compromiso", help="Formato sugerido: AAAA-MM-DD"),
                        "estatus": st.column_config.SelectboxColumn("Estatus", options=["Por iniciar", "En seguimiento", "Materializado", "Mitigado / cerrado"], width="medium"),
                        "observaciones": st.column_config.TextColumn("Observaciones", width="large"),
                        "puntaje": st.column_config.NumberColumn("Puntaje", help="Probabilidad × Impacto"),
                        "nivel": st.column_config.TextColumn("Nivel"),
                        "eliminar": st.column_config.CheckboxColumn("Eliminar", help="Marca para borrar al guardar.")})
                recalculate_risks = st.form_submit_button("Actualizar puntajes y síntesis", use_container_width=True)
                current_risks = normalize_risks(edited_risks.to_dict("records"))
                st.info(risk_summary(current_risks))
        else:
            budget_pct, physical_pct, traffic, completed = 0, float(progress), current_status, 0
            refresh_metrics, import_risks, recalculate_risks, risk_file, edited_risks = False, False, False, None, None

        st.markdown("### Ficha del proyecto")
        st.caption("La ficha incluye información general, monitoreo y evidencia fotográfica. Gestión Documental no se muestra.")
        b1, b2 = st.columns(2)
        preview = b1.form_submit_button("Previsualizar ficha", use_container_width=True)
        save = b2.form_submit_button("Guardar proyecto", type="primary", use_container_width=True)

    if add_goal:
        st.session_state.execution_goal_count += 1
        st.rerun()
    if remove_goal:
        st.session_state.execution_goal_count -= 1
        st.rerun()
    if refresh_metrics:
        st.rerun()
    if import_risks:
        if risk_file is None:
            st.error("Selecciona primero el archivo Excel lleno.")
        else:
            try:
                st.session_state.risk_rows = risks_from_excel(risk_file)
                st.success(f"Se importaron {len(st.session_state.risk_rows)} riesgos.")
                st.rerun()
            except Exception as exc:
                st.error(f"No fue posible importar la plantilla: {exc}")
    if recalculate_risks:
        st.session_state.risk_rows = normalize_risks(edited_risks.to_dict("records"))
        st.rerun()

    if save or preview:
        errors = []
        if not name.strip() or not applicant.strip() or not general.strip():
            errors.append("Completa todos los campos obligatorios.")
        clean_objectives = [o.strip() for o in objectives if o.strip()]
        if not clean_objectives:
            errors.append("Agrega al menos un objetivo específico.")
        oversized = [f.name for f in photos if f.size > 5 * 1024 * 1024]
        oversized += [f.name for files in goal_evidence_groups for f in files if f.size > 5 * 1024 * 1024]
        if oversized:
            errors.append("Estas fotografías exceden 5 MB: " + ", ".join(oversized))
        if errors:
            for error in errors:
                st.error(error)
            return
        if is_projects:
            risk_rows = normalize_risks(edited_risks.to_dict("records"))
            st.session_state.risk_rows = risk_rows
            advance_data = {"presupuesto_dispersado": budget_dispersed, "porcentaje_financiero": round(budget_pct, 2),
                            "porcentaje_fisico": round(physical_pct, 2), "semaforo": traffic,
                            "metas_terminadas": completed, "metas": execution_goals,
                            "matriz_riesgos": risk_rows, "sintesis_riesgos": risk_summary(risk_rows)}
            monitoring_data = {"estatus": traffic, "responsable": "", "periodo": "", "avance": round(physical_pct),
                               "avances": f"{completed} de {len(execution_goals)} metas terminadas. Avance financiero: {budget_pct:.1f}%.",
                               "pendientes": "", "proximas_acciones": "", "observaciones": ""}
        else:
            advance_data = {}
            monitoring_data = {"estatus": status, "responsable": responsible.strip(), "periodo": period.strip(),
                               "avance": int(progress), "avances": monitoring_progress.strip(),
                               "pendientes": pending.strip(), "proximas_acciones": next_actions.strip(),
                               "observaciones": observations.strip()}
        payload = {"direccion": direction, "nombre": name.strip(), "solicitante": applicant.strip(),
                   "municipio": municipality, "anio_inicio": int(year), "monto": amount,
                   "objetivo_general": general.strip(), "objetivos_especificos": clean_objectives,
                   "monitoreo": monitoring_data, "avance_proyecto": advance_data,
                   "creado_por": st.session_state.user["id"]}
        if preview:
            photo_bytes = [photo.getvalue() for photo in photos]
            photo_bytes += [file.getvalue() for files in goal_evidence_groups for file in files
                            if (file.type or "").startswith("image/")]
            st.session_state.ficha_data = payload
            st.session_state.ficha_photos = photo_bytes
            st.session_state.ficha_pdf = build_pdf(payload, photo_bytes, "assets/logo_coinvierte.jpeg")
            st.session_state.ficha_docx = build_docx(payload, photo_bytes, "assets/logo_coinvierte.jpeg")
            st.rerun()
        if not configured():
            st.success("Proyecto validado correctamente (modo demostración; aún no se guarda en base de datos).")
            st.json(payload)
            return
        try:
            client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
            if project:
                result = client.table("proyectos").update(payload).eq("id", project["id"]).execute()
            else:
                result = client.table("proyectos").insert(payload).execute()
            project_id = str(result.data[0]["id"])
            groups = {"juridica": legal, "auxiliar": auxiliary, "acta_comite": [committee] if committee else [],
                      "acta_junta": [board] if board else [], "convenio": [agreement] if agreement else [],
                      "fotografia": photos}
            for category, files in groups.items():
                upload_files(client, project_id, category, files)
            for files in goal_evidence_groups:
                upload_files(client, project_id, "evidencia_meta", files)
            st.success("Proyecto guardado correctamente.")
        except Exception as exc:
            st.error(f"No fue posible guardar el proyecto: {exc}")

    if st.session_state.get("ficha_data"):
        render_project_preview(st.session_state.ficha_data, st.session_state.get("ficha_photos", []))


def render_project_preview(data: dict, photos: list[bytes]):
    st.markdown("---")
    st.markdown("## Previsualización de la ficha")
    st.caption("Gestión Documental se excluye intencionalmente de esta ficha.")
    st.markdown(f'''<div class="card card-blue">
      <div class="card-icon">FP</div><h3>{data.get("nombre") or "Proyecto sin nombre"}</h3>
      <p class="muted"><b>{data.get("direccion", "")}</b><br>{data.get("solicitante", "")} · {data.get("municipio", "")} · {data.get("anio_inicio", "")}</p>
      <p><b>Monto:</b> ${float(data.get("monto", 0)):,.2f} MXN</p></div>''', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    with c1:
        st.markdown("### Información General")
        st.markdown(f"**Objetivo general**\n\n{data.get('objetivo_general') or 'Sin información'}")
        st.markdown("**Objetivos específicos**")
        for objective in data.get("objetivos_especificos", []):
            st.markdown(f"- {objective}")
    with c2:
        advance = data.get("avance_proyecto", {}) or {}
        if advance:
            st.markdown("### Avance del Proyecto")
            st.progress(int(advance.get("porcentaje_fisico", 0)), text=f"Avance físico: {advance.get('porcentaje_fisico', 0)}%")
            st.markdown(f"**Presupuesto dispersado:** ${float(advance.get('presupuesto_dispersado',0)):,.2f} MXN  \n"
                        f"**Avance financiero:** {advance.get('porcentaje_financiero',0)}%  \n"
                        f"**Semáforo:** {advance.get('semaforo','Sin información')}")
            for index, goal in enumerate(advance.get("metas", []), 1):
                color = {"Por iniciar":"🔴","En progreso":"🟡","Terminada":"🟢"}.get(goal.get("estatus"),"⚪")
                st.markdown(f"**{color} Meta {index}: {goal.get('nombre') or 'Sin nombre'}**  \n"
                            f"{goal.get('descripcion') or 'Sin descripción'}  \n"
                            f"Fecha objetivo: {goal.get('fecha_objetivo') or 'Sin fecha'}")
        else:
            st.markdown("### Monitoreo y Seguimiento")
            monitoring = data.get("monitoreo", {})
            st.progress(int(monitoring.get("avance", 0)), text=f"Avance: {monitoring.get('avance', 0)}%")
            st.markdown(f"**Estatus:** {monitoring.get('estatus', 'Sin información')}  \n"
                        f"**Responsable:** {monitoring.get('responsable') or 'Sin información'}  \n"
                        f"**Periodo:** {monitoring.get('periodo') or 'Sin información'}")
            for label, key in [("Principales avances","avances"),("Pendientes o riesgos","pendientes"),
                               ("Próximas acciones","proximas_acciones"),("Observaciones","observaciones")]:
                if monitoring.get(key):
                    st.markdown(f"**{label}**\n\n{monitoring[key]}")
    st.markdown("### Evidencia Fotográfica")
    if photos:
        columns = st.columns(2)
        for index, photo in enumerate(photos):
            columns[index % 2].image(photo, caption=f"Fotografía {index + 1}", use_container_width=True)
    else:
        st.info("No se cargó evidencia fotográfica para esta ficha.")
    name = "_".join((data.get("nombre") or "proyecto").lower().split())[:60]
    d1, d2, d3 = st.columns([1,1,2])
    d1.download_button("Descargar PDF", st.session_state.ficha_pdf, f"ficha_{name}.pdf", "application/pdf", use_container_width=True)
    d2.download_button("Descargar Word", st.session_state.ficha_docx, f"ficha_{name}.docx",
                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
    if d3.button("Cerrar previsualización", use_container_width=True):
        for key in ["ficha_data", "ficha_photos", "ficha_pdf", "ficha_docx"]:
            st.session_state.pop(key, None)
        st.rerun()


def project_is_active(project: dict) -> bool:
    if project.get("direccion") == "Dirección de Proyectos":
        goals = (project.get("avance_proyecto") or {}).get("metas", [])
        return not goals or any(goal.get("estatus") != "Terminada" for goal in goals)
    return (project.get("monitoreo") or {}).get("estatus") != "Concluido"


def render_readonly_project(project: dict, photos: list[bytes]):
    name = html.escape(project.get("nombre") or "Proyecto sin nombre")
    st.markdown(f'''<div class="card card-blue"><div class="card-icon">VP</div><h3>{name}</h3>
      <p class="muted"><b>{html.escape(project.get("solicitante") or "")}</b><br>
      {html.escape(project.get("municipio") or "")} · {project.get("anio_inicio", "")}</p>
      <p><b>Monto total:</b> ${float(project.get("monto",0)):,.2f} MXN</p></div>''', unsafe_allow_html=True)

    is_projects = project.get("direccion") == "Dirección de Proyectos"
    tab_labels = ["Ficha general", "Avance", "Matriz de riesgos"] if is_projects else ["Ficha general", "Monitoreo y seguimiento"]
    tabs = st.tabs(tab_labels)
    general_tab, progress_tab = tabs[0], tabs[1]
    risks_tab = tabs[2] if is_projects else None
    with general_tab:
        st.markdown("### Información General")
        c1, c2, c3 = st.columns(3)
        c1.metric("Municipio", project.get("municipio") or "Sin información")
        c2.metric("Año de inicio", project.get("anio_inicio") or "Sin información")
        c3.metric("Monto", f'${float(project.get("monto",0)):,.2f}')
        st.markdown("#### Objetivo general")
        st.write(project.get("objetivo_general") or "Sin información")
        st.markdown("#### Objetivos específicos")
        for objective in project.get("objetivos_especificos", []):
            st.markdown(f"- {objective}")
        st.markdown("#### Evidencia fotográfica")
        if photos:
            columns = st.columns(2)
            for index, photo in enumerate(photos):
                columns[index % 2].image(photo, caption=f"Evidencia {index + 1}", use_container_width=True)
        else:
            st.info("Este proyecto todavía no tiene fotografías disponibles.")

    with progress_tab:
        if is_projects:
            advance = project.get("avance_proyecto", {}) or {}
            total_goals = len(advance.get("metas", []))
            st.markdown(f'''<div class="metric-grid">
              <div class="metric-box metric-blue"><div class="metric-label">Avance financiero</div><div class="metric-value">{advance.get("porcentaje_financiero",0)}%</div></div>
              <div class="metric-box metric-green"><div class="metric-label">Metas terminadas</div><div class="metric-value">{advance.get("metas_terminadas",0)}/{total_goals}</div></div>
              <div class="metric-box metric-purple"><div class="metric-label">Avance físico</div><div class="metric-value">{advance.get("porcentaje_fisico",0)}%</div></div>
              <div class="metric-box metric-orange"><div class="metric-label">Semáforo</div><div class="metric-value">{advance.get("semaforo","Sin datos")}</div></div>
            </div>''', unsafe_allow_html=True)
            st.markdown(f"**Presupuesto dispersado:** ${float(advance.get('presupuesto_dispersado',0)):,.2f} MXN")
            _render_financial_progress(float(project.get("monto") or 0),
                                       float(advance.get("presupuesto_dispersado", 0) or 0),
                                       float(advance.get("porcentaje_financiero", 0) or 0))
            if not advance.get("metas"):
                st.info("Todavía no se han registrado metas de ejecución.")
            for index, goal in enumerate(advance.get("metas", []), 1):
                icon = {"Por iniciar":"🔴", "En progreso":"🟡", "Terminada":"🟢"}.get(goal.get("estatus"), "⚪")
                with st.expander(f"{icon} Meta {index}: {goal.get('nombre') or 'Sin nombre'} · {goal.get('estatus','')}", expanded=True):
                    st.write(goal.get("descripcion") or "Sin descripción")
                    st.caption(f"Fecha objetivo: {goal.get('fecha_objetivo') or 'Sin fecha'}")
                    evidence_names = goal.get("evidencias_nombres", [])
                    st.write("**Evidencias:** " + (", ".join(evidence_names) if evidence_names else "Sin evidencia"))
        else:
            monitoring = project.get("monitoreo", {}) or {}
            st.progress(int(monitoring.get("avance", 0)), text=f"Avance: {monitoring.get('avance',0)}%")
            for label, key in [("Estatus","estatus"),("Responsable","responsable"),("Periodo","periodo"),
                               ("Principales avances","avances"),("Pendientes o riesgos","pendientes"),
                               ("Próximas acciones","proximas_acciones"),("Observaciones","observaciones")]:
                st.markdown(f"**{label}:** {monitoring.get(key) or 'Sin información'}")

    if risks_tab:
        with risks_tab:
            risks = normalize_risks((project.get("avance_proyecto") or {}).get("matriz_riesgos", []))
            st.markdown("### Síntesis de riesgos")
            st.info(risk_summary(risks))
            if risks:
                display = pd.DataFrame(risks).drop(columns=["id", "eliminar"], errors="ignore")
                st.dataframe(display, use_container_width=True, hide_index=True)
            else:
                st.info("Todavía no se han registrado riesgos.")

    pdf = build_pdf(project, photos, "assets/logo_coinvierte.jpeg")
    docx = build_docx(project, photos, "assets/logo_coinvierte.jpeg")
    safe_name = "_".join((project.get("nombre") or "proyecto").lower().split())[:60]
    d1, d2 = st.columns(2)
    d1.download_button("Descargar ficha en PDF", pdf, f"ficha_{safe_name}.pdf", "application/pdf", use_container_width=True)
    d2.download_button("Descargar ficha en Word", docx, f"ficha_{safe_name}.docx",
                       "application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)


def view_active_projects(direction: str):
    if not configured():
        st.info("La visualización de proyectos estará disponible al conectar Supabase.")
        return
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
    rows = client.table("proyectos").select("*").eq("direccion", direction).order("updated_at", desc=True).execute().data
    active = [project for project in rows or [] if project_is_active(project)]
    selected_id = st.session_state.get("view_project_id")
    if selected_id:
        project = next((item for item in active if str(item["id"]) == str(selected_id)), None)
        if not project:
            st.session_state.pop("view_project_id", None)
            st.rerun()
        if st.button("← Volver a proyectos activos"):
            st.session_state.pop("view_project_id", None)
            st.session_state.pop("view_project_photos", None)
            st.rerun()
        if "view_project_photos" not in st.session_state:
            st.session_state.view_project_photos = download_project_images(client, str(project["id"]))
        master_delete_control(
            "proyecto", str(project["id"]), f"view_project_{project['id']}",
            lambda: (delete_project_master(client, str(project["id"])),
                     st.session_state.pop("view_project_id", None),
                     st.session_state.pop("view_project_photos", None)),
        )
        render_readonly_project(project, st.session_state.view_project_photos)
        if direction == "Dirección de Proyectos":
            project_financial_documents(project)
        return

    st.markdown("## Proyectos activos")
    st.caption(f"{direction} · Selecciona un proyecto para consultar su información")
    if not active:
        st.info("No hay proyectos activos registrados en esta dirección.")
        return
    labels = {f"{p['nombre']} — {p['municipio']} ({p['anio_inicio']})": p for p in active}
    selected_label = st.selectbox("Proyecto", list(labels.keys()))
    selected = labels[selected_label]
    advance = selected.get("avance_proyecto", {}) or {}
    st.markdown(f'''<div class="card card-green"><div class="card-icon">{len(active)}</div>
      <h3>{html.escape(selected.get("nombre") or "")}</h3><p class="muted">{html.escape(selected.get("solicitante") or "")} · 
      {html.escape(selected.get("municipio") or "")}</p><p><b>Avance:</b> {advance.get("porcentaje_fisico", (selected.get("monitoreo") or {}).get("avance",0))}%</p></div>''', unsafe_allow_html=True)
    if st.button("Ver información del proyecto", type="primary", use_container_width=True):
        st.session_state.view_project_id = str(selected["id"])
        st.session_state.pop("view_project_photos", None)
        st.rerun()


def user_management():
    st.title("Gestión de usuarios")
    if st.session_state.user.get("rol") != "administrador":
        st.error("No tienes permisos para acceder a este módulo.")
        return
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
    create_tab, users_tab = st.tabs(["Generar código temporal", "Usuarios autorizados"])
    with create_tab:
        st.markdown("### Autorizar a una persona")
        with st.form("create_access_code"):
            name = st.text_input("Nombre de la persona")
            email = st.text_input("Correo institucional", placeholder="nombre@jalisco.gob.mx")
            direction = st.selectbox("Dirección de adscripción", USER_DIRECTIONS)
            modules = st.multiselect("Módulos visibles en el inicio", ALL_MODULES, default=[])
            project_directions = st.multiselect("Áreas visibles dentro de Proyectos", PROJECT_DIRECTIONS,
                                                help="Sólo se aplican cuando se concede acceso a Programas / Proyectos.")
            hours = st.selectbox("Vigencia del código", [24, 48, 72, 168],
                                 format_func=lambda value: "7 días" if value == 168 else f"{value} horas")
            create_code = st.form_submit_button("Generar código de acceso", type="primary", use_container_width=True)
        if create_code:
            if not valid_official_email(email):
                st.error("Sólo se pueden autorizar correos @jalisco.gob.mx.")
            elif not modules:
                st.error("Selecciona al menos un módulo visible.")
            elif MODULE_PROJECTS in modules and not project_directions:
                st.error("Selecciona Operaciones, Dirección de Proyectos o ambas.")
            else:
                try:
                    result = client.rpc("crear_codigo_acceso", {"p_email": email.lower().strip(),
                        "p_nombre": name.strip(), "p_horas": hours, "p_direccion": direction,
                        "p_modulos": modules, "p_direcciones_proyectos": project_directions}).execute().data
                    record = result[0] if isinstance(result, list) else result
                    st.session_state.generated_code = record
                    st.session_state.generated_email = email.lower().strip()
                except Exception as exc:
                    st.error(f"No fue posible generar el código: {exc}")
        if st.session_state.get("generated_code"):
            record = st.session_state.generated_code
            st.success(f"Código generado para {st.session_state.generated_email}")
            st.code(record.get("codigo", ""), language=None)
            st.caption(f"Vence: {record.get('vence', '')}. Compártelo únicamente con la persona autorizada.")

    with users_tab:
        rows = client.table("usuarios_autorizados").select("id,email,nombre,rol,activo,direccion,modulos,direcciones_proyectos,ultimo_acceso,created_at").order("created_at", desc=True).execute().data
        if not rows:
            st.info("Todavía no hay usuarios registrados.")
        else:
            st.dataframe([{ "Nombre": row.get("nombre"), "Correo": row.get("email"), "Dirección": row.get("direccion") or "Sin asignar",
                            "Módulos": ", ".join(row.get("modulos") or []) or "Ninguno",
                            "Proyectos": ", ".join(row.get("direcciones_proyectos") or []) or "Sin acceso",
                            "Rol": row.get("rol"),
                            "Estado": "Activo" if row.get("activo") else "Suspendido / pendiente",
                            "Último acceso": row.get("ultimo_acceso") or "Sin acceso"} for row in rows],
                         use_container_width=True, hide_index=True)
            manageable = [row for row in rows if row.get("rol") != "administrador"]
            if manageable:
                labels = {f"{row.get('nombre') or 'Sin nombre'} — {row['email']}": row for row in manageable}
                selected_label = st.selectbox("Administrar usuario", list(labels.keys()))
                selected = labels[selected_label]
                with st.form(f"permissions_{selected['id']}"):
                    edit_direction = st.selectbox("Dirección de adscripción", USER_DIRECTIONS,
                        index=USER_DIRECTIONS.index(selected.get("direccion")) if selected.get("direccion") in USER_DIRECTIONS else 0)
                    edit_modules = st.multiselect("Módulos visibles", ALL_MODULES, default=selected.get("modulos") or [])
                    edit_project_directions = st.multiselect("Áreas visibles dentro de Proyectos", PROJECT_DIRECTIONS,
                        default=selected.get("direcciones_proyectos") or [])
                    save_permissions = st.form_submit_button("Guardar dirección y permisos", type="primary", use_container_width=True)
                if save_permissions:
                    if not edit_modules:
                        st.error("Selecciona al menos un módulo.")
                    elif MODULE_PROJECTS in edit_modules and not edit_project_directions:
                        st.error("Selecciona al menos un área de Proyectos.")
                    else:
                        client.table("usuarios_autorizados").update({"direccion": edit_direction, "modulos": edit_modules,
                            "direcciones_proyectos": edit_project_directions, "updated_at": datetime.now().isoformat()}).eq("id", selected["id"]).execute()
                        st.success("Dirección y permisos actualizados.")
                        st.rerun()
                c1, c2 = st.columns(2)
                if selected.get("activo"):
                    if c1.button("Suspender acceso", use_container_width=True):
                        client.table("usuarios_autorizados").update({"activo": False}).eq("id", selected["id"]).execute()
                        st.success("Acceso suspendido.")
                        st.rerun()
                else:
                    if c1.button("Reactivar acceso", use_container_width=True):
                        client.table("usuarios_autorizados").update({"activo": True}).eq("id", selected["id"]).execute()
                        st.success("Acceso reactivado.")
                        st.rerun()
                if c2.button("Generar nuevo código", use_container_width=True):
                    result = client.rpc("crear_codigo_acceso", {"p_email": selected["email"],
                        "p_nombre": selected.get("nombre") or "", "p_horas": 24,
                        "p_direccion": selected.get("direccion") or USER_DIRECTIONS[0],
                        "p_modulos": selected.get("modulos") or [],
                        "p_direcciones_proyectos": selected.get("direcciones_proyectos") or []}).execute().data
                    st.session_state.generated_code = result[0] if isinstance(result, list) else result
                    st.session_state.generated_email = selected["email"]
                    st.success("Nuevo código generado. Consúltalo en la primera pestaña.")
                st.divider()
                confirm_remove = st.checkbox(f"Confirmo que deseo remover el acceso de {selected.get('nombre') or selected['email']}",
                                             key=f"confirm_remove_{selected['id']}")
                if st.button("Remover usuario", disabled=not confirm_remove, key=f"remove_user_{selected['id']}"):
                    client.rpc("remover_usuario_autorizado", {"p_usuario_id": selected["id"]}).execute()
                    st.success("El acceso fue removido. Sus aportaciones y documentos históricos se conservaron.")
                    st.rerun()


BOARD_YEARS = [2025, 2026, 2027, 2028, 2029, 2030]
PRESET_BOARD_SESSIONS = ["Primera (1era)", "Segunda (2da)", "Tercera (3ra)"]
COMMITTEE_CATALOG = [
    ("Comité de Ética", "Principios, integridad y conducta institucional.", "var(--blue)"),
    ("Comité de Igualdad de Género", "Igualdad sustantiva y transversalización institucional.", "var(--purple)"),
    ("Comité de Archivo", "Gestión documental, conservación y cumplimiento archivístico.", "var(--orange)"),
    ("Comité de Control Interno", "Control institucional, riesgos y mejora continua.", "var(--green)"),
]
BOARD_AREAS = ["Dirección Jurídica", "Dirección General", "Dirección de Administración", "Dirección de Operaciones", "Dirección de Planeación", "Órgano Interno de Control"]


def board_year_label(year: int) -> str:
    return str(year)


def board_year_selector():
    logo_header()
    st.markdown('<h1 class="choice-title">Junta de Gobierno</h1>', unsafe_allow_html=True)
    st.markdown('<p class="choice-subtitle">Selecciona el año que deseas consultar</p>', unsafe_allow_html=True)
    colors = ["var(--blue)", "var(--green)", "var(--teal)", "var(--purple)", "var(--orange)", "var(--gray)"]
    for start in (0, 3):
        columns = st.columns(3, gap="large")
        for column, year, color in zip(columns, BOARD_YEARS[start:start + 3], colors[start:start + 3]):
            with column:
                label = board_year_label(year)
                st.markdown(f'<div class="year-card" style="--accent:{color}"><h2>{label}</h2><p>Sesiones y acuerdos</p></div>', unsafe_allow_html=True)
                if st.button(f"Abrir {label}", key=f"board_year_{year}", use_container_width=True, type="primary"):
                    st.session_state.board_year = year
                    st.session_state.pop("board_session", None)
                    st.rerun()
    st.markdown('''<div class="analytics-banner"><h3>Analítica de datos</h3>
        <p>Numeralia de sesiones y acuerdos, así como distribución del avance por áreas responsables.</p></div>''', unsafe_allow_html=True)
    if st.button("Abrir Analítica de datos", key="open_board_analytics", use_container_width=True, type="primary"):
        st.session_state.board_analytics = True
        st.session_state.pop("board_year", None)
        st.session_state.pop("board_session", None)
        st.rerun()


def board_analytics_dashboard():
    top1, top2 = st.columns([1, 5])
    if top1.button("← Junta", use_container_width=True):
        st.session_state.pop("board_analytics", None)
        st.rerun()
    top2.markdown("## Analítica de datos · Junta de Gobierno")
    scope = st.selectbox("Periodo de análisis", ["Todos los años"] + [str(year) for year in BOARD_YEARS])
    selected_year = None if scope == "Todos los años" else int(scope)
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token) if configured() else None
    if not client:
        st.error("Primero debes conectar Supabase.")
        return
    try:
        session_query = client.table("sesiones_junta").select("id,anio,fecha_sesion")
        if selected_year: session_query = session_query.eq("anio", selected_year)
        sessions = session_query.execute().data or []
        session_ids = [row["id"] for row in sessions]
        agreements = []
        if session_ids:
            agreements = (client.table("acuerdos_junta").select("id,sesion_id,tipo_registro,estatus,resultado,areas")
                          .in_("sesion_id", session_ids).execute().data or [])
    except Exception as exc:
        st.error(f"No fue posible construir la analítica. Ejecuta la migración nueva y vuelve a intentar: {exc}")
        return
    agreements = [row for row in agreements if row.get("tipo_registro") == "Acuerdo"]
    celebrated = sum(1 for row in sessions if row.get("fecha_sesion") and str(row["fecha_sesion"])[:10] <= date.today().isoformat())
    counts = {
        "Juntas de Gobierno celebradas": celebrated,
        "Acuerdos aprobados": sum(row.get("resultado") == "Aprobado" for row in agreements),
        "Acuerdos rechazados": sum(row.get("resultado") == "Rechazado" for row in agreements),
        "Acuerdos por iniciar": sum(row.get("estatus") == "Por iniciar" for row in agreements),
        "Acuerdos en progreso": sum(row.get("estatus") == "En proceso" for row in agreements),
        "Acuerdos terminados": sum(row.get("estatus") == "Terminada" for row in agreements),
    }
    tones = ["#0798cf", "#009b4c", "#b85c62", "#f68b08", "#c5a44a", "#16ad8f"]
    cards = "".join(f'''<div class="analytics-metric" style="--tone:{tone}"><div class="analytics-value">{value}</div>
        <div class="analytics-label">{html.escape(label)}</div></div>''' for (label, value), tone in zip(counts.items(), tones))
    st.markdown("### Numeralia")
    st.caption("Cifras acumuladas de 2025–2030." if selected_year is None else f"Cifras correspondientes a {selected_year}.")
    st.markdown(f'<div class="analytics-metrics">{cards}</div>', unsafe_allow_html=True)
    pending_result = sum(row.get("resultado") in (None, "Pendiente") for row in agreements)
    if pending_result:
        st.info(f"Hay {pending_result} acuerdo(s) cuyo resultado todavía debe clasificarse como aprobado o rechazado.")
    st.markdown("### Avance por áreas responsables")
    st.caption("Cantidad de acuerdos asignados a cada dirección. Se excluyen Dirección General y Órgano Interno de Control.")
    included_areas = [area for area in BOARD_AREAS if area not in ("Dirección General", "Órgano Interno de Control")]
    area_data = [{"Dirección": area, "Acuerdos": sum(area in (row.get("areas") or []) for row in agreements)} for area in included_areas]
    area_df = pd.DataFrame(area_data)
    chart = (alt.Chart(area_df).mark_bar(cornerRadiusTopLeft=7, cornerRadiusTopRight=7, size=54)
             .encode(x=alt.X("Dirección:N", sort=included_areas, title=None, axis=alt.Axis(labelAngle=0, labelLimit=260)),
                     y=alt.Y("Acuerdos:Q", title="Cantidad de acuerdos", axis=alt.Axis(tickMinStep=1)),
                     color=alt.Color("Dirección:N", scale=alt.Scale(domain=included_areas,
                         range=["#0798cf", "#009b4c", "#a990c7"]), legend=None),
                     tooltip=[alt.Tooltip("Dirección:N"), alt.Tooltip("Acuerdos:Q", format=".0f")]))
    labels = alt.Chart(area_df).mark_text(dy=-12, fontSize=15, fontWeight="bold", color="#35434b").encode(
        x=alt.X("Dirección:N", sort=included_areas), y="Acuerdos:Q", text=alt.Text("Acuerdos:Q", format=".0f"))
    st.altair_chart((chart + labels).properties(height=390), use_container_width=True)
def _extract_board_text(uploaded) -> str:
    data = uploaded.getvalue()
    suffix = Path(uploaded.name).suffix.lower()
    if suffix == ".docx":
        doc = Document(io.BytesIO(data))
        blocks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        blocks += [" ".join(c.text.strip() for c in row.cells if c.text.strip()) for table in doc.tables for row in table.rows]
        return "\n".join(blocks)
    if suffix == ".pptx":
        prs = Presentation(io.BytesIO(data))
        return "\n".join(shape.text.strip() for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
    if suffix == ".pdf":
        text = "\n".join((page.extract_text() or "").strip() for page in PdfReader(io.BytesIO(data)).pages).strip()
        if len(text) >= 80:
            return text
        pdf = fitz.open(stream=data, filetype="pdf")
        pages = []
        for page in pdf:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = Image.open(io.BytesIO(pix.tobytes("png")))
            pages.append(pytesseract.image_to_string(image, lang="spa"))
        return "\n".join(pages).strip()
    raise ValueError("Formato no compatible.")


def _board_items_from_text(text: str) -> list[dict]:
    lines = [re.sub(r"\s+", " ", line).strip(" •\t") for line in re.sub(r"\r", "", text or "").split("\n")]
    lines = [line for line in lines if line]
    start = next((i for i, line in enumerate(lines) if re.search(r"aprobaci[oó]n.{0,20}orden.{0,12}d[ií]a", line, re.I)), -1)
    if start < 0:
        first_substantive = next((i for i, line in enumerate(lines) if re.search(r"presentaci[oó]n|\binforme\b", line, re.I)), -1)
        start = first_substantive - 1 if first_substantive >= 0 else -1
    end = next((i for i, line in enumerate(lines) if i > start and re.search(r"\basuntos\s+varios\b", line, re.I)), len(lines))
    body, items, current = (lines[start + 1:end] if start >= 0 else []), [], ""
    for line in body:
        marker = re.match(r"^(?:punto\s+)?(?:\d+[.)-]|[IVXLCDM]+[.)-])\s*(.+)$", line, re.I)
        begins = bool(marker or re.search(r"\b(presentaci[oó]n|informe)\b", line, re.I))
        if begins:
            if current: items.append(current.strip())
            current = marker.group(1).strip() if marker else line
        elif current:
            current += " " + line
    if current: items.append(current.strip())
    results = []
    for item in items:
        if re.search(r"(lista de asistencia|declaraci[oó]n de qu[oó]rum|lectura del acta|clausura(?:\s+de)?\s+la\s+sesi[oó]n|asuntos\s+varios)", item, re.I): continue
        is_report = bool(re.search(r"\binforme\b", item, re.I)) and not bool(re.search(r"aprobaci[oó]n", item, re.I))
        results.append({"tipo_registro": "Informe" if is_report else "Acuerdo", "titulo": item, "texto": item, "Eliminar": False})
    return results


def _session_number(name: str) -> int:
    match = re.search(r"(\d+)\s*(?:era|da|ra|ta|ma|va|na)?", name or "", re.I)
    if match: return int(match.group(1))
    words = {"primera": 1, "segunda": 2, "tercera": 3, "cuarta": 4, "quinta": 5, "sexta": 6,
             "séptima": 7, "septima": 7, "octava": 8, "novena": 9, "décima": 10, "decima": 10}
    low = (name or "").lower()
    return next((number for word, number in words.items() if word in low), 0)


def _agreement_code(session: dict, consecutive: int) -> str:
    kind = "ORD" if session.get("tipo") == "Ordinaria" else "EXT"
    return f"JG-{int(session.get('anio'))}-{kind}-{_session_number(session.get('nombre')):02d}-{consecutive:03d}"


def _next_agreement_number(session: dict, rows: list[dict]) -> str:
    consecutives = []
    for row in rows:
        match = re.search(r"-(\d{3})$", row.get("numero") or "")
        if match:
            consecutives.append(int(match.group(1)))
    return _agreement_code(session, max(consecutives, default=0) + 1)


def _deadline_label(value, status: str) -> str:
    if not value: return "Sin fecha compromiso"
    days = (date.fromisoformat(str(value)[:10]) - date.today()).days
    if status != "Terminada" and days < 0: return f"Vencida hace {abs(days)} día(s)"
    if days == 0: return "Vence hoy"
    if days > 0: return f"Faltan {days} día(s)"
    return "Terminada"


def _pdf_preview(data: bytes, height: int = 650):
    """Muestra el PDF como páginas renderizadas sin depender de st.pdf.

    El componente PDF nativo puede fallar en Streamlit Cloud por diferencias
    entre la versión del frontend y el complemento instalado. PyMuPDF ya forma
    parte del proyecto y produce una vista previa estable en todos los casos.
    """
    try:
        document = fitz.open(stream=data, filetype="pdf")
        if document.page_count == 0:
            st.info("El PDF no contiene páginas para previsualizar.")
            return
        for page_number, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.35, 1.35), alpha=False)
            st.image(
                pixmap.tobytes("png"),
                caption=f"Página {page_number} de {document.page_count}",
                use_container_width=True,
            )
        document.close()
    except Exception:
        st.info("No fue posible generar la vista previa de este PDF. Puedes descargarlo para consultarlo.")


def _upload_junta_document(client, path: str, uploaded) -> None:
    """Usa carga reanudable por bloques para documentos mayores a 6 MB."""
    data = uploaded.getvalue()
    mime_type = uploaded.type or "application/octet-stream"
    if len(data) <= 6 * 1024 * 1024:
        client.storage.from_("expedientes").upload(path, data, {"content-type": mime_type})
        return
    endpoint = f"{st.secrets['SUPABASE_URL'].rstrip('/')}/storage/v1/upload/resumable"
    encode_meta = lambda value: base64.b64encode(str(value).encode("utf-8")).decode("ascii")
    headers = {
        "Authorization": f"Bearer {st.session_state.access_token}", "apikey": public_key(),
        "tus-resumable": "1.0.0", "upload-length": str(len(data)), "x-upsert": "false",
        "upload-metadata": ",".join([
            f"bucketName {encode_meta('expedientes')}", f"objectName {encode_meta(path)}",
            f"contentType {encode_meta(mime_type)}", f"cacheControl {encode_meta('3600')}",
        ]),
    }
    response = requests.post(endpoint, headers=headers, timeout=30)
    response.raise_for_status()
    upload_url = response.headers.get("Location") or response.headers.get("location")
    if not upload_url:
        raise RuntimeError("No se recibió la ubicación de la carga.")
    if upload_url.startswith("/"):
        upload_url = f"{st.secrets['SUPABASE_URL'].rstrip('/')}{upload_url}"
    offset, chunk_size = 0, 6 * 1024 * 1024
    while offset < len(data):
        chunk = data[offset:offset + chunk_size]
        chunk_headers = {
            "Authorization": f"Bearer {st.session_state.access_token}", "apikey": public_key(),
            "tus-resumable": "1.0.0", "upload-offset": str(offset),
            "content-type": "application/offset+octet-stream",
        }
        result = requests.patch(upload_url, headers=chunk_headers, data=chunk, timeout=120)
        result.raise_for_status()
        offset = int(result.headers.get("Upload-Offset") or result.headers.get("upload-offset") or offset + len(chunk))


@st.cache_data(show_spinner=False, ttl=3600)
def _office_to_pdf(data: bytes, filename: str) -> tuple[bytes | None, str]:
    """Convierte Office a PDF con un perfil aislado y devuelve diagnóstico."""
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        return None, "LibreOffice Impress no está instalado en el servidor."
    try:
        with tempfile.TemporaryDirectory() as folder:
            work = Path(folder)
            source = work / (safe_name(Path(filename).stem) + Path(filename).suffix.lower())
            output_dir = work / "output"
            profile_dir = work / "profile"
            output_dir.mkdir()
            profile_dir.mkdir()
            source.write_bytes(data)
            command = [
                executable, f"-env:UserInstallation={profile_dir.as_uri()}",
                "--headless", "--nologo", "--nodefault", "--nofirststartwizard",
                "--convert-to", "pdf:impress_pdf_Export", "--outdir", str(output_dir), str(source),
            ]
            result = subprocess.run(command, capture_output=True, text=True, timeout=300, check=False)
            converted = output_dir / f"{source.stem}.pdf"
            if result.returncode != 0 or not converted.exists():
                detail = (result.stderr or result.stdout or "La conversión no produjo un PDF.").strip()
                return None, detail[-500:]
            return converted.read_bytes(), ""
    except subprocess.TimeoutExpired:
        return None, "La conversión excedió cinco minutos."
    except Exception as exc:
        return None, str(exc)


def _document_preview(data: bytes, filename: str, height: int = 650) -> bool:
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        _pdf_preview(data, height)
        return True
    if suffix not in (".docx", ".pptx"):
        return False
    converted, diagnostic = _office_to_pdf(data, filename)
    if converted:
        _pdf_preview(converted, height)
        return True
    if suffix == ".pptx":
        st.warning("No fue posible generar las imágenes de las diapositivas.")
        st.caption(f"Diagnóstico: {diagnostic}")
    return False


def _saved_main_document(client, document: dict, column, key_prefix: str) -> None:
    """Acciones visibles para un documento principal que ya está guardado."""
    filename = document.get("nombre_archivo") or "documento"
    label = document.get("nombre_visible") or Path(filename).stem
    column.success(f"Guardado: {label}")
    try:
        data = client.storage.from_("expedientes").download(document["ruta_storage"])
        view_col, download_col = column.columns(2)
        show = view_col.toggle("Previsualizar", key=f"{key_prefix}_view_{document['id']}")
        download_col.download_button(
            "Descargar", data, file_name=filename,
            mime=document.get("mime_type") or "application/octet-stream",
            key=f"{key_prefix}_download_{document['id']}", use_container_width=True,
        )
        if show:
            with column.container(border=True):
                if not _document_preview(data, filename, 520):
                    st.info("No fue posible generar la vista previa, pero el archivo sí puede descargarse.")
    except Exception:
        column.error("El registro existe, pero el archivo no pudo recuperarse del almacenamiento.")


def _document_card(client, document: dict, key_prefix: str, table_name: str):
    """Tarjeta compacta con metadatos, descarga y vista previa."""
    filename = document.get("nombre_archivo") or "documento"
    suffix = Path(filename).suffix.upper().lstrip(".") or "ARCHIVO"
    created = str(document.get("created_at") or "")[:16].replace("T", " · ")
    with st.container(border=True):
        info, actions = st.columns([4.5, 2.5], vertical_alignment="center")
        info.markdown(f"**{html.escape(document.get('nombre_visible') or Path(filename).stem)}**")
        info.caption(
            f"{suffix} · {html.escape(document.get('tipo_documento') or 'Documento')}  \n"
            f"Subido por {html.escape(document.get('autor_nombre') or 'Usuario')} · {created or 'Fecha no disponible'}"
        )
        try:
            data = client.storage.from_("expedientes").download(document["ruta_storage"])
            action_columns = actions.columns(3 if is_master_admin() else 2)
            view_col, download_col = action_columns[0], action_columns[1]
            delete_col = action_columns[2] if is_master_admin() else None
            show_preview = view_col.toggle("Ver", key=f"{key_prefix}_view_{document['id']}")
            download_col.download_button(
                "Descargar", data, file_name=filename,
                mime=document.get("mime_type") or "application/octet-stream",
                key=f"{key_prefix}_download_{document['id']}", use_container_width=True,
            )
            pending_key = f"{key_prefix}_pending_delete_{document['id']}"
            if delete_col and delete_col.button("Eliminar", key=f"{key_prefix}_delete_{document['id']}", use_container_width=True):
                st.session_state[pending_key] = True
                st.rerun()
            if is_master_admin() and st.session_state.get(pending_key):
                st.warning(f"¿Eliminar definitivamente “{document.get('nombre_visible') or filename}” del expediente?")
                confirm_col, cancel_col, _ = st.columns([1, 1, 4])
                if confirm_col.button("Sí, eliminar", key=f"{key_prefix}_confirm_{document['id']}", type="primary"):
                    client.storage.from_("expedientes").remove([document["ruta_storage"]])
                    client.table(table_name).delete().eq("id", document["id"]).execute()
                    if table_name == "documentos_sesion_junta" and document.get("tipo_documento") == "Acta firmada":
                        client.table("sesiones_junta").update({"acta_firmada_nombre": None, "acta_firmada_ruta": None}).eq("id", document.get("sesion_id")).execute()
                        if st.session_state.get("board_session", {}).get("id") == document.get("sesion_id"):
                            st.session_state.board_session.update({"acta_firmada_nombre": None, "acta_firmada_ruta": None})
                    st.session_state.pop(pending_key, None)
                    st.success("Documento eliminado.")
                    st.rerun()
                if cancel_col.button("Cancelar", key=f"{key_prefix}_cancel_{document['id']}"):
                    st.session_state.pop(pending_key, None)
                    st.rerun()
            if show_preview:
                st.markdown("##### Vista previa")
                if not _document_preview(data, filename, 650):
                    st.info("La vista previa no está disponible para este formato.")
        except Exception:
            actions.caption("No fue posible abrir el archivo.")


def _agreement_ficha_docx(agreement: dict, session: dict, comments: list, history: list, files: list) -> bytes:
    doc = Document(); section = doc.sections[0]
    section.top_margin = section.bottom_margin = section.left_margin = section.right_margin = Inches(.85)
    normal = doc.styles["Normal"]; normal.font.name = "Arial"; normal.font.size = Pt(10.5)
    normal.paragraph_format.space_after = Pt(6)
    logo = Path("assets/logo_coinvierte.jpeg")
    if logo.exists():
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(str(logo), width=Inches(5.8))
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"{session.get('nombre','')} SESIÓN {session.get('tipo','').upper()}"); r.bold = True; r.font.size = Pt(16); r.font.color.rgb = RGBColor(53,67,75)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(f"Junta de Gobierno · {session.get('fecha_sesion') or session.get('anio','Fecha no registrada')}").italic = True
    p = doc.add_paragraph(); r = p.add_run(f"FICHA DE {agreement.get('titulo') or 'Punto de acuerdo'}"); r.bold = True; r.font.size = Pt(13); r.font.color.rgb = RGBColor(103,80,164)
    def add_rows(title, rows):
        h = doc.add_paragraph(); rr = h.add_run(title); rr.bold = True; rr.font.size = Pt(13); rr.font.color.rgb = RGBColor(103,80,164)
        table = doc.add_table(rows=0, cols=2); table.style = "Table Grid"
        for label, value in rows:
            cells = table.add_row().cells; cells[0].text = str(label); cells[1].text = str(value or "Sin información")
            cells[0].paragraphs[0].runs[0].bold = True
    add_rows("Datos del acuerdo", [("Código",agreement.get("numero")),("Tipo",agreement.get("tipo_registro")),("Texto",agreement.get("texto") or agreement.get("titulo"))])
    add_rows("Seguimiento", [("Responsables",", ".join(agreement.get("areas") or []) or "Sin responsable"),
        ("Estatus","En progreso" if agreement.get("estatus")=="En proceso" else agreement.get("estatus")),
        ("Fecha compromiso",agreement.get("fecha_compromiso")),("Fecha de cierre",agreement.get("fecha_cierre")),("Cumplimiento",agreement.get("cumplimiento"))])
    for title, records, field in [("Historial de seguimiento",history,"descripcion"),("Comentarios",comments,"comentario")]:
        h=doc.add_paragraph(); rr=h.add_run(title); rr.bold=True; rr.font.size=Pt(13); rr.font.color.rgb=RGBColor(7,152,207)
        if not records: doc.add_paragraph("Sin registros.")
        for item in records: doc.add_paragraph(f"{str(item.get('created_at') or '')[:16]} · {item.get('autor_nombre') or 'Usuario'}: {item.get(field) or ''}")
    h=doc.add_paragraph(); rr=h.add_run("Archivos relacionados"); rr.bold=True; rr.font.size=Pt(13); rr.font.color.rgb=RGBColor(7,152,207)
    doc.add_paragraph(", ".join(f.get("nombre_archivo","") for f in files) if files else "Sin archivos adjuntos.")
    output=io.BytesIO(); doc.save(output); return output.getvalue()


def _agreement_ficha_pdf(agreement: dict, session: dict, comments: list, history: list, files: list) -> bytes:
    output=io.BytesIO(); doc=SimpleDocTemplate(output,pagesize=letter,leftMargin=.7*inch,rightMargin=.7*inch,topMargin=.55*inch,bottomMargin=.55*inch)
    styles=getSampleStyleSheet(); body=ParagraphStyle("fb",parent=styles["BodyText"],fontSize=9.2,leading=12,textColor=colors.HexColor("#35434B"))
    heading=ParagraphStyle("fh",parent=body,fontName="Helvetica-Bold",fontSize=13,textColor=colors.HexColor("#6750A4"),spaceBefore=10,spaceAfter=6)
    title=ParagraphStyle("ft",parent=body,fontName="Helvetica-Bold",fontSize=15,leading=18,alignment=TA_CENTER,spaceAfter=4)
    story=[]; logo=Path("assets/logo_coinvierte.jpeg")
    if logo.exists(): story += [RLImage(str(logo),width=5.6*inch,height=1.08*inch),Spacer(1,6)]
    story += [Paragraph(html.escape(f"{session.get('nombre','')} SESIÓN {session.get('tipo','').upper()}"),title),
              Paragraph(html.escape(f"Junta de Gobierno · {session.get('fecha_sesion') or session.get('anio','Fecha no registrada')}"),ParagraphStyle("fc",parent=body,alignment=TA_CENTER,spaceAfter=12)),
              Paragraph(html.escape(f"FICHA DE {agreement.get('titulo') or 'Punto de acuerdo'}"),heading)]
    def add_table(rows):
        data=[[Paragraph(html.escape(str(a)),body),Paragraph(html.escape(str(b or "Sin información")),body)] for a,b in rows]
        t=Table(data,colWidths=[1.55*inch,5.25*inch]); t.setStyle(TableStyle([("GRID",(0,0),(-1,-1),.3,colors.HexColor("#DDE4E6")),("BACKGROUND",(0,0),(0,-1),colors.HexColor("#F1EFF8")),("VALIGN",(0,0),(-1,-1),"TOP"),("LEFTPADDING",(0,0),(-1,-1),6),("RIGHTPADDING",(0,0),(-1,-1),6),("TOPPADDING",(0,0),(-1,-1),6),("BOTTOMPADDING",(0,0),(-1,-1),6)])); story.append(t)
    add_table([("Código",agreement.get("numero")),("Tipo",agreement.get("tipo_registro")),("Texto",agreement.get("texto") or agreement.get("titulo"))])
    story.append(Paragraph("Seguimiento",heading)); add_table([("Responsables",", ".join(agreement.get("areas") or []) or "Sin responsable"),("Estatus","En progreso" if agreement.get("estatus")=="En proceso" else agreement.get("estatus")),("Fecha compromiso",agreement.get("fecha_compromiso")),("Fecha de cierre",agreement.get("fecha_cierre")),("Cumplimiento",agreement.get("cumplimiento"))])
    for title_text,records,field in [("Historial de seguimiento",history,"descripcion"),("Comentarios",comments,"comentario")]:
        story.append(Paragraph(title_text,heading))
        if not records: story.append(Paragraph("Sin registros.",body))
        for item in records: story.append(Paragraph(html.escape(f"{str(item.get('created_at') or '')[:16]} · {item.get('autor_nombre') or 'Usuario'}: {item.get(field) or ''}"),body))
    story.append(Paragraph("Archivos relacionados",heading)); story.append(Paragraph(html.escape(", ".join(f.get("nombre_archivo","") for f in files) if files else "Sin archivos adjuntos."),body))
    doc.build(story); return output.getvalue()




def _active_notification_users(client, required_module: str) -> list[dict]:
    """Usuarios activos que pueden ser responsables según su nivel de acceso."""
    try:
        rows = (
            client.table("usuarios_autorizados")
            .select("id,nombre,email,activo,rol,direccion,modulos")
            .eq("activo", True)
            .order("nombre")
            .execute()
            .data
            or []
        )
    except Exception:
        return []

    allowed = []
    for row in rows:
        email = str(row.get("email") or "").strip()
        if not email:
            continue
        modules = row.get("modulos") or []
        if row.get("rol") == "administrador" or required_module in modules:
            allowed.append(row)
    return allowed


def _legacy_responsibles_from_row(row: dict) -> list[dict]:
    current = row.get("responsables_notificacion")
    if isinstance(current, list):
        clean = []
        for item in current:
            if not isinstance(item, dict):
                continue
            email = str(item.get("email") or "").strip()
            if email:
                clean.append({
                    "id": item.get("id"),
                    "nombre": item.get("nombre") or email,
                    "email": email,
                    "direccion": item.get("direccion"),
                })
        if clean:
            return clean

    # Compatibilidad con la versión anterior de responsable único.
    if row.get("responsable_email"):
        return [{
            "id": row.get("responsable_usuario_id"),
            "nombre": row.get("responsable_nombre") or row.get("responsable_email"),
            "email": row.get("responsable_email"),
            "direccion": None,
        }]
    return []


def _responsibles_selector(
    client,
    *,
    required_module: str,
    current_responsibles: list[dict] | None,
    key: str,
):
    users = _active_notification_users(client, required_module)
    by_id = {str(row["id"]): row for row in users}

    current_ids = []
    for item in current_responsibles or []:
        item_id = str(item.get("id") or "")
        if item_id in by_id:
            current_ids.append(item_id)

    def _label(value):
        row = by_id.get(value, {})
        name = row.get("nombre") or row.get("email") or "Usuario"
        email = row.get("email") or ""
        direction = row.get("direccion") or "Sin dirección"
        return f"{name} · {direction} · {email}"

    selected_ids = st.multiselect(
        "Personas responsables",
        options=list(by_id.keys()),
        default=current_ids,
        format_func=_label,
        key=key,
        help=(
            f"Sólo aparecen usuarios activos con acceso a {required_module}. "
            "Cada persona seleccionada recibirá su propio recordatorio por correo."
        ),
    )

    selected = []
    for user_id in selected_ids:
        row = by_id[user_id]
        selected.append({
            "id": row.get("id"),
            "nombre": row.get("nombre") or row.get("email"),
            "email": row.get("email"),
            "direccion": row.get("direccion"),
        })
    return selected


def _notification_status_caption(deadline, responsibles, enabled, status):
    if status == "Terminada":
        return "El acuerdo está terminado; no se enviarán recordatorios."
    if not enabled:
        return "Recordatorio por correo desactivado."
    if not deadline:
        return "Falta definir la fecha compromiso para programar el aviso."
    if not responsibles:
        return "Falta asignar al menos una persona responsable con correo."
    try:
        days = (deadline - date.today()).days
    except Exception:
        days = None

    recipients = len(responsibles)
    suffix = f" · {recipients} destinatario(s)"
    if days is None:
        return "Recordatorio configurado para 3 días antes del vencimiento" + suffix
    if days > 3:
        return f"Correo programado para 3 días antes del vencimiento · faltan {days} días" + suffix
    if days == 3:
        return "El correo corresponde enviarse hoy" + suffix
    if 0 <= days < 3:
        return f"Faltan {days} día(s). El aviso de 3 días ya debió generarse" + suffix
    return "La fecha compromiso ya venció" + suffix


def _drive_video_preview_url(url: str) -> str | None:
    url = str(url or "").strip()
    patterns = [
        r"drive\.google\.com/file/d/([^/]+)",
        r"drive\.google\.com/open\?id=([^&]+)",
        r"drive\.google\.com/uc\?.*?[?&]id=([^&]+)",
    ]
    for pattern in patterns:
        match = re.search(pattern, url, re.I)
        if match:
            return f"https://drive.google.com/file/d/{match.group(1)}/preview"
    return None


def _preview_recording_url(url: str, key: str):
    url = str(url or "").strip()
    if not url:
        return

    st.markdown("#### Previsualización de la videograbación")
    drive_preview = _drive_video_preview_url(url)

    try:
        if drive_preview:
            components.iframe(drive_preview, height=520, scrolling=False)
        elif re.search(r"(youtube\.com|youtu\.be|vimeo\.com)", url, re.I):
            st.video(url)
        elif re.search(r"\.(mp4|webm|mov)(?:\?|$)", url, re.I):
            st.video(url)
        else:
            st.info(
                "El proveedor de esta URL puede impedir la reproducción embebida. "
                "Puedes intentar abrirla directamente."
            )
        st.link_button(
            "Abrir videograbación en una pestaña nueva",
            url,
            use_container_width=True,
        )
    except Exception:
        st.warning("No fue posible previsualizar esta URL dentro de la app.")
        st.link_button(
            "Abrir videograbación",
            url,
            use_container_width=True,
        )


def board_session_detail(session: dict):
    year = st.session_state.board_year
    year_label = board_year_label(int(year))
    if st.button(f"← Volver a sesiones de {year_label}"):
        st.session_state.pop("board_session", None)
        st.rerun()
    st.markdown(f"## {session.get('nombre')}")
    st.caption("Versión Junta Visor PPT V12 · láminas con diseño completo")
    st.caption(f"{session.get('tipo')} · Junta de Gobierno · {year_label}")
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
    if not str(session.get("id", "")).startswith("preset-"):
        master_delete_control(
            "sesión de Junta de Gobierno", str(session["id"]), f"board_detail_{session['id']}",
            lambda: (delete_board_session_master(client, str(session["id"])),
                     st.session_state.pop("board_session", None)),
        )
    session_documents = (client.table("documentos_sesion_junta").select("*").eq("sesion_id", session["id"])
                         .order("created_at").execute().data or [])
    latest_document = {}
    for saved_document in session_documents:
        latest_document[saved_document.get("tipo_documento")] = saved_document
    session_date = st.date_input("Fecha de la sesión", value=date.fromisoformat(session["fecha_sesion"][:10]) if session.get("fecha_sesion") else None,
                                 key=f"session_date_{session['id']}")
    if st.button("Guardar fecha de la sesión", key=f"save_session_date_{session['id']}"):
        client.table("sesiones_junta").update({"fecha_sesion": session_date.isoformat() if session_date else None}).eq("id", session["id"]).execute()
        st.session_state.board_session["fecha_sesion"] = session_date.isoformat() if session_date else None
        st.success("Fecha de la sesión guardada.")
    st.markdown("### Documentación principal de la sesión")
    st.caption("Integra aquí los documentos generales. Todos quedarán identificados por usuario y fecha de carga.")
    media1, media_ppt = st.columns(2)
    uploaded = media1.file_uploader("Convocatoria u orden del día", type=["docx", "pdf"], key=f"board_ingest_{session['id']}")
    session_ppt = media_ppt.file_uploader("Presentación de la sesión", type=["pptx"], key=f"board_presentation_{session['id']}")
    if latest_document.get("Convocatoria / orden del día"):
        _saved_main_document(client, latest_document["Convocatoria / orden del día"], media1, "main_agenda")
    if latest_document.get("Presentación de la sesión"):
        _saved_main_document(client, latest_document["Presentación de la sesión"], media_ppt, "main_presentation")
    media2, media3 = st.columns(2)
    signed_minutes = media2.file_uploader("Acta firmada", type=["pdf"], key=f"signed_minutes_{session['id']}")
    if latest_document.get("Acta firmada"):
        _saved_main_document(client, latest_document["Acta firmada"], media2, "main_minutes")
    video_url = media3.text_input("URL de la videograbación", value=session.get("videograbacion_url") or "",
                                  placeholder="https://…", key=f"video_url_{session['id']}")
    if media3.button("Guardar URL", key=f"save_video_{session['id']}", use_container_width=True):
        client.table("sesiones_junta").update({"videograbacion_url": video_url.strip() or None}).eq("id", session["id"]).execute()
        st.session_state.board_session["videograbacion_url"] = video_url.strip() or None
        st.success("URL de la videograbación guardada.")
    if video_url.strip():
        with media3.expander("Previsualizar videograbación", expanded=False):
            _preview_recording_url(video_url.strip(), f"board_recording_{session['id']}")
    if uploaded:
        with media1.expander("Previsualizar convocatoria"):
            if not _document_preview(uploaded.getvalue(), uploaded.name, 480):
                st.info("La vista previa no está disponible para este formato.")
    if signed_minutes:
        with media2.expander("Previsualizar acta"):
            _document_preview(signed_minutes.getvalue(), signed_minutes.name, 480)
    if session_ppt:
        with media_ppt.expander("Previsualizar presentación"):
            if not _document_preview(session_ppt.getvalue(), session_ppt.name, 520):
                st.info("No fue posible convertir esta presentación para previsualizarla.")
    selected_documents = [item for item in [uploaded, session_ppt, signed_minutes] if item]
    if st.button("Guardar documentación seleccionada", type="primary", use_container_width=True,
                 disabled=not selected_documents, key=f"save_main_documents_{session['id']}"):
        document_specs = [
            (uploaded, "Convocatoria / orden del día", "documentos"),
            (session_ppt, "Presentación de la sesión", "presentacion"),
            (signed_minutes, "Acta firmada", "acta_firmada"),
        ]
        saved_count = 0
        try:
            with st.spinner(f"Guardando {len(selected_documents)} documento(s)…"):
                for item, document_type, folder in document_specs:
                    if not item:
                        continue
                    path = f"junta/{session['id']}/{folder}/{uuid.uuid4().hex}_{safe_name(item.name)}"
                    _upload_junta_document(client, path, item)
                    client.table("documentos_sesion_junta").insert({"sesion_id": session["id"], "tipo_documento": document_type,
                        "nombre_visible": Path(item.name).stem, "nombre_archivo": item.name, "ruta_storage": path,
                        "mime_type": item.type, "tamano_bytes": item.size, "subido_por": st.session_state.user["id"],
                        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email")}).execute()
                    if document_type == "Acta firmada":
                        client.table("sesiones_junta").update({"acta_firmada_nombre": item.name, "acta_firmada_ruta": path}).eq("id", session["id"]).execute()
                        st.session_state.board_session.update({"acta_firmada_nombre": item.name, "acta_firmada_ruta": path})
                    saved_count += 1
            st.success(f"Se guardaron {saved_count} documento(s) en el expediente.")
            st.rerun()
        except Exception:
            st.error(f"Se guardaron {saved_count} documento(s), pero uno no pudo completarse. Los ya guardados permanecen en el expediente.")
    if session.get("acta_firmada_ruta"):
        try:
            minutes_data = client.storage.from_("expedientes").download(session["acta_firmada_ruta"])
            media2.download_button(f"Descargar · {session.get('acta_firmada_nombre') or 'Acta firmada'}", minutes_data,
                                   file_name=session.get("acta_firmada_nombre") or "acta_firmada.pdf", mime="application/pdf",
                                   key=f"download_minutes_{session['id']}", use_container_width=True)
        except Exception:
            media2.caption("El acta está registrada, pero no pudo descargarse.")
    if uploaded and st.button("Analizar y separar puntos", type="primary", use_container_width=True):
        try:
            text = _extract_board_text(uploaded)
            if not text.strip(): st.error("El documento no contiene texto seleccionable. Si es un PDF escaneado, conviértelo a PDF con OCR o súbelo en Word.")
            else: st.session_state[f"board_draft_{session['id']}"] = _board_items_from_text(text)
        except Exception as exc: st.error(f"No fue posible leer el archivo: {exc}")
    draft_key = f"board_draft_{session['id']}"
    if draft_key in st.session_state:
        st.markdown("#### Revisión antes de guardar")
        add_col, help_col = st.columns([1, 4])
        if add_col.button("＋ Agregar punto", use_container_width=True):
            st.session_state[draft_key].append({"tipo_registro": "Acuerdo", "titulo": "", "texto": "", "Eliminar": False})
            st.rerun()
        help_col.caption("Marca “Eliminar” en los renglones que no deben guardarse. Puedes editar cualquier texto directamente en la tabla.")
        edited = st.data_editor(pd.DataFrame(st.session_state[draft_key]), use_container_width=True, hide_index=True, num_rows="dynamic",
            column_config={"tipo_registro": st.column_config.SelectboxColumn("Tipo", options=["Acuerdo", "Informe"], required=True),
                           "titulo": st.column_config.TextColumn("Punto del orden del día", width="large", required=True),
                           "texto": st.column_config.TextColumn("Texto / descripción", width="large"),
                           "Eliminar": st.column_config.CheckboxColumn("Eliminar", default=False)}, key=f"board_draft_editor_{session['id']}")
        delete_col, _ = st.columns([1, 4])
        if delete_col.button("Borrar marcados", use_container_width=True):
            st.session_state[draft_key] = [row for row in edited.to_dict("records") if not row.get("Eliminar")]
            st.rerun()
        if st.button("Guardar filas en la sesión", type="primary", use_container_width=True):
            existing = client.table("acuerdos_junta").select("id").eq("sesion_id", session["id"]).execute().data or []
            payload = []
            approved_rows = [row for row in edited.to_dict("records") if not row.get("Eliminar") and str(row.get("titulo") or "").strip()]
            for offset, row in enumerate(approved_rows, len(existing) + 1):
                payload.append({"sesion_id": session["id"], "numero": _agreement_code(session, offset), "tipo_registro": row.get("tipo_registro") or "Acuerdo",
                    "titulo": str(row["titulo"]).strip(), "texto": str(row.get("texto") or "").strip(), "areas": [],
                    "estatus": "Por iniciar", "fecha_compromiso": None, "notificar_email": True})
            if payload:
                client.table("acuerdos_junta").insert(payload).execute(); st.session_state.pop(draft_key, None); st.rerun()
    session_documents = client.table("documentos_sesion_junta").select("*").eq("sesion_id", session["id"]).order("created_at").execute().data or []
    if session_documents:
        st.markdown(f"#### Expediente general de la sesión · {len(session_documents)} documento(s)")
        for document in session_documents:
            _document_card(client, document, "session_doc", "documentos_sesion_junta")
    st.divider()
    rows = client.table("acuerdos_junta").select("*").eq("sesion_id", session["id"]).order("numero").execute().data or []
    rows = [row for row in rows if not re.search(r"clausura(?:\s+de)?\s+la\s+sesi[oó]n|asuntos\s+varios", row.get("titulo") or "", re.I)]
    st.markdown(f"### Acuerdos e informes ({len(rows)})")
    if not rows: st.info("Esta sesión todavía no tiene puntos registrados.")
    for row in rows:
        is_report = (row.get("tipo_registro") == "Informe")
        status = row.get("estatus") or "Por iniciar"; color = {"Por iniciar": "red", "En proceso": "yellow", "Terminada": "green"}.get(status, "gray")
        areas = ", ".join(row.get("areas") or []) or "Sin responsable"
        status_ribbon = {"Por iniciar": ":red[▌]", "En proceso": ":orange[▌]", "Terminada": ":green[▌]"}.get(status, ":gray[▌]")
        with st.expander(
            f"{status_ribbon} {row.get('numero') or 'Sin código'} · {row.get('tipo_registro') or 'Acuerdo'} · {row.get('titulo')}",
            expanded=str(row.get("id")) == str(st.session_state.get("deep_link_agreement_id") or ""),
        ):
            master_reset_agreement_control(
                client,
                agreement_id=str(row["id"]),
                agreement_number=row.get("numero") or "Acuerdo",
                kind="board",
                key=f"board_{row['id']}",
            )
            widget_nonce = int(st.session_state.get(f"agreement_widget_nonce_{row['id']}", 0))
            summary_status = "En progreso" if status == "En proceso" else status
            summary = areas if is_report else f"{summary_status} · {areas}"
            st.markdown(f'<div class="goal-heading status-{"gray" if is_report else color}">{html.escape(summary)}</div>', unsafe_allow_html=True)
            if not is_report: st.caption(f"_{_deadline_label(row.get('fecha_compromiso'), status)}_")
            if row.get("texto"): st.write(row["texto"])
            c1, c2, c3 = st.columns([2, 1, 1])
            new_areas = c1.multiselect("Áreas responsables", BOARD_AREAS, default=row.get("areas") or [], key=f"areas_{row['id']}_{widget_nonce}")
            statuses = ["Por iniciar", "En proceso", "Terminada"]
            display_statuses = {"Por iniciar": "Por iniciar", "En proceso": "En progreso", "Terminada": "Terminada"}
            new_status = status if is_report else c2.selectbox("Estatus", statuses, index=statuses.index(status), format_func=lambda value: display_statuses[value], key=f"status_{row['id']}_{widget_nonce}")
            new_date = None if is_report else c3.date_input("Fecha compromiso", value=date.fromisoformat(row["fecha_compromiso"][:10]) if row.get("fecha_compromiso") else None, key=f"date_{row['id']}_{widget_nonce}")
            responsibles = []
            notify_email = False
            if not is_report:
                st.markdown("##### Responsables y notificación")
                nr1, nr2 = st.columns([3, 1])
                with nr1:
                    responsibles = _responsibles_selector(
                        client,
                        required_module=MODULE_BOARD,
                        current_responsibles=_legacy_responsibles_from_row(row),
                        key=f"board_responsibles_{row['id']}_{widget_nonce}",
                    )
                with nr2:
                    notify_email = st.checkbox(
                        "Enviar recordatorio por correo",
                        value=True if row.get("notificar_email") is None else bool(row.get("notificar_email")),
                        key=f"board_notify_email_{row['id']}_{widget_nonce}",
                    )
                st.caption(
                    _notification_status_caption(
                        new_date,
                        responsibles,
                        notify_email,
                        new_status,
                    )
                )
            result_options = ["Pendiente", "Aprobado", "Rechazado"]
            current_result = row.get("resultado") or "Pendiente"
            new_result = current_result if is_report else st.selectbox("Resultado del acuerdo", result_options,
                index=result_options.index(current_result), key=f"result_{row['id']}_{widget_nonce}")
            if not is_report: c3.caption(_deadline_label(new_date.isoformat() if new_date else None, new_status))
            def save_follow_up():
                close_date = row.get("fecha_cierre")
                compliance = row.get("cumplimiento")
                if new_status == "Terminada":
                    close_date = close_date or date.today().isoformat()
                    compliance = "En tiempo" if new_date and date.fromisoformat(str(close_date)[:10]) <= new_date else ("Extemporáneo" if new_date else "Sin fecha compromiso")
                else:
                    close_date, compliance = None, None
                client.table("acuerdos_junta").update({
                    "areas": new_areas,
                    "estatus": new_status,
                    "resultado": new_result,
                    "fecha_compromiso": new_date.isoformat() if new_date else None,
                    "fecha_cierre": close_date,
                    "cumplimiento": compliance,
                    "responsables_notificacion": responsibles,
                    "responsable_usuario_id": responsibles[0].get("id") if responsibles else None,
                    "responsable_nombre": responsibles[0].get("nombre") if responsibles else None,
                    "responsable_email": responsibles[0].get("email") if responsibles else None,
                    "notificar_email": bool(notify_email),
                    "updated_at": datetime.now().isoformat(),
                }).eq("id", row["id"]).execute()
                description = f"Estatus: {display_statuses[new_status]}; responsables: {', '.join(new_areas) or 'sin responsable'}; fecha compromiso: {new_date.isoformat() if new_date else 'sin fecha'}"
                if responsibles:
                    description += "; personas responsables: " + ", ".join(
                        f"{item.get('nombre') or item.get('email')} ({item.get('email')})"
                        for item in responsibles
                    )
                description += f"; aviso por correo: {'sí' if notify_email else 'no'}"
                if compliance: description += f"; cumplimiento: {compliance}"
                client.table("historial_acuerdo").insert({"acuerdo_id": row["id"], "autor_id": st.session_state.user["id"],
                    "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"), "descripcion": description}).execute()
                st.success("Seguimiento guardado."); st.rerun()
            if row.get("cumplimiento"):
                st.info(f"Resultado de cumplimiento: {row['cumplimiento']} · cierre {row.get('fecha_cierre')}")
            comments = client.table("comentarios_acuerdo").select("*").eq("acuerdo_id", row["id"]).order("created_at").execute().data or []
            history = client.table("historial_acuerdo").select("*").eq("acuerdo_id", row["id"]).order("created_at").execute().data or []
            for comment in comments:
                st.markdown(f"**{html.escape(comment.get('autor_nombre') or 'Usuario')}** · {str(comment.get('created_at') or '')[:16]}"); st.write(comment.get("comentario"))
            with st.form(f"comment_{row['id']}", clear_on_submit=True):
                comment_text = st.text_area("Agregar comentario"); add_comment = st.form_submit_button("Publicar comentario")
            if add_comment and comment_text.strip():
                try:
                    client.table("comentarios_acuerdo").insert({"acuerdo_id": row["id"], "autor_id": st.session_state.user["id"],
                        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"), "comentario": comment_text.strip()}).execute()
                    st.rerun()
                except Exception as exc:
                    st.error(f"No fue posible publicar el comentario: {exc}")
            upload_nonce_key = f"upload_nonce_{row['id']}"
            upload_nonce = st.session_state.get(upload_nonce_key, 0)
            stored_files = client.table("archivos_acuerdo").select("*").eq("acuerdo_id", row["id"]).order("created_at").execute().data or []
            st.markdown(f"#### Expediente documental · {len(stored_files)} archivo(s)")
            st.caption("Documentos, presentaciones y evidencias vinculadas específicamente con este acuerdo.")
            with st.container(border=True):
                st.markdown("##### Incorporar documentos")
                files = st.file_uploader("Seleccionar archivos", accept_multiple_files=True, key=f"files_{row['id']}_{upload_nonce}")
                visible_names = {}
                for file_index, item in enumerate(files or []):
                    visible_names[item.name] = st.text_input(f"Nombre descriptivo · {item.name}", value=Path(item.name).stem,
                                                             key=f"file_title_{row['id']}_{upload_nonce}_{file_index}")
                if files and st.button("Agregar al expediente", key=f"upload_{row['id']}", type="secondary"):
                    records = []
                    for item in files:
                        path = f"junta/{session['id']}/{row['id']}/{uuid.uuid4().hex}_{Path(item.name).name}"
                        client.storage.from_("expedientes").upload(path, item.getvalue(), {"content-type": item.type or "application/octet-stream"})
                        records.append({"acuerdo_id": row["id"], "nombre_visible": visible_names.get(item.name) or Path(item.name).stem,
                                        "nombre_archivo": item.name, "ruta_storage": path, "mime_type": item.type,
                                        "tamano_bytes": item.size, "subido_por": st.session_state.user["id"],
                                        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email")})
                    client.table("archivos_acuerdo").insert(records).execute()
                    st.session_state[upload_nonce_key] = upload_nonce + 1
                    st.success("Archivos incorporados al expediente."); st.rerun()
            for stored in stored_files:
                agreement_document = {**stored, "tipo_documento": "Documento del acuerdo"}
                _document_card(client, agreement_document, "agreement_doc", "archivos_acuerdo")
            if st.button("Guardar cambios de seguimiento", key=f"save_follow_{row['id']}", type="primary", use_container_width=True):
                save_follow_up()
            st.markdown("#### Ficha del acuerdo")
            agreement_docx = _agreement_ficha_docx(row, session, comments, history, stored_files)
            agreement_pdf = _agreement_ficha_pdf(row, session, comments, history, stored_files)
            with st.expander("Previsualizar ficha del acuerdo"):
                _pdf_preview(agreement_pdf)
            d1, d2 = st.columns(2)
            safe_code = row.get("numero") or "acuerdo"
            d1.download_button("Descargar ficha en PDF", agreement_pdf, file_name=f"Ficha_{safe_code}.pdf", mime="application/pdf", use_container_width=True, key=f"agreement_pdf_{row['id']}")
            d2.download_button("Descargar ficha en Word", agreement_docx, file_name=f"Ficha_{safe_code}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True, key=f"agreement_docx_{row['id']}")
    st.markdown("### Agregar acuerdo manualmente")
    with st.expander("＋ Agregar acuerdo que no fue detectado"):
        with st.form(f"manual_agreement_{session['id']}", clear_on_submit=True):
            manual_type = st.selectbox("Tipo", ["Acuerdo", "Informe"])
            manual_title = st.text_input("Nombre del punto")
            manual_text = st.text_area("Texto o descripción")
            save_manual = st.form_submit_button("Agregar a la sesión", type="primary", use_container_width=True)
        if save_manual:
            if not manual_title.strip():
                st.error("Escribe el nombre del punto.")
            else:
                all_rows = client.table("acuerdos_junta").select("numero").eq("sesion_id", session["id"]).execute().data or []
                client.table("acuerdos_junta").insert({"sesion_id": session["id"], "numero": _next_agreement_number(session, all_rows),
                    "tipo_registro": manual_type, "titulo": manual_title.strip(), "texto": manual_text.strip(), "areas": [],
                    "estatus": "Por iniciar", "fecha_compromiso": None, "notificar_email": True}).execute()
                st.success("Acuerdo agregado."); st.rerun()


def board_search(client, year: int, term: str):
    if not term.strip():
        return
    st.markdown("### Resultados de búsqueda")
    if not configured():
        st.info("El buscador quedará activo al conectar la tabla de sesiones y acuerdos en Supabase.")
        return
    try:
        rows = (client.table("acuerdos_junta")
                .select("id,numero,titulo,texto,sesion_id,sesiones_junta!inner(id,anio,nombre,tipo)")
                .eq("sesiones_junta.anio", year)
                .or_(f"titulo.ilike.%{term.strip()}%,texto.ilike.%{term.strip()}%,numero.ilike.%{term.strip()}%")
                .execute().data or [])
        if not rows:
            st.info(f'No se encontraron acuerdos con “{term.strip()}” en {board_year_label(year)}.')
            return
        for row in rows:
            session = row.get("sesiones_junta") or {}
            st.markdown(f'''<div class="session-card" style="--accent:var(--orange)"><h4>{html.escape(row.get("numero") or "Acuerdo")} · {html.escape(row.get("titulo") or "Sin título")}</h4>
            <p>{html.escape(session.get("nombre") or "Sesión sin identificar")} · {html.escape(session.get("tipo") or "")}</p></div>''', unsafe_allow_html=True)
            st.write(row.get("texto") or "Sin descripción")
    except Exception as exc:
        st.warning(f"El buscador estará disponible después de actualizar el esquema de Supabase: {exc}")


def board_year_dashboard(year: int):
    top1, top2 = st.columns([1, 5])
    if top1.button("← Años", use_container_width=True):
        st.session_state.pop("board_year", None)
        st.session_state.pop("board_session", None)
        st.rerun()
    year_label = board_year_label(year)
    top2.markdown(f"## Junta de Gobierno · {year_label}")
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token) if configured() else None
    term = st.text_input("Buscar acuerdos por palabra clave", placeholder="Ej. convenio, presupuesto, inmueble, municipio…", key=f"board_search_{year}")
    board_search(client, year, term)
    sessions = {"Ordinaria": [], "Extraordinaria": []}
    if configured():
        try:
            rows = client.table("sesiones_junta").select("id,anio,tipo,nombre,fecha_sesion,videograbacion_url,acta_firmada_nombre,acta_firmada_ruta").eq("anio", year).order("created_at").execute().data or []
            for row in rows:
                sessions.setdefault(row.get("tipo"), []).append(row)
        except Exception as exc:
            st.warning(f"Ejecuta el esquema actualizado de Supabase para guardar sesiones: {exc}")
    if year in (2025, 2026) and not any(sessions.values()):
        for session_type in sessions:
            sessions[session_type] = [{"id": f"preset-{year}-{session_type}-{index}", "anio": year,
                                       "tipo": session_type, "nombre": name}
                                      for index, name in enumerate(PRESET_BOARD_SESSIONS, 1)]

    ordinary_col, extraordinary_col = st.columns(2, gap="large")
    for column, session_type, accent in [(ordinary_col, "Ordinaria", "var(--blue)"),
                                          (extraordinary_col, "Extraordinaria", "var(--purple)")]:
        with column:
            st.markdown(f'<div class="session-column"><h3>Sesiones {session_type.lower()}s</h3></div>', unsafe_allow_html=True)
            for index, session in enumerate(sessions.get(session_type, [])):
                session_date_label = session.get("fecha_sesion") or "Fecha pendiente"
                st.markdown(f'<div class="session-card" style="--accent:{accent}"><h4>{html.escape(session["nombre"])}</h4><p>Sesión {session_type.lower()} · <b>{html.escape(session_date_label)}</b></p></div>', unsafe_allow_html=True)
                with st.container(key=f'session_open_{session_type}_{session["id"]}'):
                    open_session = st.button(f'Abrir {session["nombre"]}', key=f'open_board_{session_type}_{session["id"]}', use_container_width=True)
                if open_session:
                    st.session_state.board_session = session
                    st.rerun()
                if not str(session.get("id", "")).startswith("preset-"):
                    master_delete_control(
                        "sesión de Junta de Gobierno", str(session["id"]),
                        f"board_card_{session_type}_{session['id']}",
                        lambda session_id=str(session["id"]): delete_board_session_master(client, session_id),
                    )
            with st.expander(f"＋ Agregar sesión {session_type.lower()}"):
                    with st.form(f"new_board_{year}_{session_type}"):
                        session_name = st.text_input("Nombre de la sesión", placeholder="Ej. Primera sesión ordinaria")
                        add_session = st.form_submit_button("Agregar sesión", type="primary", use_container_width=True)
                    if add_session:
                        if not session_name.strip():
                            st.error("Escribe el nombre de la sesión.")
                        elif not configured():
                            st.error("Primero debes conectar Supabase.")
                        else:
                            try:
                                client.table("sesiones_junta").insert({"anio": year, "tipo": session_type,
                                    "nombre": session_name.strip(), "creado_por": st.session_state.user["id"]}).execute()
                                st.success("Sesión agregada.")
                                st.rerun()
                            except Exception as exc:
                                st.error(f"No fue posible agregar la sesión: {exc}")


def _committee_items_from_text(text: str) -> list[dict]:
    """Separa el orden del día y clasifica como acuerdo sólo la fórmula aprobatoria."""
    items = _board_items_from_text(text)
    results = []
    approval_pattern = re.compile(
        r"presentaci[oó]n\s+y\s*,?\s*en\s+su\s+caso\s*,?\s*aprobaci[oó]n",
        re.I,
    )
    for item in items:
        title = item.get("titulo") or ""
        text_value = item.get("texto") or title
        is_agreement = bool(approval_pattern.search(f"{title} {text_value}"))
        results.append({
            "tipo_registro": "Acuerdo" if is_agreement else "Informe",
            "titulo": title,
            "texto": text_value,
            "Eliminar": False,
        })
    return results


def committee_session_detail(session: dict, client):
    """Expediente, puntos y seguimiento de una sesión de comité."""
    if st.button("← Volver a sesiones", key=f"back_committee_session_{session['id']}"):
        st.session_state.pop("committee_session", None)
        st.rerun()
    st.markdown(f"## {session.get('nombre') or 'Sesión'}")
    st.caption(f"{session.get('comite')} · {session.get('tipo')} · {session.get('fecha_sesion') or 'Fecha pendiente'}")
    st.caption("Versión Comités V15 · expediente y documentos de seguimiento")
    master_delete_control(
        "sesión de Comité", str(session["id"]), f"committee_detail_{session['id']}",
        lambda: (delete_committee_session_master(client, str(session["id"])),
                 st.session_state.pop("committee_session", None)),
    )

    try:
        documents = (client.table("documentos_sesion_comite").select("*").eq("sesion_id", session["id"])
                     .order("created_at").execute().data or [])
    except Exception as exc:
        st.error(f"Falta preparar el expediente de Comités en Supabase: {exc}")
        return
    latest = {}
    for document in documents:
        latest[document.get("tipo_documento")] = document

    st.markdown("### Documentación de la sesión")
    st.caption("La convocatoria permite generar los puntos; el acta queda integrada al expediente de la sesión.")
    agenda_col, minutes_col = st.columns(2, gap="large")
    agenda = agenda_col.file_uploader(
        "Convocatoria u orden del día", type=["pdf", "docx"],
        key=f"committee_agenda_{session['id']}",
    )
    minutes = minutes_col.file_uploader(
        "Acta de la sesión", type=["pdf", "docx"],
        key=f"committee_minutes_{session['id']}",
    )
    if latest.get("Convocatoria / orden del día"):
        _saved_main_document(client, latest["Convocatoria / orden del día"], agenda_col, "committee_agenda_saved")
    if latest.get("Acta de la sesión"):
        _saved_main_document(client, latest["Acta de la sesión"], minutes_col, "committee_minutes_saved")
    if agenda:
        with agenda_col.expander("Previsualizar archivo seleccionado"):
            if not _document_preview(agenda.getvalue(), agenda.name, 520):
                st.info("No hay vista previa para este formato; puedes guardarlo y descargarlo desde el expediente.")
    if minutes:
        with minutes_col.expander("Previsualizar archivo seleccionado"):
            if not _document_preview(minutes.getvalue(), minutes.name, 520):
                st.info("No hay vista previa para este formato; puedes guardarlo y descargarlo desde el expediente.")

    selected_documents = [item for item in (agenda, minutes) if item]
    if st.button("Guardar documentación seleccionada", type="primary", use_container_width=True,
                 disabled=not selected_documents, key=f"save_committee_docs_{session['id']}"):
        specs = [(agenda, "Convocatoria / orden del día", "convocatoria"),
                 (minutes, "Acta de la sesión", "acta")]
        saved = 0
        try:
            with st.spinner("Guardando documentación…"):
                for item, document_type, folder in specs:
                    if not item:
                        continue
                    path = f"comites/{session['id']}/{folder}/{uuid.uuid4().hex}_{safe_name(item.name)}"
                    _upload_junta_document(client, path, item)
                    client.table("documentos_sesion_comite").insert({
                        "sesion_id": session["id"], "tipo_documento": document_type,
                        "nombre_visible": Path(item.name).stem, "nombre_archivo": item.name,
                        "ruta_storage": path, "mime_type": item.type,
                        "tamano_bytes": item.size, "subido_por": st.session_state.user["id"],
                        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"),
                    }).execute()
                    saved += 1
            st.success(f"Se guardaron {saved} documento(s).")
            st.rerun()
        except Exception as exc:
            st.error(f"Se guardaron {saved} documento(s), pero otro no pudo completarse: {exc}")

    if agenda and st.button("Analizar y separar puntos de la convocatoria", type="primary",
                            use_container_width=True, key=f"analyze_committee_{session['id']}"):
        try:
            extracted = _extract_board_text(agenda)
            if not extracted.strip():
                st.error("No se encontró texto legible en la convocatoria.")
            else:
                st.session_state[f"committee_draft_{session['id']}"] = _committee_items_from_text(extracted)
                st.rerun()
        except Exception as exc:
            st.error(f"No fue posible leer la convocatoria: {exc}")

    draft_key = f"committee_draft_{session['id']}"
    if draft_key in st.session_state:
        st.markdown("#### Revisión de puntos antes de guardar")
        st.caption("Se marca como acuerdo únicamente el punto que contiene “Presentación y, en su caso, aprobación”. Puedes corregir la clasificación y el texto.")
        edited = st.data_editor(
            pd.DataFrame(st.session_state[draft_key]), use_container_width=True, hide_index=True,
            num_rows="dynamic", key=f"committee_draft_editor_{session['id']}",
            column_config={
                "tipo_registro": st.column_config.SelectboxColumn("Tipo", options=["Acuerdo", "Informe"], required=True),
                "titulo": st.column_config.TextColumn("Punto del orden del día", width="large", required=True),
                "texto": st.column_config.TextColumn("Texto / descripción", width="large"),
                "Eliminar": st.column_config.CheckboxColumn("Eliminar", default=False),
            },
        )
        if st.button("Guardar puntos en la sesión", type="primary", use_container_width=True,
                     key=f"save_committee_items_{session['id']}"):
            rows = [row for row in edited.to_dict("records") if not row.get("Eliminar") and str(row.get("titulo") or "").strip()]
            existing = client.table("acuerdos_comite").select("id").eq("sesion_id", session["id"]).execute().data or []
            payload = []
            for offset, row in enumerate(rows, len(existing) + 1):
                payload.append({
                    "sesion_id": session["id"], "numero": f"COM-{session.get('anio')}-{offset:03d}",
                    "tipo_registro": row.get("tipo_registro") or "Informe",
                    "titulo": str(row.get("titulo") or "").strip(),
                    "texto": str(row.get("texto") or "").strip(), "areas": [],
                    "estatus": "Por iniciar", "resultado": "Pendiente", "notificar_email": True,
                })
            if payload:
                client.table("acuerdos_comite").insert(payload).execute()
            st.session_state.pop(draft_key, None)
            st.success("Puntos guardados.")
            st.rerun()

    documents = (client.table("documentos_sesion_comite").select("*").eq("sesion_id", session["id"])
                 .order("created_at").execute().data or [])
    if documents:
        st.markdown(f"#### Expediente documental · {len(documents)} documento(s)")
        for document in documents:
            _document_card(client, document, "committee_doc", "documentos_sesion_comite")

    st.divider()
    items = (client.table("acuerdos_comite").select("*").eq("sesion_id", session["id"])
             .order("numero").execute().data or [])
    st.markdown(f"### Acuerdos e informes ({len(items)})")
    if not items:
        st.info("Todavía no hay puntos guardados. Sube y analiza la convocatoria.")
    for item in items:
        is_report = item.get("tipo_registro") == "Informe"
        status = item.get("estatus") or "Por iniciar"
        tone = {"Por iniciar": "🔴", "En proceso": "🟡", "Terminada": "🟢"}.get(status, "⚪")
        with st.expander(f"{tone} {item.get('numero')} · {item.get('titulo')}"):
            master_reset_agreement_control(
                client,
                agreement_id=str(item["id"]),
                agreement_number=item.get("numero") or "Acuerdo",
                kind="committee",
                key=f"committee_{item['id']}",
            )
            widget_nonce = int(st.session_state.get(f"agreement_widget_nonce_{item['id']}", 0))
            st.caption(item.get("tipo_registro") or "Punto")
            st.write(item.get("texto") or "Sin descripción adicional.")
            if is_report:
                st.info("Punto informativo: se conserva en el orden del día, pero no requiere seguimiento de acuerdo.")
                continue
            area_value = item.get("areas") or []
            areas = st.multiselect("Áreas responsables", BOARD_AREAS, default=[a for a in area_value if a in BOARD_AREAS],
                                   key=f"committee_areas_{item['id']}_{widget_nonce}")
            c1, c2, c3 = st.columns(3)
            status_value = c1.selectbox("Estatus", ["Por iniciar", "En proceso", "Terminada"],
                                        index=["Por iniciar", "En proceso", "Terminada"].index(status if status in ["Por iniciar", "En proceso", "Terminada"] else "Por iniciar"),
                                        key=f"committee_status_{item['id']}_{widget_nonce}")
            result_value = c2.selectbox("Resultado", ["Pendiente", "Aprobado", "Rechazado"],
                                        index=["Pendiente", "Aprobado", "Rechazado"].index(item.get("resultado") if item.get("resultado") in ["Pendiente", "Aprobado", "Rechazado"] else "Pendiente"),
                                        key=f"committee_result_{item['id']}_{widget_nonce}")
            current_date = date.fromisoformat(str(item["fecha_compromiso"])[:10]) if item.get("fecha_compromiso") else None
            deadline = c3.date_input("Fecha compromiso", value=current_date, key=f"committee_deadline_{item['id']}_{widget_nonce}")
            st.markdown("##### Responsables y notificación")
            nr1, nr2 = st.columns([3, 1])
            with nr1:
                responsibles = _responsibles_selector(
                    client,
                    required_module=MODULE_COMMITTEES,
                    current_responsibles=_legacy_responsibles_from_row(item),
                    key=f"committee_responsibles_{item['id']}_{widget_nonce}",
                )
            with nr2:
                notify_email = st.checkbox(
                    "Enviar recordatorio por correo",
                    value=True if item.get("notificar_email") is None else bool(item.get("notificar_email")),
                    key=f"committee_notify_email_{item['id']}_{widget_nonce}",
                )
            st.caption(
                _notification_status_caption(
                    deadline,
                    responsibles,
                    notify_email,
                    status_value,
                )
            )
            comment = st.text_area("Comentario de seguimiento", value=item.get("comentario_seguimiento") or "",
                                   key=f"committee_comment_{item['id']}_{widget_nonce}")
            st.markdown("#### Documentos de seguimiento")
            st.caption("Adjunta evidencias, oficios, informes o entregables vinculados exclusivamente a este acuerdo.")
            file_name_col, file_col = st.columns([2, 3])
            visible_name = file_name_col.text_input(
                "Nombre descriptivo del documento",
                placeholder="Ej. Evidencia de cumplimiento del acuerdo",
                key=f"committee_followup_name_{item['id']}",
            )
            followup_file = file_col.file_uploader(
                "Archivo", type=["pdf", "docx", "xlsx", "xls", "pptx", "jpg", "jpeg", "png"],
                key=f"committee_followup_file_{item['id']}",
            )
            if st.button(
                "Agregar documento al seguimiento", use_container_width=True,
                key=f"upload_committee_followup_{item['id']}",
                disabled=not followup_file or not visible_name.strip(),
            ):
                try:
                    path = f"comites/acuerdos/{item['id']}/{uuid.uuid4().hex}_{safe_name(followup_file.name)}"
                    _upload_junta_document(client, path, followup_file)
                    client.table("archivos_acuerdo_comite").insert({
                        "acuerdo_id": item["id"], "nombre_visible": visible_name.strip(),
                        "nombre_archivo": followup_file.name, "ruta_storage": path,
                        "mime_type": followup_file.type, "tamano_bytes": followup_file.size,
                        "subido_por": st.session_state.user["id"],
                        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"),
                    }).execute()
                    st.success("Documento incorporado al seguimiento.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"No fue posible guardar el documento: {exc}")

            followup_documents = (client.table("archivos_acuerdo_comite").select("*")
                                  .eq("acuerdo_id", item["id"]).order("created_at").execute().data or [])
            if followup_documents:
                st.caption(f"Expediente del acuerdo · {len(followup_documents)} documento(s)")
                for document in followup_documents:
                    document["tipo_documento"] = "Seguimiento del acuerdo"
                    _document_card(client, document, "committee_followup_doc", "archivos_acuerdo_comite")
            else:
                st.info("Este acuerdo todavía no tiene documentos de seguimiento.")
            if st.button("Guardar cambios de seguimiento", type="primary", use_container_width=True,
                         key=f"save_committee_followup_{item['id']}"):
                client.table("acuerdos_comite").update({
                    "areas": areas,
                    "estatus": status_value,
                    "resultado": result_value,
                    "fecha_compromiso": deadline.isoformat() if deadline else None,
                    "comentario_seguimiento": comment.strip() or None,
                    "responsables_notificacion": responsibles,
                    "responsable_usuario_id": responsibles[0].get("id") if responsibles else None,
                    "responsable_nombre": responsibles[0].get("nombre") if responsibles else None,
                    "responsable_email": responsibles[0].get("email") if responsibles else None,
                    "notificar_email": bool(notify_email),
                    "actualizado_por": st.session_state.user["id"],
                }).eq("id", item["id"]).execute()
                st.success("Seguimiento guardado.")
                st.rerun()


def committee_analytics_dashboard(committee_name: str):
    """Analítica aislada para el comité seleccionado."""
    top1, top2 = st.columns([1, 5])
    if top1.button("← Años", use_container_width=True, key="back_committee_analytics"):
        st.session_state.pop("committee_analytics", None)
        st.rerun()
    top2.markdown(f"## Analítica de datos · {committee_name}")
    st.caption("Versión Comités V16 · analítica individual por comité")
    scope = st.selectbox(
        "Periodo de análisis", ["Todos los años"] + [str(year) for year in BOARD_YEARS],
        key=f"committee_analytics_scope_{committee_name}",
    )
    selected_year = None if scope == "Todos los años" else int(scope)
    if not configured():
        st.error("Primero debes conectar Supabase.")
        return
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
    try:
        session_query = (client.table("sesiones_comite").select("id,anio,fecha_sesion")
                         .eq("comite", committee_name))
        if selected_year:
            session_query = session_query.eq("anio", selected_year)
        sessions = session_query.execute().data or []
        session_ids = [row["id"] for row in sessions]
        agreements = []
        if session_ids:
            agreements = (client.table("acuerdos_comite")
                          .select("id,sesion_id,tipo_registro,estatus,resultado,areas")
                          .in_("sesion_id", session_ids).execute().data or [])
    except Exception as exc:
        st.error(f"No fue posible construir la analítica del comité: {exc}")
        return
    agreements = [row for row in agreements if row.get("tipo_registro") == "Acuerdo"]
    today = date.today().isoformat()
    celebrated = sum(1 for row in sessions if row.get("fecha_sesion") and str(row["fecha_sesion"])[:10] <= today)
    counts = {
        "Sesiones de Comité celebradas": celebrated,
        "Acuerdos aprobados": sum(row.get("resultado") == "Aprobado" for row in agreements),
        "Acuerdos rechazados": sum(row.get("resultado") == "Rechazado" for row in agreements),
        "Acuerdos por iniciar": sum(row.get("estatus") == "Por iniciar" for row in agreements),
        "Acuerdos en progreso": sum(row.get("estatus") == "En proceso" for row in agreements),
        "Acuerdos terminados": sum(row.get("estatus") == "Terminada" for row in agreements),
    }
    tones = ["#0798cf", "#009b4c", "#b85c62", "#f68b08", "#c5a44a", "#16ad8f"]
    cards = "".join(
        f'''<div class="analytics-metric" style="--tone:{tone}">
        <div class="analytics-value">{value}</div><div class="analytics-label">{html.escape(label)}</div></div>'''
        for (label, value), tone in zip(counts.items(), tones)
    )
    st.markdown("### Numeralia")
    st.caption(
        f"Cifras acumuladas de {committee_name} para 2025–2030."
        if selected_year is None else f"Cifras de {committee_name} correspondientes a {selected_year}."
    )
    st.markdown(f'<div class="analytics-metrics">{cards}</div>', unsafe_allow_html=True)
    pending_result = sum(row.get("resultado") in (None, "Pendiente") for row in agreements)
    if pending_result:
        st.info(f"Hay {pending_result} acuerdo(s) pendientes de clasificar como aprobado o rechazado.")

    st.markdown("### Avance por áreas responsables")
    st.caption("Cantidad de acuerdos asignados a cada dirección. Se excluyen Dirección General y Órgano Interno de Control.")
    included_areas = [area for area in BOARD_AREAS if area not in ("Dirección General", "Órgano Interno de Control")]
    area_data = [
        {"Dirección": area, "Acuerdos": sum(area in (row.get("areas") or []) for row in agreements)}
        for area in included_areas
    ]
    area_df = pd.DataFrame(area_data)
    chart = (alt.Chart(area_df).mark_bar(cornerRadiusTopLeft=7, cornerRadiusTopRight=7, size=54)
             .encode(
                 x=alt.X("Dirección:N", sort=included_areas, title=None, axis=alt.Axis(labelAngle=0, labelLimit=260)),
                 y=alt.Y("Acuerdos:Q", title="Cantidad de acuerdos", axis=alt.Axis(tickMinStep=1)),
                 color=alt.Color("Dirección:N", scale=alt.Scale(
                     domain=included_areas, range=["#0798cf", "#009b4c", "#a990c7"]), legend=None),
                 tooltip=[alt.Tooltip("Dirección:N"), alt.Tooltip("Acuerdos:Q", format=".0f")],
             ))
    labels = alt.Chart(area_df).mark_text(dy=-12, fontSize=15, fontWeight="bold", color="#35434b").encode(
        x=alt.X("Dirección:N", sort=included_areas), y="Acuerdos:Q", text=alt.Text("Acuerdos:Q", format=".0f"))
    st.altair_chart((chart + labels).properties(height=390), use_container_width=True)


def committees():
    if not user_can(MODULE_COMMITTEES):
        st.error("No tienes permisos para acceder a Comités.")
        return
    selected_committee = st.session_state.get("committee_name")
    selected_year = st.session_state.get("committee_year")
    selected_session = st.session_state.get("committee_session")
    if not selected_committee:
        st.markdown('<h1 class="choice-title">Comités</h1>', unsafe_allow_html=True)
        st.markdown('<p class="choice-subtitle">Selecciona el comité que deseas consultar</p>', unsafe_allow_html=True)
        for start in range(0, len(COMMITTEE_CATALOG), 2):
            columns = st.columns(2, gap="large")
            for column, (name, description, accent) in zip(columns, COMMITTEE_CATALOG[start:start + 2]):
                with column:
                    st.markdown(f'''<div class="card" style="--accent:{accent};min-height:185px"><div class="card-icon">C</div>
                        <h3>{html.escape(name)}</h3><p class="muted">{html.escape(description)}</p></div>''', unsafe_allow_html=True)
                    if st.button(f"Abrir {name}", key=f"open_committee_{name}", use_container_width=True, type="primary"):
                        st.session_state.committee_name = name
                        st.session_state.pop("committee_year", None)
                        st.session_state.pop("committee_analytics", None)
                        st.rerun()
        return
    if st.session_state.get("committee_analytics"):
        committee_analytics_dashboard(selected_committee)
        return
    if selected_year is None:
        top1, top2 = st.columns([1, 5])
        if top1.button("← Comités", use_container_width=True):
            st.session_state.pop("committee_name", None)
            st.session_state.pop("committee_session", None)
            st.session_state.pop("committee_analytics", None)
            st.rerun()
        top2.markdown(f"## {selected_committee}")
        st.markdown('<p class="choice-subtitle">Selecciona el año de trabajo</p>', unsafe_allow_html=True)
        year_columns = st.columns(3, gap="large")
        year_colors = ["var(--blue)", "var(--green)", "var(--teal)", "var(--purple)", "var(--orange)", "var(--gray)"]
        for index, year in enumerate(range(2025, 2031)):
            with year_columns[index % 3]:
                st.markdown(f'''<div class="year-card" style="--accent:{year_colors[index]}">
                    <h2>{year}</h2><p>Sesiones, acuerdos y documentación</p></div>''', unsafe_allow_html=True)
                if st.button(f"Abrir {year}", key=f"committee_year_{selected_committee}_{year}", use_container_width=True):
                    st.session_state.committee_year = year
                    st.session_state.pop("committee_session", None)
                    st.rerun()
        st.markdown(f'''<div class="analytics-banner"><h3>Analítica de datos · {html.escape(selected_committee)}</h3>
            <p>Numeralia de sesiones y acuerdos, así como distribución del avance por áreas responsables.</p></div>''',
            unsafe_allow_html=True)
        if st.button("Abrir Analítica de datos", key=f"open_committee_analytics_{selected_committee}",
                     use_container_width=True, type="primary"):
            st.session_state.committee_analytics = True
            st.session_state.pop("committee_year", None)
            st.session_state.pop("committee_session", None)
            st.rerun()
        return
    top1, top2 = st.columns([1, 5])
    if top1.button("← Años", use_container_width=True):
        st.session_state.pop("committee_year", None)
        st.session_state.pop("committee_session", None)
        st.rerun()
    top2.markdown(f"## {selected_committee} · {selected_year}")
    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token) if configured() else None
    if selected_session:
        committee_session_detail(selected_session, client)
        return
    sessions = []
    if configured():
        try:
            sessions = (client.table("sesiones_comite").select("*").eq("comite", selected_committee)
                        .eq("anio", selected_year).order("fecha_sesion").order("created_at").execute().data or [])
        except Exception as exc:
            st.warning(f"Ejecuta la migración de Comités para guardar sesiones: {exc}")
    st.markdown("### Sesiones")
    if not sessions:
        st.info("Todavía no hay sesiones registradas para este comité y año.")
    for session in sessions:
        date_label = session.get("fecha_sesion") or "Fecha pendiente"
        with st.container(border=True):
            info, action = st.columns([5, 1], vertical_alignment="center")
            info.markdown(f"#### {html.escape(session.get('nombre') or 'Sesión')}")
            info.caption(f"{session.get('tipo') or 'Sesión'} · {date_label}")
            if action.button("Abrir sesión", key=f"open_committee_session_{session['id']}",
                             type="primary", use_container_width=True):
                st.session_state.committee_session = session
                st.rerun()
            master_delete_control(
                "sesión de Comité", str(session["id"]), f"committee_card_{session['id']}",
                lambda session_id=str(session["id"]): delete_committee_session_master(client, session_id),
            )
    st.markdown("---")
    with st.expander("＋ Crear nueva sesión", expanded=not sessions):
        with st.form(f"new_committee_session_{selected_committee}_{selected_year}", clear_on_submit=True):
            c1, c2 = st.columns(2)
            session_name = c1.text_input("Nombre de la sesión", placeholder="Ej. Primera sesión ordinaria")
            session_type = c2.selectbox("Tipo de sesión", ["Ordinaria", "Extraordinaria"])
            session_date = st.date_input("Fecha de la sesión", value=None)
            add_session = st.form_submit_button("Crear sesión", type="primary", use_container_width=True)
        if add_session:
            if not session_name.strip():
                st.error("Escribe el nombre de la sesión.")
            elif not configured():
                st.error("Primero debes conectar Supabase.")
            else:
                try:
                    client.table("sesiones_comite").insert({"comite": selected_committee, "anio": selected_year,
                        "tipo": session_type, "nombre": session_name.strip(),
                        "fecha_sesion": session_date.isoformat() if session_date else None,
                        "creado_por": st.session_state.user["id"]}).execute()
                    st.success("Sesión creada correctamente.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"No fue posible crear la sesión: {exc}")


def _normalize_office_link_key(value) -> str:
    """Normaliza números de oficio para cruzar vínculos de Drive de forma tolerante."""
    raw = str(value or "").upper().strip()
    raw = raw.replace("\\", "/")
    raw = re.sub(r"\s+", "", raw)
    raw = raw.replace("DGE", "DG/E").replace("DGI", "DG/I")
    raw = raw.replace("DG-E", "DG/E").replace("DG-I", "DG/I")
    raw = re.sub(r"/+", "/", raw)
    raw = re.sub(r"-+", "-", raw)
    return raw.strip("-/")


def _parse_drive_links_file(uploaded) -> tuple[list[dict], list[str]]:
    """Acepta Excel o CSV con numero_oficio, nombre_archivo, url_drive."""
    filename = (getattr(uploaded, "name", "") or "").lower()
    if filename.endswith(".csv"):
        frame = pd.read_csv(uploaded, dtype=object)
    else:
        frame = pd.read_excel(uploaded, dtype=object)

    frame.columns = [str(c).strip().lower() for c in frame.columns]

    aliases = {
        "folio_control": ["folio_control", "folio control", "folio", "no", "n°", "numero_control", "número control"],
        "numero_oficio": ["numero_oficio", "número de oficio", "numero de oficio", "oficio"],
        "nombre_archivo": ["nombre_archivo", "nombre archivo", "archivo", "nombre del archivo"],
        "url_drive": ["url_drive", "url drive", "liga", "liga drive", "url", "enlace", "enlace drive"],
    }

    resolved = {}
    for target, options in aliases.items():
        for option in options:
            if option in frame.columns:
                resolved[target] = option
                break

    if "url_drive" not in resolved or ("numero_oficio" not in resolved and "folio_control" not in resolved):
        raise ValueError(
            "El archivo debe incluir 'url_drive' y al menos una de estas columnas: "
            "'folio_control' o 'numero_oficio'. 'nombre_archivo' es opcional."
        )

    records, warnings = [], []
    seen = set()

    for idx, row in frame.iterrows():
        office_number = ""
        if "numero_oficio" in resolved:
            office_number = str(row.get(resolved["numero_oficio"], "") or "").strip()
        folio_control = ""
        if "folio_control" in resolved:
            folio_control = str(row.get(resolved["folio_control"], "") or "").strip()
        url = str(row.get(resolved["url_drive"], "") or "").strip()
        file_name = ""
        if "nombre_archivo" in resolved:
            file_name = str(row.get(resolved["nombre_archivo"], "") or "").strip()

        excel_row = idx + 2
        if not office_number and not folio_control and not url:
            continue
        if not office_number and not folio_control:
            warnings.append(f"Fila {excel_row}: falta folio_control y numero_oficio.")
            continue
        if not url:
            warnings.append(f"Fila {excel_row}: falta url_drive.")
            continue
        if "drive.google.com" not in url:
            warnings.append(
                f"Fila {excel_row}: la URL no parece ser de Google Drive: "
                f"{office_number or folio_control}."
            )

        folio_norm = re.sub(r"\.0$", "", folio_control.upper().strip())
        folio_norm = re.sub(r"\s+", "", folio_norm)
        key = _normalize_office_link_key(office_number) if office_number else ""

        dedupe_key = folio_norm or key
        if dedupe_key in seen:
            warnings.append(
                f"Fila {excel_row}: folio/número repetido en el archivo: "
                f"{folio_control or office_number}."
            )
        seen.add(dedupe_key)

        records.append({
            "folio_control": folio_control or None,
            "folio_normalizado": folio_norm,
            "numero_oficio": office_number or None,
            "numero_normalizado": key,
            "nombre_archivo": file_name or None,
            "url_drive": url,
            "fila_origen": excel_row,
        })

    return records, warnings


def _import_drive_links(client, uploaded) -> dict:
    records, warnings = _parse_drive_links_file(uploaded)

    offices = client.table("oficios_direccion_general").select(
        "id,numero_oficio,folio_control,drive_url"
    ).execute().data or []

    index = {}
    folio_index = {}
    for row in offices:
        key = _normalize_office_link_key(row.get("numero_oficio"))
        if key:
            index.setdefault(key, []).append(row)

        folio = str(row.get("folio_control") or "").upper().strip()
        folio = re.sub(r"\.0$", "", folio)
        folio = re.sub(r"\s+", "", folio)
        if folio:
            folio_index.setdefault(folio, []).append(row)

    linked = ambiguous = unmatched = 0
    now_iso = datetime.now().isoformat()
    details = []

    for record in records:
        # 1) Empieza por folio_control / NO del Excel.
        #    OJO: el mismo folio se repite entre años, por lo que un folio
        #    no debe declararse ambiguo hasta intentar desambiguarlo con
        #    el número completo del oficio.
        candidates = []
        folio_candidates = []
        if record.get("folio_normalizado"):
            folio_candidates = folio_index.get(record["folio_normalizado"], [])

        # 2) Si el folio produjo varias coincidencias (por ejemplo 2024/2025/2026),
        #    intenta quedarte con la que coincide también con numero_oficio.
        if folio_candidates and record.get("numero_normalizado"):
            exact_within_folio = [
                row for row in folio_candidates
                if _normalize_office_link_key(row.get("numero_oficio")) == record["numero_normalizado"]
            ]
            if exact_within_folio:
                candidates = exact_within_folio
            elif len(folio_candidates) == 1:
                candidates = folio_candidates
        elif len(folio_candidates) == 1:
            candidates = folio_candidates

        # 3) Si el folio no resolvió de manera inequívoca, intenta el número
        #    completo contra toda la base.
        if len(candidates) != 1 and record.get("numero_normalizado"):
            exact_global = index.get(record["numero_normalizado"], [])
            if exact_global:
                candidates = exact_global

        # 4) Fallback por folio + año. Esto rescata archivos cuyo nombre
        #    tiene errores de formato, mes incompleto o variantes DGE/DGI.
        if len(candidates) != 1 and record.get("folio_normalizado"):
            year_hint = None
            month_hint = None
            type_hint = None

            raw_num = str(record.get("numero_oficio") or "")
            raw_name = str(record.get("nombre_archivo") or "")
            combined = (raw_num + " " + raw_name).upper()

            ym = re.search(r"(20\d{2})", combined)
            if ym:
                year_hint = ym.group(1)

            mm = re.search(r"(?:-|/|\.)(0[1-9]|1[0-2])(?:-|/|\.|\s)*(?:20\d{2}|$)", combined)
            if mm:
                month_hint = mm.group(1)
            else:
                mm2 = re.search(r"(0[1-9]|1[0-2])\s*20\d{2}", combined)
                if mm2:
                    month_hint = mm2.group(1)

            if re.search(r"\b(?:DG[\s._/-]*I|DGI)\b", combined):
                type_hint = "I"
            elif re.search(r"\b(?:DG[\s._/-]*E|DGE)\b", combined):
                type_hint = "E"

            same_folio = folio_index.get(record["folio_normalizado"], [])
            narrowed = same_folio

            if year_hint:
                by_year = [
                    row for row in narrowed
                    if year_hint in str(row.get("numero_oficio") or "")
                ]
                if by_year:
                    narrowed = by_year

            if month_hint and len(narrowed) > 1:
                by_month = []
                for row in narrowed:
                    norm = _normalize_office_link_key(row.get("numero_oficio"))
                    if re.search(rf"\D{re.escape(month_hint)}\D+20\d{{2}}$", norm):
                        by_month.append(row)
                if by_month:
                    narrowed = by_month

            if type_hint and len(narrowed) > 1:
                by_type = []
                for row in narrowed:
                    norm = _normalize_office_link_key(row.get("numero_oficio"))
                    if type_hint == "I" and "DG/I" in norm:
                        by_type.append(row)
                    elif type_hint == "E" and "DG/E" in norm:
                        by_type.append(row)
                if by_type:
                    narrowed = by_type

            if len(narrowed) == 1:
                candidates = narrowed

        # 5) Fallback por consecutivo + mes + año si el formato difiere.
        if len(candidates) != 1 and record.get("numero_normalizado"):
            m = re.search(r"(\d{1,3}(?:-BIS)?)\D+(\d{2})\D+(20\d{2})$", record["numero_normalizado"])
            if m:
                seq, month, year = m.groups()
                seq_norm = seq.upper()
                fallback = []
                for row in offices:
                    norm = _normalize_office_link_key(row.get("numero_oficio"))
                    m2 = re.search(r"(\d{1,3}(?:-BIS)?)\D+(\d{2})\D+(20\d{2})$", norm)
                    if m2 and m2.groups() == (seq_norm, month, year):
                        fallback.append(row)
                if fallback:
                    candidates = fallback

        if len(candidates) == 1:
            office = candidates[0]
            client.table("oficios_direccion_general").update({
                "drive_url": record["url_drive"],
                "drive_nombre_archivo": record["nombre_archivo"],
                "drive_vinculado_at": now_iso,
                "updated_at": now_iso,
            }).eq("id", office["id"]).execute()
            linked += 1
            details.append({
                "folio_control": record.get("folio_control"),
                "numero_oficio": record.get("numero_oficio"),
                "estado": "Vinculado",
                "url_drive": record["url_drive"],
            })
        elif len(candidates) > 1:
            ambiguous += 1
            details.append({
                "folio_control": record.get("folio_control"),
                "numero_oficio": record.get("numero_oficio"),
                "estado": "Ambiguo",
                "url_drive": record["url_drive"],
            })
        else:
            unmatched += 1
            details.append({
                "folio_control": record.get("folio_control"),
                "numero_oficio": record.get("numero_oficio"),
                "estado": "Sin coincidencia",
                "url_drive": record["url_drive"],
            })

    # Guardar bitácora de importación.
    try:
        client.table("ingestas_vinculos_drive_oficios_dg").insert({
            "nombre_archivo": getattr(uploaded, "name", None),
            "registros_detectados": len(records),
            "registros_vinculados": linked,
            "registros_ambiguos": ambiguous,
            "registros_sin_coincidencia": unmatched,
            "importado_por": st.session_state.user["id"],
            "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"),
            "detalle": details,
            "created_at": now_iso,
        }).execute()
    except Exception:
        pass

    return {
        "detectados": len(records),
        "vinculados": linked,
        "ambiguos": ambiguous,
        "sin_coincidencia": unmatched,
        "warnings": warnings,
        "details": details,
    }


def _latest_drive_link_import(client):
    try:
        rows = (client.table("ingestas_vinculos_drive_oficios_dg").select("*")
                .order("created_at", desc=True).limit(1).execute().data or [])
        return rows[0] if rows else None
    except Exception:
        return None


OFFICIAL_LETTER_YEARS = list(range(2024, 2031))
MONTHS_ES = [
    (1, "Enero"), (2, "Febrero"), (3, "Marzo"), (4, "Abril"),
    (5, "Mayo"), (6, "Junio"), (7, "Julio"), (8, "Agosto"),
    (9, "Septiembre"), (10, "Octubre"), (11, "Noviembre"), (12, "Diciembre"),
]



OFFICIAL_THEME_CATALOG = [
    "Comité de Adquisiciones",
    "Junta de Gobierno y otros comités",
    "Adquisiciones y contrataciones",
    "Presupuesto y finanzas",
    "Personal y organización",
    "Jurídico y litigios",
    "Patrimonio, inmuebles y arrendamientos",
    "Proyectos e inversión",
    "Auditoría, control y cumplimiento",
    "Comisiones y viáticos",
    "Archivo, transparencia y gestión documental",
    "Coordinación institucional y enlaces",
    "Otros / por clasificar",
]


def _theme_text(value) -> str:
    raw = str(value or "").strip().upper()
    if not raw:
        return ""
    raw = unicodedata.normalize("NFKD", raw)
    raw = "".join(ch for ch in raw if not unicodedata.combining(ch))
    raw = re.sub(r"\s+", " ", raw)
    return raw


def _contains_any(text_value: str, terms: tuple[str, ...] | list[str]) -> bool:
    return any(term in text_value for term in terms)


def _classify_official_letter(row: dict) -> dict:
    """Clasificación temática determinística para los oficios de Dirección General.

    Se apoya principalmente en ASUNTO y, como contexto, en dependencia,
    destinatario y cargo. La clasificación se guarda en Supabase para que la
    analítica no tenga que recalcularse en cada carga.
    """
    subject = _theme_text(row.get("asunto"))
    dependency = _theme_text(row.get("dependencia"))
    recipient = _theme_text(row.get("destinatario"))
    role = _theme_text(row.get("cargo"))
    text_value = " | ".join(v for v in (subject, dependency, recipient, role) if v)

    def result(theme: str, subtheme: str, confidence: float):
        return {
            "tema": theme,
            "subtema": subtheme,
            "clasificacion_confianza": confidence,
            "clasificacion_fuente": "reglas_v1",
        }

    # 1. Comité de Adquisiciones: se separa porque representa una carga propia.
    if _contains_any(text_value, (
        "COMITE DE ADQUISICIONES", "COMITÉ DE ADQUISICIONES",
        "SESION ORDINARIA COMITE", "SESION EXTRAORDINARIA COMITE",
        "SESIÓN ORDINARIA COMITÉ", "SESIÓN EXTRAORDINARIA COMITÉ",
        "SIN CONCURRENCIA DEL COMITE", "SIN CONCURRENCIA DEL COMITÉ",
    )):
        if "CONVOCATORIA" in text_value:
            sub = "Convocatorias"
        elif _contains_any(text_value, ("FALLO", "APERTURA", "BASES", "PROPUESTAS")):
            sub = "Sesiones, bases y fallos"
        elif _contains_any(text_value, ("PRESIDENTE SUPLENTE", "SECRETARIO TECNICO", "SECRETARIO TÉCNICO")):
            sub = "Integración del comité"
        else:
            sub = "Gestión del Comité de Adquisiciones"
        return result("Comité de Adquisiciones", sub, 0.98)

    # 2. Junta de Gobierno y otros órganos colegiados.
    if _contains_any(text_value, (
        "JUNTA DE GOBIERNO", "ORDEN DEL DIA", "ORDEN DEL DÍA",
        "SESION ORDINARIA", "SESIÓN ORDINARIA", "SESION EXTRAORDINARIA", "SESIÓN EXTRAORDINARIA",
        "COMITE DE TRANSPARENCIA", "COMITÉ DE TRANSPARENCIA",
        "COMITE DE CONTROL INTERNO", "COMITÉ DE CONTROL INTERNO",
        "COMITE JALISCO ATRAE", "COMITÉ JALISCO ATRAE",
        "UNIDAD DE IGUALDAD DE GENERO", "UNIDAD DE IGUALDAD DE GÉNERO",
    )):
        if "JUNTA DE GOBIERNO" in text_value or "ORDEN DEL DIA" in text_value or "ORDEN DEL DÍA" in text_value:
            sub = "Junta de Gobierno"
        elif "TRANSPARENCIA" in text_value:
            sub = "Comité de Transparencia"
        elif "CONTROL INTERNO" in text_value:
            sub = "Comité de Control Interno"
        elif "IGUALDAD" in text_value:
            sub = "Igualdad de Género"
        elif "JALISCO ATRAE" in text_value:
            sub = "Comité Jalisco Atrae"
        else:
            sub = "Otros órganos colegiados"
        return result("Junta de Gobierno y otros comités", sub, 0.92)

    # 3. Comisiones y viáticos.
    if _contains_any(subject, (
        "COMISION", "COMISIÓN", "VIATICO", "VIÁTICO", "VIAJE", "TRANSPORTE AEREO",
        "TRANSPORTE AÉREO", "VISITA A ", "REUNION EN ", "REUNIÓN EN ",
    )):
        if _contains_any(subject, ("VIATICO", "VIÁTICO")):
            sub = "Viáticos"
        elif _contains_any(subject, ("TRANSPORTE AEREO", "TRANSPORTE AÉREO", "EXTRANJERO")):
            sub = "Traslados y transporte"
        else:
            sub = "Comisiones de trabajo"
        return result("Comisiones y viáticos", sub, 0.97)

    # 4. Presupuesto y finanzas.
    if _contains_any(text_value, (
        "PRESUPUEST", "ADECUACION", "ADECUACIÓN", "AMPLIACION", "AMPLIACIÓN",
        "CUENTA BANCARIA", "CUENTAS BANCARIAS", "SUBSIDIO", "MINISTRACION", "MINISTRACIÓN",
        "INSTRUCCION DE PAGO", "INSTRUCCIÓN DE PAGO", "SIIF", "NOMINA", "NÓMINA",
        "ISR", "CUENTA PUBLICA", "CUENTA PÚBLICA", "ESTADOS FINANCIEROS", "EDOS FINANCIEROS",
    )):
        if _contains_any(text_value, ("ADECUACION", "ADECUACIÓN", "AMPLIACION", "AMPLIACIÓN", "MODIFICACION PRESUPUEST", "MODIFICACIÓN PRESUPUEST")):
            sub = "Adecuaciones y ampliaciones presupuestales"
        elif "CUENTA BANCARIA" in text_value or "CUENTAS BANCARIAS" in text_value:
            sub = "Cuentas bancarias"
        elif _contains_any(text_value, ("SUBSIDIO", "PAGO", "MINISTRACION", "MINISTRACIÓN")):
            sub = "Pagos y ministraciones"
        elif _contains_any(text_value, ("CUENTA PUBLICA", "CUENTA PÚBLICA", "ESTADOS FINANCIEROS", "EDOS FINANCIEROS")):
            sub = "Información financiera y cuenta pública"
        elif "SIIF" in text_value:
            sub = "Sistemas financieros"
        else:
            sub = "Planeación presupuestal"
        return result("Presupuesto y finanzas", sub, 0.92)

    # 5. Personal y organización.
    if _contains_any(text_value, (
        "CONTRATACION DE PERSONAL", "CONTRATACIÓN DE PERSONAL", "ALTA DE PERSONAL",
        "BAJA DE PERSONAL", "CONSTANCIA DE BAJA", "HOJA DE SERVICIO", "PLANTILLA",
        "ESTRUCTURA ORGANICA", "ESTRUCTURA ORGÁNICA", "REGLAMENTO INTERNO", "MOP",
        "ENTREGA RECEPCION", "ENTREGA RECEPCIÓN", "SUPLENCIA", "DESIGNA SUPLENTE",
        "NOMBRAMIENTO", "VACANTE", "PERSONAL EVENTUAL",
    )):
        if _contains_any(text_value, ("CONTRATACION", "CONTRATACIÓN", "ALTA DE PERSONAL", "VACANTE", "PERSONAL EVENTUAL")):
            sub = "Contratación y altas"
        elif _contains_any(text_value, ("BAJA DE PERSONAL", "CONSTANCIA DE BAJA", "HOJA DE SERVICIO")):
            sub = "Bajas y movimientos"
        elif _contains_any(text_value, ("PLANTILLA", "ESTRUCTURA ORGANICA", "ESTRUCTURA ORGÁNICA", "REGLAMENTO INTERNO", "MOP")):
            sub = "Estructura, plantilla y normativa interna"
        elif _contains_any(text_value, ("ENTREGA RECEPCION", "ENTREGA RECEPCIÓN")):
            sub = "Entrega-recepción"
        else:
            sub = "Nombramientos y suplencias"
        return result("Personal y organización", sub, 0.93)

    # 6. Jurídico y litigios.
    if _contains_any(text_value, (
        "JUZGADO", "AMPARO", "RECURSO DE REVOCACION", "RECURSO DE REVOCACIÓN",
        "ACTOS PREJUDICIALES", "EMPLAZAMIENTO", "CARPETA DE INVESTIGACION", "CARPETA DE INVESTIGACIÓN",
        "CONVENIO JUDICIAL", "OPINION JURIDICA", "OPINIÓN JURÍDICA", "INCIDENTE",
        "COMPARECE", "NOTARIA", "NOTARÍA", "TESTIMONIO", "CARTA DE INSTRUCCION", "CARTA DE INSTRUCCIÓN",
    )):
        if _contains_any(text_value, ("JUZGADO", "AMPARO", "RECURSO", "EMPLAZAMIENTO", "ACTOS PREJUDICIALES", "INCIDENTE", "COMPARECE")):
            sub = "Litigios y procedimientos"
        elif _contains_any(text_value, ("CONVENIO JUDICIAL", "OPINION JURIDICA", "OPINIÓN JURÍDICA")):
            sub = "Opiniones y convenios jurídicos"
        elif _contains_any(text_value, ("NOTARIA", "NOTARÍA", "TESTIMONIO", "CARTA DE INSTRUCCION", "CARTA DE INSTRUCCIÓN")):
            sub = "Gestiones notariales"
        else:
            sub = "Asuntos jurídicos"
        return result("Jurídico y litigios", sub, 0.91)

    # 7. Auditoría, control y cumplimiento.
    if _contains_any(text_value, (
        "AUDITORIA", "AUDITORÍA", "OIC", "CONTROL INTERNO", "ANTICORRUPCION", "ANTICORRUPCIÓN",
        "OBSERVACIONES", "ENTE FISCALIZADOR", "ENTES FISCALIZADORES", "CUMPLIMIENTO A DISPOSICIONES",
    )):
        if _contains_any(text_value, ("AUDITORIA", "AUDITORÍA", "OBSERVACIONES", "ENTE FISCALIZADOR", "ENTES FISCALIZADORES")):
            sub = "Auditorías y observaciones"
        elif "CONTROL INTERNO" in text_value:
            sub = "Control interno"
        elif _contains_any(text_value, ("ANTICORRUPCION", "ANTICORRUPCIÓN")):
            sub = "Anticorrupción"
        else:
            sub = "Cumplimiento"
        return result("Auditoría, control y cumplimiento", sub, 0.94)

    # 8. Archivo, transparencia y gestión documental.
    if _contains_any(text_value, (
        "ARCHIVO", "GESTION DOCUMENTAL", "GESTIÓN DOCUMENTAL", "OFICIALIA DE PARTES", "OFICIALÍA DE PARTES",
        "TRANSFERENCIA PRIMARIA", "TRANFERECIA PRIMARIA", "PUBLICACION", "PUBLICACIÓN",
        "PERIODICO OFICIAL", "PERIÓDICO OFICIAL", "TRANSPARENCIA",
    )):
        if _contains_any(text_value, ("ARCHIVO", "TRANSFERENCIA PRIMARIA", "TRANFERECIA PRIMARIA")):
            sub = "Archivo institucional"
        elif _contains_any(text_value, ("PUBLICACION", "PUBLICACIÓN", "PERIODICO OFICIAL", "PERIÓDICO OFICIAL")):
            sub = "Publicaciones oficiales"
        elif "TRANSPARENCIA" in text_value:
            sub = "Transparencia"
        else:
            sub = "Gestión documental"
        return result("Archivo, transparencia y gestión documental", sub, 0.90)

    # 9. Patrimonio, inmuebles y arrendamientos.
    if _contains_any(text_value, (
        "PREDIO", "INMUEBLE", "INMUEBLES", "COMODATO", "ARRENDAMIENTO", "ARRENDAMIENTOS",
        "CATASTRO", "CATASTRAL", "AVALUO", "AVALÚO", "AVALUOS", "AVALÚOS",
        "SUBDIVISION", "SUBDIVISIÓN", "COMPRA VENTA", "COMPRAVENTA", "PROPIEDAD",
        "BIENES MUEBLES", "BIENES INMUEBLES", "RECTIFICACION DE COLINDANCIA", "RECTIFICACIÓN DE COLINDANCIA",
        "USO DE SUELO",
    )):
        if "COMODATO" in text_value:
            sub = "Comodatos"
        elif "ARREND" in text_value:
            sub = "Arrendamientos"
        elif _contains_any(text_value, ("AVALUO", "AVALÚO", "AVALUOS", "AVALÚOS", "CATASTRO", "CATASTRAL")):
            sub = "Avalúos y catastro"
        elif _contains_any(text_value, ("COMPRA VENTA", "COMPRAVENTA")):
            sub = "Compraventa"
        elif _contains_any(text_value, ("PREDIO", "SUBDIVISION", "SUBDIVISIÓN", "COLINDANCIA", "USO DE SUELO")):
            sub = "Predios y ordenamiento"
        else:
            sub = "Patrimonio"
        return result("Patrimonio, inmuebles y arrendamientos", sub, 0.91)

    # 10. Adquisiciones y contrataciones.
    if _contains_any(text_value, (
        "LICITACION", "LICITACIÓN", "ADJUDICACION", "ADJUDICACIÓN", "APROVISIONAMIENTO",
        "ADQUISICION", "ADQUISICIÓN", "SEGURO", "VIGILANCIA", "MANTENIMIENTO",
        "FOTOCOPIADO", "GASOLINA", "DESPENSA", "LUMINARIA", "VEHICULO", "VEHÍCULO",
        "EQUIPO DE COMPUTO", "EQUIPO DE CÓMPUTO", "SERVICIO",
    )):
        if _contains_any(text_value, ("LICITACION", "LICITACIÓN", "TIEMPOS RECORTADOS")):
            sub = "Licitaciones y procedimientos"
        elif "APROVISIONAMIENTO" in text_value:
            sub = "Aprovisionamientos"
        elif _contains_any(text_value, ("ADJUDICACION", "ADJUDICACIÓN")):
            sub = "Adjudicación directa"
        elif _contains_any(text_value, ("VEHICULO", "VEHÍCULO", "EQUIPO", "LUMINARIA")):
            sub = "Bienes y equipamiento"
        else:
            sub = "Servicios y contrataciones"
        return result("Adquisiciones y contrataciones", sub, 0.89)

    # 11. Proyectos e inversión.
    if _contains_any(text_value, (
        "FIMJA", "PROYECTO", "PARQUE INDUSTRIAL", "CENTRO LOGISTICO", "CENTRO LOGÍSTICO",
        "CLJ", "PUERTO SECO", "CEA", "INFRAESTRUCTURA", "OBRA PUBLICA", "OBRA PÚBLICA",
        "PLANTA DE TRATAMIENTO",
    )):
        if "FIMJA" in text_value:
            sub = "FIMJA"
        elif _contains_any(text_value, ("PARQUE INDUSTRIAL", "CENTRO LOGISTICO", "CENTRO LOGÍSTICO", "CLJ", "PUERTO SECO")):
            sub = "Parques industriales y CLJ"
        elif "CEA" in text_value or "PLANTA DE TRATAMIENTO" in text_value:
            sub = "Infraestructura hídrica"
        else:
            sub = "Proyectos estratégicos"
        return result("Proyectos e inversión", sub, 0.86)

    # 12. Coordinación institucional y enlaces.
    if _contains_any(text_value, (
        "DESIGNACION DE ENLACE", "DESIGNACIÓN DE ENLACE", "ENLACE INSTITUCIONAL",
        "SE DESIGNA ENLACE", "SOLICITUD DE INFORMACION", "SOLICITUD DE INFORMACIÓN",
        "SE SOLICITA INFORMACION", "SE SOLICITA INFORMACIÓN", "COLABORACION", "COLABORACIÓN",
        "CONVENIO DE COLABORACION", "CONVENIO DE COLABORACIÓN",
    )):
        if "ENLACE" in text_value:
            sub = "Enlaces institucionales"
        elif _contains_any(text_value, ("CONVENIO", "COLABORACION", "COLABORACIÓN")):
            sub = "Coordinación y colaboración"
        else:
            sub = "Intercambio de información"
        return result("Coordinación institucional y enlaces", sub, 0.82)

    return result("Otros / por clasificar", "Asunto genérico o insuficiente", 0.40)


def _classify_existing_official_letters(client, years=(2025, 2026)) -> dict:
    """Backfill temático para oficios ya existentes, preservando clasificación manual."""
    total = 0
    by_theme = {}
    for year in years:
        rows = (
            client.table("oficios_direccion_general")
            .select("id,asunto,dependencia,destinatario,cargo,tema,subtema,clasificacion_manual")
            .eq("anio", year)
            .execute()
            .data or []
        )
        updates = []
        for row in rows:
            if bool(row.get("clasificacion_manual")):
                continue
            cls = _classify_official_letter(row)
            updates.append({
                "id": row["id"],
                **cls,
                "clasificado_at": datetime.now().isoformat(),
            })
            by_theme[cls["tema"]] = by_theme.get(cls["tema"], 0) + 1
        for offset in range(0, len(updates), 100):
            batch = updates[offset:offset + 100]
            for item in batch:
                row_id = item.pop("id")
                client.table("oficios_direccion_general").update(item).eq("id", row_id).execute()
                total += 1
    return {"actualizados": total, "por_tema": by_theme}


def _official_theme_rows(rows: list[dict]) -> list[dict]:
    """Garantiza tema/subtema incluso antes del backfill físico."""
    themed = []
    for row in rows or []:
        copy_row = dict(row)
        if not str(copy_row.get("tema") or "").strip():
            copy_row.update(_classify_official_letter(copy_row))
        themed.append(copy_row)
    return themed


def _control_text(value) -> str:
    if value is None:
        return ""
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def _control_bool(value):
    if value is None:
        return None
    try:
        if pd.isna(value):
            return None
    except Exception:
        pass
    if isinstance(value, bool):
        return value
    text_value = str(value).strip().lower()
    if text_value in ("true", "sí", "si", "1", "x"):
        return True
    if text_value in ("false", "no", "0"):
        return False
    return None


def _control_date(value):
    if value is None:
        return None
    try:
        if pd.isna(value):
            return None
    except Exception:
        pass
    try:
        parsed = pd.to_datetime(value, errors="coerce")
        if pd.isna(parsed):
            return None
        return parsed.date().isoformat()
    except Exception:
        return None


def _office_period_from_number(office_number: str) -> tuple[int, int] | None:
    match = re.search(r"(\d{2})\D*(\d{4})\s*$", str(office_number or "").strip())
    if not match:
        return None
    month, year = int(match.group(1)), int(match.group(2))
    if month < 1 or month > 12 or year < 2024 or year > 2030:
        return None
    return month, year


def _parse_dg_control_excel(uploaded) -> tuple[list[dict], list[str]]:
    frame = pd.read_excel(uploaded, sheet_name="ENVIADOS DG", header=1, dtype=object)
    frame.columns = [str(column).strip().upper() for column in frame.columns]
    expected = ["NO", "OFICIO", "DIRIGIDO A", "CARGO", "DEPENDENCIA", "ASUNTO", "FECHA", "FIRMA", "SOLICITADO POR", "STATUS", "A. FISICO", "A.DIGITAL"]
    normalized = {name: name for name in expected if name in frame.columns}
    if "OFICIO" not in normalized:
        raise ValueError("La hoja ENVIADOS DG no contiene la columna OFICIO esperada.")
    records, warnings = [], []
    seen_control_keys = {}
    for index, row in frame.iterrows():
        excel_row = int(index) + 3
        office_number = _control_text(row.get(normalized["OFICIO"]))
        if not office_number:
            continue
        period = _office_period_from_number(office_number)
        if not period:
            warnings.append(f"Fila {excel_row}: no se pudo identificar mes/año en {office_number}.")
            continue
        month, year = period
        control_folio = _control_text(row.get(normalized.get("NO"))) if normalized.get("NO") else ""
        key_suffix = control_folio or str(excel_row)
        base_control_key = f"ENVIADOS DG|{year}|{key_suffix}|{office_number}".upper()
        duplicate_number = seen_control_keys.get(base_control_key, 0) + 1
        seen_control_keys[base_control_key] = duplicate_number
        control_key = base_control_key if duplicate_number == 1 else f"{base_control_key}|DUP{duplicate_number}"
        if duplicate_number > 1:
            warnings.append(
                f"Fila {excel_row}: el oficio {office_number} comparte número/control con otro registro; "
                f"se conservará como registro independiente (duplicado {duplicate_number})."
            )
        base_record = {
            "anio": year,
            "mes": month,
            "numero_oficio": office_number,
            "folio_control": control_folio or None,
            "destinatario": (_control_text(row.get(normalized.get("DIRIGIDO A"))) or None) if normalized.get("DIRIGIDO A") else None,
            "cargo": (_control_text(row.get(normalized.get("CARGO"))) or None) if normalized.get("CARGO") else None,
            "dependencia": (_control_text(row.get(normalized.get("DEPENDENCIA"))) or None) if normalized.get("DEPENDENCIA") else None,
            "asunto": (_control_text(row.get(normalized.get("ASUNTO"))) or None) if normalized.get("ASUNTO") else None,
            "fecha_control": _control_date(row.get(normalized.get("FECHA"))) if normalized.get("FECHA") else None,
            "firma": (_control_text(row.get(normalized.get("FIRMA"))) or None) if normalized.get("FIRMA") else None,
            "solicitado_por": (_control_text(row.get(normalized.get("SOLICITADO POR"))) or None) if normalized.get("SOLICITADO POR") else None,
            "status_control": (_control_text(row.get(normalized.get("STATUS"))) or None) if normalized.get("STATUS") else None,
            "archivo_fisico": _control_bool(row.get(normalized.get("A. FISICO"))) if normalized.get("A. FISICO") else None,
            "archivo_digital": _control_bool(row.get(normalized.get("A.DIGITAL"))) if normalized.get("A.DIGITAL") else None,
            "origen": "control_excel",
            "hoja_origen": "ENVIADOS DG",
            "fila_origen": excel_row,
            "clave_control": control_key,
        }
        classification = _classify_official_letter(base_record)
        base_record.update(classification)
        base_record["clasificado_at"] = datetime.now().isoformat()
        base_record["clasificacion_manual"] = False
        records.append(base_record)
    return records, warnings


def _ingest_dg_control(client, uploaded) -> tuple[dict, list[str]]:
    raw = uploaded.getvalue()
    author_name = st.session_state.user.get("nombre") or st.session_state.user.get("email")
    ingestion = client.table("ingestas_oficios_dg").insert({
        "nombre_archivo": uploaded.name,
        "hash_archivo": hashlib.sha256(raw).hexdigest(),
        "hoja": "ENVIADOS DG",
        "estado": "Procesando",
        "ingestado_por": st.session_state.user["id"],
        "autor_nombre": author_name,
    }).execute().data[0]
    ingestion_id = str(ingestion["id"])
    try:
        records, warnings = _parse_dg_control_excel(io.BytesIO(raw))
        existing_rows = (client.table("oficios_direccion_general").select("clave_control").execute().data or [])
        existing_keys = {row.get("clave_control") for row in existing_rows if row.get("clave_control")}
        new_count = sum(1 for row in records if row["clave_control"] not in existing_keys)
        updated_count = len(records) - new_count
        payload = [{**record, "ingesta_id": ingestion_id, "registrado_por": st.session_state.user["id"], "registrado_por_nombre": author_name} for record in records]
        for offset in range(0, len(payload), 100):
            client.table("oficios_direccion_general").upsert(payload[offset:offset + 100], on_conflict="clave_control").execute()
        summary = {
            "registros_detectados": len(records), "registros_nuevos": new_count,
            "registros_actualizados": updated_count, "registros_omitidos": len(warnings),
            "estado": "Completada", "completed_at": datetime.now().isoformat(),
        }
        client.table("ingestas_oficios_dg").update(summary).eq("id", ingestion_id).execute()
        return {**ingestion, **summary}, warnings
    except Exception as exc:
        client.table("ingestas_oficios_dg").update({"estado": "Error", "detalle_error": str(exc), "completed_at": datetime.now().isoformat()}).eq("id", ingestion_id).execute()
        raise


def _latest_dg_ingestion(client):
    try:
        rows = client.table("ingestas_oficios_dg").select("*").order("created_at", desc=True).limit(1).execute().data or []
        return rows[0] if rows else None
    except Exception:
        return None


def _delete_official_letter(client, document: dict) -> None:
    if document.get("ruta_storage"):
        _remove_storage_paths(client, [document["ruta_storage"]])
    client.table("oficios_direccion_general").delete().eq("id", document["id"]).execute()


def _attach_signed_office(client, document: dict, uploaded) -> None:
    old_path = document.get("ruta_storage")
    path = f"oficios_direccion_general/{document['anio']}/{int(document['mes']):02d}/{uuid.uuid4().hex}_{safe_name(uploaded.name)}"
    _upload_junta_document(client, path, uploaded)
    client.table("oficios_direccion_general").update({
        "nombre_archivo": uploaded.name, "ruta_storage": path, "mime_type": uploaded.type,
        "tamano_bytes": uploaded.size, "subido_por": st.session_state.user["id"],
        "autor_nombre": st.session_state.user.get("nombre") or st.session_state.user.get("email"),
        "updated_at": datetime.now().isoformat(),
    }).eq("id", document["id"]).execute()
    if old_path and old_path != path:
        _remove_storage_paths(client, [old_path])


def _official_letter_card(client, document: dict):
    filename = document.get("nombre_archivo") or ""
    office_number = document.get("numero_oficio") or "Sin número"
    subject = document.get("asunto") or "Sin asunto en control"
    month_name = dict(MONTHS_ES).get(int(document.get("mes") or 0), "")
    date_label = str(document.get("fecha_oficio") or "")[:10] or f"{month_name} {document.get('anio') or ''}".strip()
    recipient = document.get("destinatario") or "Sin destinatario"
    has_supabase_file = bool(document.get("ruta_storage"))
    has_drive_link = bool(document.get("drive_url"))

    with st.container(border=True):
        info, actions = st.columns([4.6, 2.1], vertical_alignment="center")
        info.markdown(f"**{html.escape(office_number)} · {html.escape(subject)}**")
        info.caption(
            f"{html.escape(date_label)} · Destinatario: {html.escape(recipient)}  \n"
            f"{html.escape(document.get('cargo') or 'Cargo no registrado')} · "
            f"{html.escape(document.get('dependencia') or 'Dependencia no registrada')}"
        )
        status = document.get("status_control") or "Sin estatus"
        source = "Control ENVIADOS DG" if document.get("origen") == "control_excel" else "Registro manual"
        info.caption(
            f"{source} · Estatus: {html.escape(status)} · "
            f"Solicitado por: {html.escape(document.get('solicitado_por') or 'Sin dato')}"
        )

        # Tema y subtema: usa la clasificación guardada; si aún no existe,
        # la calcula en memoria para no dejar el oficio sin categoría visual.
        themed_document = dict(document)
        if not str(themed_document.get("tema") or "").strip():
            themed_document.update(_classify_official_letter(themed_document))
        theme = themed_document.get("tema") or "Otros / por clasificar"
        subtheme = themed_document.get("subtema") or "Sin subtema"

        info.markdown(
            f"""
            <div style="display:flex;gap:8px;flex-wrap:wrap;margin-top:8px;">
                <span style="
                    background:#eaf4fb;color:#175a84;border:1px solid #cfe5f2;
                    border-radius:999px;padding:5px 10px;font-size:12px;font-weight:700;">
                    Tema: {html.escape(str(theme))}
                </span>
                <span style="
                    background:#f3eef9;color:#694d91;border:1px solid #e2d8f0;
                    border-radius:999px;padding:5px 10px;font-size:12px;font-weight:700;">
                    Subtema: {html.escape(str(subtheme))}
                </span>
            </div>
            """,
            unsafe_allow_html=True,
        )

        if has_drive_link:
            actions.success("Firmado en Drive")
            actions.link_button(
                "Ver oficio en Drive",
                document["drive_url"],
                use_container_width=True,
            )
            if document.get("drive_nombre_archivo"):
                actions.caption(document["drive_nombre_archivo"])
        elif has_supabase_file:
            try:
                data = client.storage.from_("expedientes").download(document["ruta_storage"])
                action_cols = actions.columns(3 if is_master_admin() else 2)
                show = action_cols[0].toggle("Ver", key=f"official_view_{document['id']}")
                action_cols[1].download_button(
                    "Descargar", data, file_name=filename or "oficio_firmado",
                    mime=document.get("mime_type") or "application/octet-stream",
                    key=f"official_download_{document['id']}", use_container_width=True
                )
                if is_master_admin() and action_cols[2].button(
                    "Eliminar", key=f"official_delete_{document['id']}", use_container_width=True
                ):
                    _delete_official_letter(client, document)
                    st.rerun()
                if show and not _document_preview(data, filename, 650):
                    st.info("La vista previa no está disponible para este formato.")
            except Exception:
                actions.caption("El archivo está registrado, pero no fue posible recuperarlo.")
        else:
            actions.warning("Firmado pendiente")

        with st.expander("Expediente del oficio", expanded=False):
            if has_drive_link:
                st.caption("Documento vinculado desde Google Drive.")
                st.write(document.get("drive_url"))
                if document.get("drive_vinculado_at"):
                    st.caption(
                        f"Vínculo actualizado: "
                        f"{str(document['drive_vinculado_at'])[:16].replace('T', ' · ')}"
                    )
            elif not has_supabase_file:
                st.caption("No hay documento firmado vinculado todavía.")

            # Carga manual sigue disponible como respaldo.
            if not has_drive_link:
                signed_upload = st.file_uploader(
                    "Carga manual de respaldo (opcional)",
                    type=["pdf", "docx", "jpg", "jpeg", "png"],
                    key=f"signed_office_{document['id']}",
                )
                if st.button(
                    "Guardar firmado manualmente",
                    key=f"save_signed_office_{document['id']}",
                    disabled=not signed_upload,
                    use_container_width=True,
                ):
                    _attach_signed_office(client, document, signed_upload)
                    st.success("Oficio firmado guardado.")
                    st.rerun()




def _normalize_folio_key(value) -> str:
    """Normaliza el folio de control sin alterar su identidad lógica."""
    if value is None:
        return ""
    s = str(value).strip().upper()
    if not s:
        return ""

    # Excel puede entregar 90 como 90.0
    s = re.sub(r"\.0$", "", s)

    # Normaliza espacios y variantes de BIS.
    s = re.sub(r"\s+", "", s)
    s = s.replace("_", "-")
    s = re.sub(r"-+", "-", s)

    # 064BIS / 64-BIS / 64 BIS -> 64-BIS
    m = re.fullmatch(r"0*(\d+)-?BIS", s)
    if m:
        return f"{int(m.group(1))}-BIS"

    # Folio numérico simple: 006 -> 6
    if re.fullmatch(r"0*\d+", s):
        try:
            return str(int(s))
        except Exception:
            return s

    return s



def _canonical_office_number(row: dict) -> str:
    """Convierte variantes del número de oficio a una clave lógica común.

    Ejemplos equivalentes:
    DG-E-233-08-2025
    DGE-233-082025
    DG.E.233.08.2025
    DG/E-233-08/2025
    """
    raw = str(row.get("numero_oficio") or "").strip().upper()
    if not raw:
        return ""

    # Quita sufijos internos creados por ingestas para aceptar duplicados técnicos.
    raw = re.sub(r"\|DUP\d+$", "", raw)

    # Homologa separadores y familias DGE / DGI.
    s = raw
    s = s.replace("\\", "/")
    s = re.sub(r"\s+", "", s)
    s = s.replace("_", "-")
    s = s.replace(".", "-")
    s = s.replace("/", "-")
    s = re.sub(r"-+", "-", s)

    # DGE / DGI / DG1 y variantes.
    s = re.sub(r"^DG-?E", "DG-E", s)
    s = re.sub(r"^DG-?I", "DG-I", s)
    s = re.sub(r"^DG-?1", "DG-I", s)

    # Busca tipo.
    type_code = ""
    if s.startswith("DG-E"):
        type_code = "E"
        tail = s[4:]
    elif s.startswith("DG-I"):
        type_code = "I"
        tail = s[4:]
    else:
        # Respaldo para capturas extrañas: DGE..., DGI...
        mtype = re.match(r"^DG([EI])", s)
        if mtype:
            type_code = mtype.group(1)
            tail = s[mtype.end():]
        else:
            return ""

    tail = tail.strip("-")

    # Año: preferir el que venga en el número; si no, usar la columna anio.
    year_match = re.search(r"(20\d{2})$", tail)
    if year_match:
        year = year_match.group(1)
        before_year = tail[:year_match.start()].strip("-")
    else:
        year = str(row.get("anio") or "").strip()
        if not re.fullmatch(r"20\d{2}", year):
            return ""
        before_year = tail

    # Normaliza formatos MMYYYY pegados, por ejemplo 082025.
    before_year = re.sub(rf"(\d{{2}}){re.escape(year)}$", r"\1", before_year)

    # Extrae todas las partes numéricas / BIS antes del año.
    parts = [p for p in before_year.split("-") if p]
    if not parts:
        return ""

    # Consecutivo y BIS.
    seq = None
    bis = False
    month = None

    # Casos como 233-BIS-08 o 233-08.
    for i, p in enumerate(parts):
        if seq is None and re.fullmatch(r"\d{1,4}", p):
            seq = str(int(p))
            if i + 1 < len(parts) and parts[i + 1] == "BIS":
                bis = True
                # mes después de BIS
                if i + 2 < len(parts) and re.fullmatch(r"\d{1,2}", parts[i + 2]):
                    month = f"{int(parts[i + 2]):02d}"
            elif i + 1 < len(parts) and re.fullmatch(r"\d{1,2}", parts[i + 1]):
                month = f"{int(parts[i + 1]):02d}"
            break

    # Casos compactos tipo 233BIS o 233-BIS.
    if seq is None:
        compact = re.search(r"(\d{1,4})(BIS)?", before_year)
        if compact:
            seq = str(int(compact.group(1)))
            bis = bool(compact.group(2))

    # Fallback de mes: último número de 1-2 dígitos antes del año, distinto del consecutivo.
    if month is None:
        nums = re.findall(r"(?<!\d)(0?[1-9]|1[0-2])(?!\d)", before_year)
        if nums:
            month = f"{int(nums[-1]):02d}"

    if not seq or not month:
        return ""

    seq_part = f"{seq}-BIS" if bis else seq
    return f"DG/{type_code}-{seq_part}-{month}/{year}"


def _office_unique_key(row: dict) -> str:
    """Identidad lógica de un oficio para conteos y analítica."""
    canonical = _canonical_office_number(row)
    if canonical:
        return f"NUM|{canonical}"

    # Respaldo sólo cuando el número no puede reconstruirse.
    year = str(row.get("anio") or "").strip()
    folio = _normalize_folio_key(row.get("folio_control"))
    if folio:
        return f"FOLIO|{year}|{folio}"

    number = _normalize_office_link_key(row.get("numero_oficio"))
    if number:
        number = re.sub(r"\|DUP\d+$", "", number)
        return f"RAW|{year}|{number}"

    return f"ID|{row.get('id') or id(row)}"


def _office_row_score(row: dict) -> int:
    """Prefiere, entre duplicados, la fila con mayor información y vínculo firmado."""
    score = 0
    if row.get("drive_url"):
        score += 100
    if row.get("ruta_storage"):
        score += 80
    for field in (
        "destinatario", "dependencia", "asunto", "solicitado_por",
        "fecha_control", "firma", "status_control", "folio_control"
    ):
        if str(row.get(field) or "").strip():
            score += 1
    return score


def _dedupe_office_rows(rows: list[dict]) -> list[dict]:
    """Devuelve una sola fila por oficio lógico."""
    chosen = {}
    order = []
    for row in rows or []:
        key = _office_unique_key(row)
        if key not in chosen:
            chosen[key] = row
            order.append(key)
        elif _office_row_score(row) > _office_row_score(chosen[key]):
            chosen[key] = row
    return [chosen[key] for key in order]



def _official_status_text(row: dict) -> str:
    for field in ("status_control", "estatus", "status", "firma"):
        value = str(row.get(field) or "").strip()
        if value:
            return value
    return ""


def _official_is_cancelled(row: dict) -> bool:
    return "CANCEL" in _official_status_text(row).upper()


def _official_active_rows(rows: list[dict]) -> list[dict]:
    return [row for row in (rows or []) if not _official_is_cancelled(row)]



def official_letters_month(year: int, month: int, month_name: str):
    top1, top2 = st.columns([1, 5])
    top1.button(
        "← Meses",
        use_container_width=True,
        key=f"official_back_months_{year}_{month}",
        on_click=_back_official_months,
    )
    top2.markdown(f"## Oficios Dirección General · {month_name} {year}")

    if not configured():
        st.error("Primero debes conectar Supabase.")
        return

    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
    rows = (client.table("oficios_direccion_general").select("*")
            .eq("anio", year).eq("mes", month)
            .order("fila_origen").order("created_at").execute().data or [])
    rows = _dedupe_office_rows(rows)

    unique_recipients = len({
        str(row.get("destinatario") or "").strip().casefold()
        for row in rows if str(row.get("destinatario") or "").strip()
    })
    unique_dependencies = len({
        str(row.get("dependencia") or "").strip().casefold()
        for row in rows if str(row.get("dependencia") or "").strip()
    })
    unique_requesters = len({
        str(row.get("solicitado_por") or "").strip().casefold()
        for row in rows if str(row.get("solicitado_por") or "").strip()
    })
    metrics_html = (
        '<div class="metric-grid">'
        f'<div class="metric-box metric-blue"><div class="metric-label">Oficios registrados</div><div class="metric-value">{len(rows)}</div></div>'
        f'<div class="metric-box metric-green"><div class="metric-label">Destinatarios únicos</div><div class="metric-value">{unique_recipients}</div></div>'
        f'<div class="metric-box metric-orange"><div class="metric-label">Dependencias únicas</div><div class="metric-value">{unique_dependencies}</div></div>'
        f'<div class="metric-box metric-purple"><div class="metric-label">Solicitado por</div><div class="metric-value">{unique_requesters}</div></div>'
        '</div>'
    )
    st.markdown(metrics_html, unsafe_allow_html=True)

    # Selector de vista
    if "official_list_mode" not in st.session_state:
        st.session_state.official_list_mode = False

    mode_col, _ = st.columns([1.4, 4.6])
    if st.session_state.official_list_mode:
        if mode_col.button("Vista modo tarjetas", key=f"official_cards_mode_{year}_{month}",
                           use_container_width=True, type="primary"):
            st.session_state.official_list_mode = False
            st.rerun()
    else:
        if mode_col.button("Vista modo lista", key=f"official_list_mode_{year}_{month}",
                           use_container_width=True, type="primary"):
            st.session_state.official_list_mode = True
            st.rerun()

    term = st.text_input(
        "Buscar dentro del mes",
        placeholder="Número, asunto, destinatario, dependencia, solicitante, tema, subtema…",
        key=f"official_search_{year}_{month}",
    ).strip().lower()

    filtered = rows
    if term:
        filtered = [
            row for row in rows
            if term in " ".join([
                str(row.get("numero_oficio") or ""),
                str(row.get("folio_control") or ""),
                str(row.get("asunto") or ""),
                str(row.get("destinatario") or ""),
                str(row.get("cargo") or ""),
                str(row.get("dependencia") or ""),
                str(row.get("solicitado_por") or ""),
                str(row.get("status_control") or ""),
                str(row.get("tema") or ""),
                str(row.get("subtema") or ""),
            ]).lower()
        ]

    st.markdown(f"### Oficios del mes · {len(filtered)}")
    if not filtered:
        st.info(f"No hay registros que mostrar en {month_name} de {year}.")
        return

    if st.session_state.official_list_mode:
        list_rows = []
        for row in filtered:
            list_rows.append({
                "Número de oficio": row.get("numero_oficio") or "Sin número",
                "Título / asunto": row.get("asunto") or "Sin asunto",
                "Mes": f"{month_name} {year}",
                "Destinatario": row.get("destinatario") or "Sin destinatario",
                "Dependencia / organización destinataria": row.get("dependencia") or "Sin información",
                "Tema": row.get("tema") or _classify_official_letter(row).get("tema") or "Otros / por clasificar",
                "Subtema": row.get("subtema") or _classify_official_letter(row).get("subtema") or "Sin subtema",
            })
        st.dataframe(
            pd.DataFrame(list_rows),
            use_container_width=True,
            hide_index=True,
            column_config={
                "Número de oficio": st.column_config.TextColumn("Número de oficio", width="medium"),
                "Título / asunto": st.column_config.TextColumn("Título / asunto", width="large"),
                "Mes": st.column_config.TextColumn("Mes", width="small"),
                "Destinatario": st.column_config.TextColumn("Destinatario", width="medium"),
                "Dependencia / organización destinataria": st.column_config.TextColumn(
                    "Dependencia / organización destinataria", width="large"
                ),
            },
        )
        st.caption("Para abrir el expediente de un oficio o adjuntar el firmado, cambia a Vista modo tarjetas.")
    else:
        for row in filtered:
            _official_letter_card(client, row)


def _go_official_year(year: int):
    st.session_state["official_year"] = int(year)
    st.session_state.pop("official_month", None)
    st.session_state.pop("official_list_mode", None)


def _go_official_month(month: int):
    st.session_state["official_month"] = int(month)
    st.session_state.pop("official_list_mode", None)


def _back_official_years():
    st.session_state.pop("official_year", None)
    st.session_state.pop("official_month", None)
    st.session_state.pop("official_list_mode", None)


def _back_official_months():
    st.session_state.pop("official_month", None)
    st.session_state.pop("official_list_mode", None)

def _official_group_counts(rows: list[dict], field: str, empty_label: str) -> pd.DataFrame:
    groups: dict[str, dict] = {}
    for row in rows:
        raw = str(row.get(field) or "").strip()
        label = re.sub(r"\s+", " ", raw) if raw else empty_label
        key = label.casefold()
        if key not in groups:
            groups[key] = {"Etiqueta": label, "Oficios": 0}
        groups[key]["Oficios"] += 1
    return pd.DataFrame(groups.values()).sort_values(
        ["Oficios", "Etiqueta"], ascending=[False, True]
    ) if groups else pd.DataFrame(columns=["Etiqueta", "Oficios"])



st.markdown("""
<style>
/* --- Analítica Oficios: integración visual con fondo general --- */
div[data-testid="stAltairChart"] {
    background: transparent !important;
}

div[data-testid="stAltairChart"] > div {
    background: transparent !important;
}

/* Fuerza que el contenedor del tablero no parezca una "hoja blanca pegada". */
[data-testid="stVerticalBlock"]:has(.oficios-dashboard-marker) {
    background: transparent !important;
}

/* Alineación general de analítica */
.oficios-analytics-wrap {
    width: 100%;
    margin: 0;
    padding: 0;
}

/* Tarjetas superiores: ocupar todo el ancho disponible */
.oficios-metric-grid {
    display: grid;
    grid-template-columns: repeat(4, minmax(0, 1fr));
    gap: 14px;
    width: 100%;
    margin: 0 0 28px 0;
}

.oficios-metric-card {
    width: 100%;
    min-width: 0;
    box-sizing: border-box;
}

@media (max-width: 1100px) {
    .oficios-metric-grid {
        grid-template-columns: repeat(2, minmax(0, 1fr));
    }
}
@media (max-width: 700px) {
    .oficios-metric-grid {
        grid-template-columns: 1fr;
    }
}
</style>
""", unsafe_allow_html=True)


def _transparent_altair(chart):
    """Integra los gráficos con el fondo gris de Streamlit."""
    try:
        return (
            chart
            .configure(background="transparent")
            .configure_view(strokeWidth=0, fill="transparent")
        )
    except Exception:
        return chart


def _official_letters_analytics(year: int, rows: list[dict]):
    # Mantiene la analítica original y añade un drill-down independiente
    # para los oficios cancelados.
    rows = _dedupe_office_rows(rows)

    cancelled_rows = [row for row in rows if _official_is_cancelled(row)]
    active_rows = _official_active_rows(rows)

    def _render_analytics_block(data_rows: list[dict], title: str, subtitle: str, cancelled_mode: bool = False):
        month_names = {month: name for month, name in MONTHS_ES}
        month_order = [name for _, name in MONTHS_ES]
        month_counts = {month: 0 for month, _ in MONTHS_ES}

        for row in data_rows:
            try:
                month_value = int(row.get("mes") or 0)
            except (TypeError, ValueError):
                month_value = 0
            if month_value in month_counts:
                month_counts[month_value] += 1

        monthly = pd.DataFrame([
            {"Mes": month_names[month], "Oficios": month_counts[month], "Orden": month}
            for month, _ in MONTHS_ES
        ])

        recipients = _official_group_counts(data_rows, "destinatario", "Sin destinatario")
        requesters = _official_group_counts(data_rows, "solicitado_por", "Sin dato")

        total = len(data_rows)
        active_months = sum(1 for value in month_counts.values() if value > 0)
        avg = (total / active_months) if active_months else 0
        unique_recipients = max(
            0,
            len(recipients) - int("Sin destinatario" in recipients["Etiqueta"].values),
        )
        unique_requesters = max(
            0,
            len(requesters) - int("Sin dato" in requesters["Etiqueta"].values),
        )

        banner_class = "analytics-banner"
        st.markdown(
            f"""<div class="{banner_class}">
            <h3>{html.escape(title)}</h3>
            <p>{html.escape(subtitle)}</p>
            </div>""",
            unsafe_allow_html=True,
        )

        metric_cards = [
            ("Oficios enviados" if not cancelled_mode else "Total cancelados", total, "#0798cf" if not cancelled_mode else "#f68b08"),
            ("Promedio por mes activo", f"{avg:.1f}", "#009b4c"),
            ("Destinatarios únicos", unique_recipients, "#16ad8f"),
            ("Solicitantes únicos", unique_requesters, "#a990c7"),
        ]
        cards = "".join(
            f"""<div class="analytics-metric oficios-metric-card" style="--tone:{tone}">
            <div class="analytics-value">{value}</div>
            <div class="analytics-label">{html.escape(label)}</div></div>"""
            for label, value, tone in metric_cards
        )
        st.markdown(
            f'<div class="oficios-metric-grid">{cards}</div>',
            unsafe_allow_html=True,
        )

        st.markdown("### Tablero interactivo de oficios" if not cancelled_mode else "### Tablero interactivo de cancelados")
        st.caption(
            "Haz clic en cualquier barra para usarla como filtro de las demás gráficas. "
            "Puedes combinar filtros; haz doble clic sobre una selección para limpiarla."
        )

        # Todas las visualizaciones comparten un mismo conjunto de datos y selecciones.
        # Esto permite que Mes, Destinatario, Dependencia, Solicitado por y Tema
        # funcionen como filtros cruzados entre sí.
        dashboard_rows = _official_theme_rows(data_rows)
        month_names = {month: name for month, name in MONTHS_ES}
        month_order = [name for _, name in MONTHS_ES]

        dashboard_df = pd.DataFrame([
            {
                "Mes": month_names.get(int(row.get("mes") or 0), "Sin mes"),
                "Destinatario": str(row.get("destinatario") or "Sin destinatario").strip() or "Sin destinatario",
                "Dependencia": str(row.get("dependencia") or "Sin dependencia").strip() or "Sin dependencia",
                "Solicitado por": str(row.get("solicitado_por") or "Sin dato").strip() or "Sin dato",
                "Tema": str(row.get("tema") or "Otros / por clasificar").strip() or "Otros / por clasificar",
            }
            for row in dashboard_rows
        ])

        if dashboard_df.empty:
            st.info("No hay información suficiente para generar el tablero interactivo.")
        else:
            # Dominios top para evitar gráficas ilegibles, manteniendo el filtrado sobre datos fila a fila.
            top_recipients = (
                dashboard_df.groupby("Destinatario").size().sort_values(ascending=False).head(12).index.tolist()
            )
            top_dependencies = (
                dashboard_df.groupby("Dependencia").size().sort_values(ascending=False).head(12).index.tolist()
            )
            top_requesters = (
                dashboard_df.groupby("Solicitado por").size().sort_values(ascending=False).head(12).index.tolist()
            )
            top_themes = (
                dashboard_df.groupby("Tema").size().sort_values(ascending=False).head(12).index.tolist()
            )

            # Selecciones persistentes por clic y resaltado por hover.
            month_sel = alt.selection_point(fields=["Mes"], on="click", clear="dblclick", empty=True, name="mes_sel")
            rec_sel = alt.selection_point(fields=["Destinatario"], on="click", clear="dblclick", empty=True, name="dest_sel")
            dep_sel = alt.selection_point(fields=["Dependencia"], on="click", clear="dblclick", empty=True, name="dep_sel")
            req_sel = alt.selection_point(fields=["Solicitado por"], on="click", clear="dblclick", empty=True, name="req_sel")
            theme_sel = alt.selection_point(fields=["Tema"], on="click", clear="dblclick", empty=True, name="tema_sel")

            month_hover = alt.selection_point(fields=["Mes"], on="pointerover", clear="pointerout", empty=False, name="mes_hover")
            rec_hover = alt.selection_point(fields=["Destinatario"], on="pointerover", clear="pointerout", empty=False, name="dest_hover")
            dep_hover = alt.selection_point(fields=["Dependencia"], on="pointerover", clear="pointerout", empty=False, name="dep_hover")
            req_hover = alt.selection_point(fields=["Solicitado por"], on="pointerover", clear="pointerout", empty=False, name="req_hover")
            theme_hover = alt.selection_point(fields=["Tema"], on="pointerover", clear="pointerout", empty=False, name="tema_hover")

            base = alt.Chart(dashboard_df)

            month_base = (
                base
                .transform_filter(rec_sel)
                .transform_filter(dep_sel)
                .transform_filter(req_sel)
                .transform_filter(theme_sel)
                .transform_aggregate(Oficios="count()", groupby=["Mes"])
            )
            month_bars = (
                month_base.mark_bar(cornerRadiusTopLeft=8, cornerRadiusTopRight=8, size=54)
                .encode(
                    x=alt.X("Mes:N", sort=month_order, title=None, axis=alt.Axis(labelAngle=0)),
                    y=alt.Y("Oficios:Q", title="Número de oficios", axis=alt.Axis(tickMinStep=1)),
                    color=alt.condition(
                        month_sel | month_hover,
                        alt.Color("Mes:N", sort=month_order, legend=None),
                        alt.value("#c3cbd1"),
                    ),
                    opacity=alt.condition(month_sel | month_hover, alt.value(1), alt.value(0.8)),
                    tooltip=[alt.Tooltip("Mes:N"), alt.Tooltip("Oficios:Q", format=".0f")],
                )
                .add_params(month_sel, month_hover)
            )
            month_labels = month_base.mark_text(dy=-10, fontSize=12, fontWeight="bold", color="#35434b").encode(
                x=alt.X("Mes:N", sort=month_order),
                y="Oficios:Q",
                text=alt.Text("Oficios:Q", format=".0f"),
            )
            month_chart = (month_bars + month_labels).properties(
                height=300,
                title=alt.TitleParams("Oficios por mes", anchor="start", fontSize=18, fontWeight="bold"),
            )

            def horizontal_chart(field, domain, selection, hover, title, color, filters):
                chart = base
                for f in filters:
                    chart = chart.transform_filter(f)
                chart = (
                    chart
                    .transform_filter(alt.FieldOneOfPredicate(field=field, oneOf=domain))
                    .transform_aggregate(Oficios="count()", groupby=[field])
                )
                bars = (
                    chart.mark_bar(cornerRadiusEnd=8)
                    .encode(
                        y=alt.Y(f"{field}:N", sort="-x", title=None, axis=alt.Axis(labelLimit=210)),
                        x=alt.X("Oficios:Q", title="Número de oficios", axis=alt.Axis(tickMinStep=1)),
                        color=alt.condition(selection | hover, alt.value(color), alt.value("#c7cdd2")),
                        opacity=alt.condition(selection | hover, alt.value(1), alt.value(0.78)),
                        tooltip=[alt.Tooltip(f"{field}:N", title=field), alt.Tooltip("Oficios:Q", format=".0f")],
                    )
                    .add_params(selection, hover)
                )
                labels = chart.mark_text(align="left", baseline="middle", dx=5, fontSize=11, fontWeight="bold", color="#35434b").encode(
                    y=alt.Y(f"{field}:N", sort="-x"),
                    x="Oficios:Q",
                    text=alt.Text("Oficios:Q", format=".0f"),
                )
                return (bars + labels).properties(
                    height=320,
                    width=310,
                    title=alt.TitleParams(title, anchor="start", fontSize=16, fontWeight="bold"),
                )

            recipient_chart = horizontal_chart(
                "Destinatario", top_recipients, rec_sel, rec_hover, "Principales destinatarios",
                "#173b63", [month_sel, dep_sel, req_sel, theme_sel],
            )
            dependency_chart = horizontal_chart(
                "Dependencia", top_dependencies, dep_sel, dep_hover, "Dependencias destinatarias",
                "#0a9b78", [month_sel, rec_sel, req_sel, theme_sel],
            )
            requester_chart = horizontal_chart(
                "Solicitado por", top_requesters, req_sel, req_hover, "Solicitado por",
                "#6750a4", [month_sel, rec_sel, dep_sel, theme_sel],
            )
            theme_chart = horizontal_chart(
                "Tema", top_themes, theme_sel, theme_hover, "Temas",
                "#f68b08", [month_sel, rec_sel, dep_sel, req_sel],
            )

            dashboard = alt.vconcat(
                month_chart,
                alt.hconcat(recipient_chart, dependency_chart).resolve_scale(x="independent"),
                alt.hconcat(requester_chart, theme_chart).resolve_scale(x="independent"),
                spacing=28,
            ).configure_view(stroke=None)

            st.altair_chart(_transparent_altair(dashboard), use_container_width=True)

    # Analítica principal: sólo oficios vigentes/no cancelados.
    # Los cancelados se muestran exclusivamente en su tarjeta y drill-down.
    _render_analytics_block(
        active_rows,
        f"Analítica de Oficios · {year}",
        "Volumen mensual, destinatarios y origen de los oficios enviados.",
        cancelled_mode=False,
    )


    # Analítica temática de los oficios enviados.
    active_rows = _official_theme_rows(active_rows)
    st.markdown("---")
    st.markdown("## Temática de los oficios")
    st.caption("Clasificación automática a partir del asunto y contexto del oficio. Los cancelados no se incluyen en este bloque.")

    theme_counts = {}
    for row in active_rows:
        theme = str(row.get("tema") or "Otros / por clasificar").strip()
        theme_counts[theme] = theme_counts.get(theme, 0) + 1

    theme_df = pd.DataFrame(
        [{"Tema": key, "Oficios": value} for key, value in theme_counts.items()]
    ).sort_values(["Oficios", "Tema"], ascending=[False, True])

    if not theme_df.empty:
        ordered_themes = theme_df["Tema"].tolist()[::-1]
        theme_chart = (
            alt.Chart(theme_df)
            .mark_bar(cornerRadiusEnd=8)
            .encode(
                y=alt.Y("Tema:N", sort=ordered_themes, title=None, axis=alt.Axis(labelLimit=320)),
                x=alt.X("Oficios:Q", title="Número de oficios", axis=alt.Axis(tickMinStep=1)),
                color=alt.value("#173b63"),
                tooltip=[
                    alt.Tooltip("Tema:N"),
                    alt.Tooltip("Oficios:Q", format=".0f"),
                ],
            )
        )
        theme_labels = (
            alt.Chart(theme_df)
            .mark_text(align="left", baseline="middle", dx=6, fontWeight="bold", color="#35434b")
            .encode(
                y=alt.Y("Tema:N", sort=ordered_themes),
                x="Oficios:Q",
                text=alt.Text("Oficios:Q", format=".0f"),
            )
        )
        st.altair_chart((theme_chart + theme_labels).properties(height=460), use_container_width=True)

        theme_options = theme_df["Tema"].tolist()
        selected_theme = st.selectbox(
            "Explorar tema",
            theme_options,
            key=f"official_theme_select_{year}",
        )
        selected_rows = [row for row in active_rows if str(row.get("tema") or "").strip() == selected_theme]

        t1, t2, t3, t4 = st.columns(4)
        theme_total = len(selected_rows)
        share = (theme_total / len(active_rows) * 100) if active_rows else 0
        active_theme_months = len({
            int(row.get("mes") or 0)
            for row in selected_rows
            if int(row.get("mes") or 0) in dict(MONTHS_ES)
        })
        subthemes = len({
            str(row.get("subtema") or "").strip()
            for row in selected_rows
            if str(row.get("subtema") or "").strip()
        })
        t1.metric("Oficios del tema", theme_total)
        t2.metric("Participación anual", f"{share:.1f}%")
        t3.metric("Meses con actividad", active_theme_months)
        t4.metric("Subtemas", subthemes)

        month_names = {month: name for month, name in MONTHS_ES}
        theme_monthly = []
        for month, month_name in MONTHS_ES:
            theme_monthly.append({
                "Mes": month_name,
                "Oficios": sum(1 for row in selected_rows if int(row.get("mes") or 0) == month),
            })
        theme_month_df = pd.DataFrame(theme_monthly)

        left_theme, right_theme = st.columns(2, gap="large")
        with left_theme:
            st.markdown("### Evolución mensual")
            st.bar_chart(theme_month_df.set_index("Mes")["Oficios"], use_container_width=True)

            sub_counts = {}
            for row in selected_rows:
                sub = str(row.get("subtema") or "Sin subtema").strip() or "Sin subtema"
                sub_counts[sub] = sub_counts.get(sub, 0) + 1
            sub_df = pd.DataFrame(
                [{"Subtema": key, "Oficios": value} for key, value in sub_counts.items()]
            ).sort_values("Oficios", ascending=False)
            st.markdown("### Subtemas")
            if sub_df.empty:
                st.info("No hay subtemas para esta categoría.")
            else:
                st.bar_chart(sub_df.set_index("Subtema")["Oficios"], use_container_width=True)

        with right_theme:
            st.markdown("### Principales destinatarios")
            selected_recipients = _official_group_counts(selected_rows, "destinatario", "Sin destinatario").head(10)
            if selected_recipients.empty:
                st.info("No hay destinatarios para este tema.")
            else:
                st.bar_chart(
                    selected_recipients.set_index("Etiqueta")["Oficios"],
                    use_container_width=True,
                )

            st.markdown("### Solicitado por")
            selected_requesters = _official_group_counts(selected_rows, "solicitado_por", "Sin dato").head(10)
            if selected_requesters.empty:
                st.info("No hay información de 'Solicitado por' para este tema.")
            else:
                st.bar_chart(
                    selected_requesters.set_index("Etiqueta")["Oficios"],
                    use_container_width=True,
                )


    # Tarjeta independiente de cancelados.
    st.markdown(
        """
        <style>
        div[data-testid="stButton"] button[kind="secondary"] {
            border-radius: 18px;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    cancel_key = f"show_cancelled_official_analytics_{year}"
    if cancel_key not in st.session_state:
        st.session_state[cancel_key] = False

    st.markdown("### Cancelados")
    st.caption("Haz clic en la tarjeta para ver la misma analítica, únicamente para los oficios cancelados.")

    cancel_col, _ = st.columns([1, 3])
    with cancel_col:
        if st.button(
            f"🟧  {len(cancelled_rows)}  Cancelados",
            key=f"btn_cancelled_officials_{year}",
            use_container_width=True,
        ):
            st.session_state[cancel_key] = not st.session_state[cancel_key]

    if st.session_state[cancel_key]:
        st.markdown("---")
        _render_analytics_block(
            cancelled_rows,
            f"Analítica de Oficios Cancelados · {year}",
            "Volumen mensual, destinatarios y origen de las solicitudes canceladas.",
            cancelled_mode=True,
        )

def official_letters_year(year: int):
    top1, top2 = st.columns([1, 5])
    top1.button(
        "← Años",
        use_container_width=True,
        key=f"official_back_years_{year}",
        on_click=_back_official_years,
    )
    top2.markdown(f"## Oficios Dirección General · {year}")

    client = client_with_token(st.session_state.access_token, st.session_state.refresh_token) if configured() else None
    rows = []
    if client:
        try:
            rows = (
                client.table("oficios_direccion_general")
                .select("*")
                .eq("anio", year)
                .execute()
                .data or []
            )
            rows = _dedupe_office_rows(rows)
        except Exception as exc:
            # Nunca vaciar silenciosamente la vista por una columna opcional.
            # Segundo intento con el núcleo histórico de la tabla.
            try:
                rows = (
                    client.table("oficios_direccion_general")
                    .select("id,anio,mes,numero_oficio,folio_control,destinatario,dependencia,solicitado_por,asunto,status_control,firma")
                    .eq("anio", year)
                    .execute()
                    .data or []
                )
                rows = _dedupe_office_rows(rows)
            except Exception:
                rows = []
                st.error(f"No fue posible cargar los oficios de {year}: {exc}")

    active_rows_year = _official_active_rows(rows)

    months_tab, analytics_tab = st.tabs(["Meses", "Analítica"])

    with months_tab:
        st.markdown(
            '<p class="choice-subtitle">Selecciona el mes que deseas consultar</p>',
            unsafe_allow_html=True,
        )
        counts = {month: 0 for month, _ in MONTHS_ES}
        for row in active_rows_year:
            try:
                month_value = int(row.get("mes") or 0)
            except (TypeError, ValueError):
                month_value = 0
            if month_value in counts:
                counts[month_value] += 1

        colors = ["var(--blue)", "var(--green)", "var(--teal)", "var(--purple)", "var(--orange)", "var(--gray)"]
        for start in range(0, 12, 3):
            columns = st.columns(3, gap="large")
            for offset, (month, month_name) in enumerate(MONTHS_ES[start:start + 3]):
                color = colors[(start + offset) % len(colors)]
                count = counts.get(month, 0)
                with columns[offset]:
                    st.markdown(
                        f'<div class="year-card" style="--accent:{color}">'
                        f'<h2 style="font-size:1.45rem">{month_name}</h2>'
                        f'<p>{count} oficio(s)</p></div>',
                        unsafe_allow_html=True,
                    )
                    st.button(
                        f"Abrir {month_name}",
                        key=f"official_month_{year}_{month}",
                        use_container_width=True,
                        type="primary",
                        on_click=_go_official_month,
                        args=(month,),
                    )

    with analytics_tab:
        if not configured():
            st.info("La analítica estará disponible al conectar Supabase.")
        elif not rows:
            st.info(f"No hay oficios registrados para {year}.")
        else:
            _official_letters_analytics(year, rows)


def official_letters():
    if not user_can(MODULE_OFFICIAL_LETTERS):
        st.error("No tienes permisos para acceder a Oficios Dirección General."); return
    selected_year = st.session_state.get("official_year")
    selected_month = st.session_state.get("official_month")
    if selected_year and selected_month:
        official_letters_month(int(selected_year), int(selected_month), dict(MONTHS_ES).get(int(selected_month), str(selected_month))); return
    if selected_year:
        official_letters_year(int(selected_year)); return
    st.markdown('<h1 class="choice-title">Oficios Dirección General</h1>', unsafe_allow_html=True)
    st.markdown('<p class="choice-subtitle">Archivo institucional de oficios · selecciona el año</p>', unsafe_allow_html=True)
    if configured():
        client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
        with st.container(border=True):
            st.markdown("### Ingestar archivo de control")
            st.caption("Por ahora se leerá exclusivamente la hoja ENVIADOS DG. Se crea un registro por cada oficio y el mes/año se obtiene del número de oficio.")
            control_file = st.file_uploader("Archivo de control (.xlsx)", type=["xlsx"], key="official_control_ingestion")
            if st.button("Ingestar ENVIADOS DG", type="primary", use_container_width=True, disabled=control_file is None, key="official_control_ingest_button"):
                try:
                    with st.spinner("Leyendo ENVIADOS DG y poblando la base de datos…"):
                        result, warnings = _ingest_dg_control(client, control_file)
                    st.success(f"Ingesta completada: {result.get('registros_detectados', 0)} oficio(s). Nuevos: {result.get('registros_nuevos', 0)} · actualizados: {result.get('registros_actualizados', 0)}.")
                    if warnings:
                        st.warning(f"Se omitieron {len(warnings)} fila(s) porque no se pudo obtener mes/año del número de oficio.")
                    st.rerun()
                except Exception as exc:
                    st.error(f"No fue posible ingestar el archivo: {exc}")
            latest = _latest_dg_ingestion(client)
            if latest:
                when = str(latest.get("completed_at") or latest.get("created_at") or "")[:16].replace("T", " · ")
                st.info(f"Última ingesta: **{latest.get('nombre_archivo') or 'Archivo'}** · {latest.get('registros_detectados') or 0} registro(s) · por **{latest.get('autor_nombre') or 'Usuario'}** · {when} · {latest.get('estado') or ''}")

            if is_master_admin():
                st.markdown(
                    """
                    <div style="background:#fff4e5;border:1px solid #ffb45c;border-radius:14px;padding:14px 16px;margin-top:12px;">
                    <b style="color:#d96900;">Administración maestra · Clasificación temática</b><br>
                    <span style="color:#59636b;">Clasifica o reclasifica los oficios existentes de 2025 y 2026 usando el catálogo temático institucional.</span>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )
                if st.button(
                    "🟧 Clasificar oficios existentes 2025–2026",
                    key="classify_existing_official_letters",
                    use_container_width=True,
                ):
                    try:
                        with st.spinner("Clasificando oficios existentes…"):
                            result = _classify_existing_official_letters(client, years=(2025, 2026))
                        st.success(f"Clasificación terminada: {result.get('actualizados', 0)} registro(s) actualizados.")
                        with st.expander("Distribución resultante"):
                            st.dataframe(
                                pd.DataFrame([
                                    {"Tema": key, "Registros": value}
                                    for key, value in sorted(
                                        (result.get("por_tema") or {}).items(),
                                        key=lambda item: item[1],
                                        reverse=True,
                                    )
                                ]),
                                use_container_width=True,
                                hide_index=True,
                            )
                        st.rerun()
                    except Exception as exc:
                        st.error(f"No fue posible clasificar los oficios existentes: {exc}")

        with st.container(border=True):
            st.markdown("### Importar vínculos de Google Drive")
            st.caption(
                "Los PDFs permanecen en Drive. Aquí sólo se vincula cada oficio con su URL, "
                "sin duplicar archivos en Supabase."
            )

            last_result = st.session_state.get("drive_links_last_result")
            if last_result:
                st.success(
                    f"Última importación: {last_result['detectados']} detectados · "
                    f"{last_result['vinculados']} vinculados · "
                    f"{last_result['ambiguos']} ambiguos · "
                    f"{last_result['sin_coincidencia']} sin coincidencia."
                )
                if last_result.get("warnings"):
                    with st.expander("Advertencias de la última importación"):
                        for warning in last_result["warnings"]:
                            st.write("•", warning)
                pending_review = [
                    row for row in last_result.get("details", [])
                    if row.get("estado") != "Vinculado"
                ]
                if pending_review:
                    with st.expander(f"Registros que requieren revisión · {len(pending_review)}"):
                        st.dataframe(
                            pd.DataFrame(pending_review),
                            use_container_width=True,
                            hide_index=True,
                        )
            link_file = st.file_uploader(
                "Sube Excel o CSV con folio_control, numero_oficio, nombre_archivo y url_drive",
                type=["xlsx", "xls", "csv"],
                key="drive_links_upload_dg",
            )
            if link_file:
                if st.button(
                    "Importar vínculos de Drive",
                    type="primary",
                    use_container_width=True,
                    key="import_drive_links_dg",
                ):
                    try:
                        with st.spinner("Cruzando vínculos con los oficios registrados…"):
                            result = _import_drive_links(client, link_file)
                        st.session_state["drive_links_last_result"] = result
                        st.session_state["drive_links_last_filename"] = getattr(link_file, "name", "Archivo")
                        st.rerun()
                    except Exception as exc:
                        st.error(f"No fue posible importar los vínculos: {exc}")

            latest_links = _latest_drive_link_import(client)
            if latest_links:
                when_links = str(latest_links.get("created_at") or "")[:16].replace("T", " · ")
                st.info(
                    f"Última importación de vínculos: **{latest_links.get('nombre_archivo') or 'Archivo'}** · "
                    f"{latest_links.get('registros_vinculados') or 0} vinculados · por "
                    f"**{latest_links.get('autor_nombre') or 'Usuario'}** · {when_links}"
                )
    else:
        st.info("La ingesta del archivo de control estará disponible al conectar Supabase.")
    year_counts = {year: 0 for year in OFFICIAL_LETTER_YEARS}
    if configured():
        for year in OFFICIAL_LETTER_YEARS:
            try:
                year_rows = (
                    client.table("oficios_direccion_general")
                    .select("id,anio,mes,numero_oficio,folio_control,status_control,firma")
                    .eq("anio", year)
                    .execute()
                    .data or []
                )
                year_rows = _dedupe_office_rows(year_rows)
                # Oficios enviados = oficios únicos no cancelados.
                year_rows = _official_active_rows(year_rows)
                year_counts[year] = len(year_rows)
            except Exception:
                year_counts[year] = 0

    colors = ["var(--blue)", "var(--green)", "var(--teal)", "var(--purple)", "var(--orange)", "var(--gray)", "var(--blue)"]
    for start in (0, 3, 6):
        batch = OFFICIAL_LETTER_YEARS[start:start + 3]
        columns = st.columns(len(batch), gap="large")
        for column, year, color in zip(columns, batch, colors[start:start + len(batch)]):
            with column:
                total_year = year_counts.get(year, 0)
                label = "oficio enviado" if total_year == 1 else "oficios enviados"
                st.markdown(
                    f'''<div class="year-card" style="--accent:{color}">
                        <h2>{year}</h2>
                        <p><b>{total_year}</b> {label}</p>
                    </div>''',
                    unsafe_allow_html=True,
                )
                st.button(
                    f"Abrir {year}",
                    key=f"official_year_{year}",
                    use_container_width=True,
                    type="primary",
                    on_click=_go_official_year,
                    args=(year,),
                )


def _apply_deep_link_from_query():
    """Abre directamente una sesión/acuerdo desde un vínculo de correo."""
    try:
        params = st.query_params
        modulo = str(params.get("modulo") or "").strip().lower()
        session_id = str(params.get("sesion_id") or "").strip()
        agreement_id = str(params.get("acuerdo_id") or "").strip()
        signature = f"{modulo}|{session_id}|{agreement_id}"

        if not modulo or not session_id:
            return
        if st.session_state.get("_deep_link_applied") == signature:
            return
        if not configured():
            return

        client = client_with_token(
            st.session_state.access_token,
            st.session_state.refresh_token,
        )

        if modulo == "junta" and user_can(MODULE_BOARD):
            rows = (
                client.table("sesiones_junta")
                .select("*")
                .eq("id", session_id)
                .limit(1)
                .execute()
                .data
                or []
            )
            if rows:
                session = rows[0]
                st.session_state.page = "Junta de Gobierno"
                st.session_state.board_year = int(session.get("anio"))
                st.session_state.board_session = session
                st.session_state.deep_link_agreement_id = agreement_id or None
                st.session_state["_deep_link_applied"] = signature

        elif modulo == "comite" and user_can(MODULE_COMMITTEES):
            rows = (
                client.table("sesiones_comite")
                .select("*")
                .eq("id", session_id)
                .limit(1)
                .execute()
                .data
                or []
            )
            if rows:
                session = rows[0]
                st.session_state.page = "Comités"
                st.session_state.committee_name = session.get("comite")
                st.session_state.committee_year = int(session.get("anio"))
                st.session_state.committee_session = session
                st.session_state.deep_link_agreement_id = agreement_id or None
                st.session_state["_deep_link_applied"] = signature
    except Exception:
        # Si un vínculo profundo falla, la app sigue funcionando normalmente.
        return

def board_government():
    if not user_can(MODULE_BOARD):
        st.error("No tienes permisos para acceder a Junta de Gobierno.")
        return
    if st.session_state.get("board_analytics"):
        board_analytics_dashboard()
    elif st.session_state.get("board_session"):
        board_session_detail(st.session_state.board_session)
    elif "board_year" in st.session_state:
        board_year_dashboard(int(st.session_state.board_year))
    else:
        board_year_selector()


def programs():
    if not user_can(MODULE_PROJECTS):
        st.error("No tienes permisos para acceder al módulo de Programas / Proyectos.")
        return
    direction = st.session_state.get("program_direction")
    action = st.session_state.get("program_action")

    if direction and not user_can_project_direction(direction):
        st.session_state.pop("program_direction", None)
        st.session_state.pop("program_action", None)
        st.error("No tienes permisos para consultar esa dirección.")
        return

    if not direction:
        st.markdown('<h1 class="choice-title">Programas / Proyectos</h1>', unsafe_allow_html=True)
        st.markdown('<p class="choice-subtitle">Selecciona la dirección responsable para continuar</p>', unsafe_allow_html=True)
        allowed = [item for item in PROJECT_DIRECTIONS if user_can_project_direction(item)]
        if not allowed:
            st.info("No tienes áreas de proyecto asignadas.")
            return
        columns = st.columns(len(allowed), gap="large")
        for column, allowed_direction in zip(columns, allowed):
            with column:
                is_operations = allowed_direction == "Dirección de Operaciones"
                st.markdown(f'''<div class="choice-card {'choice-operations' if is_operations else 'choice-projects'}"><div class="choice-icon">{'DO' if is_operations else 'DP'}</div>
                    <h3>{allowed_direction}</h3><p>Gestión de los programas y proyectos correspondientes a esta dirección.</p></div>''', unsafe_allow_html=True)
                if st.button(f"Ingresar a {'Operaciones' if is_operations else 'Proyectos'}", key=f"choose_{allowed_direction}", use_container_width=True, type="primary"):
                    st.session_state.program_direction = allowed_direction
                    st.session_state.pop("program_action", None)
                    st.rerun()
        return

    if not action:
        top1, top2 = st.columns([1, 5])
        if top1.button("← Direcciones", use_container_width=True):
            st.session_state.pop("program_direction", None)
            st.rerun()
        top2.markdown(f"### {direction}")
        st.markdown('<p class="choice-subtitle">Selecciona la acción que deseas realizar</p>', unsafe_allow_html=True)
        c1, c2, c3 = st.columns(3, gap="large")
        with c1:
            st.markdown('''<div class="choice-card choice-new"><div class="choice-icon">＋</div>
                <h3>Dar de alta nuevo proyecto</h3><p>Crear un expediente e incorporar su información general, documentos y seguimiento.</p></div>''', unsafe_allow_html=True)
            if st.button("Crear nuevo proyecto", key="choose_new", use_container_width=True, type="primary"):
                st.session_state.program_action = "new"
                st.session_state.objective_count = 1
                st.rerun()
        with c2:
            st.markdown('''<div class="choice-card choice-edit"><div class="choice-icon">✎</div>
                <h3>Editar proyecto</h3><p>Consultar un expediente existente para actualizar su información y seguimiento.</p></div>''', unsafe_allow_html=True)
            if st.button("Consultar y editar", key="choose_edit", use_container_width=True, type="primary"):
                st.session_state.program_action = "edit"
                st.rerun()
        with c3:
            st.markdown('''<div class="choice-card choice-view"><div class="choice-icon">◉</div>
                <h3>Visualizar proyectos</h3><p>Consultar los proyectos activos y acceder a su ficha general y avance.</p></div>''', unsafe_allow_html=True)
            if st.button("Ver proyectos activos", key="choose_view", use_container_width=True, type="primary"):
                st.session_state.program_action = "view"
                st.session_state.pop("view_project_id", None)
                st.rerun()
        return

    nav1, nav2 = st.columns([1, 5])
    if nav1.button("← Acciones", use_container_width=True):
        st.session_state.pop("program_action", None)
        for key in ["ficha_data", "ficha_photos", "ficha_pdf", "ficha_docx"]:
            st.session_state.pop(key, None)
        st.session_state.pop("view_project_id", None)
        st.session_state.pop("view_project_photos", None)
        st.rerun()
    nav2.markdown(f"### {direction}")

    if action == "new":
        project_form(direction)
    elif action == "view":
        try:
            view_active_projects(direction)
        except Exception as exc:
            st.error(f"No fue posible consultar los proyectos activos: {exc}")
    elif not configured():
        st.info("La consulta y edición estarán disponibles al conectar Supabase.")
    else:
        try:
            client = client_with_token(st.session_state.access_token, st.session_state.refresh_token)
            rows = client.table("proyectos").select("*").eq("direccion", direction).order("updated_at", desc=True).execute().data
            if not rows:
                st.info("Todavía no hay proyectos registrados en esta dirección.")
            else:
                labels = {f"{p['nombre']} — {p['municipio']} ({p['anio_inicio']})": p for p in rows}
                selected = st.selectbox("Selecciona el proyecto", labels)
                selected_project = labels[selected]
                master_delete_control(
                    "proyecto", str(selected_project["id"]), f"project_{selected_project['id']}",
                    lambda project_id=str(selected_project["id"]): delete_project_master(client, project_id),
                )
                project_form(direction, selected_project)
                if direction == "Dirección de Proyectos":
                    project_financial_documents(selected_project)
        except Exception as exc:
            st.error(f"No fue posible consultar los proyectos: {exc}")


def placeholder(title: str):
    st.title(title)
    st.info("Módulo preparado para desarrollarse en la siguiente etapa.")


if "user" not in st.session_state:
    login()
else:
    _apply_deep_link_from_query()
    with st.sidebar:
        st.markdown(brand_html(sidebar=True), unsafe_allow_html=True)
        if is_master_admin():
            st.success("Administrador maestro")
        if st.button("Inicio", use_container_width=True):
            st.session_state.page = "Inicio"
            st.rerun()
        if user_can(MODULE_PROJECTS):
            if st.button("Programas / Proyectos", use_container_width=True):
                st.session_state.page = "Programas / Proyectos"
                st.session_state.pop("program_direction", None)
                st.session_state.pop("program_action", None)
                st.rerun()
        if user_can(MODULE_BOARD):
            if st.button("Junta de Gobierno", use_container_width=True):
                st.session_state.page = "Junta de Gobierno"
                st.session_state.pop("board_year", None)
                st.session_state.pop("board_session", None)
                st.rerun()
        if user_can(MODULE_COMMITTEES):
            if st.button("Comités", use_container_width=True):
                st.session_state.page = "Comités"
                st.session_state.pop("committee_name", None)
                st.session_state.pop("committee_year", None)
                st.session_state.pop("committee_session", None)
                st.session_state.pop("committee_analytics", None)
                st.rerun()
        if user_can(MODULE_OFFICIAL_LETTERS):
            if st.button("Oficios Dirección General", use_container_width=True):
                st.session_state.page = "Oficios Dirección General"
                st.session_state.pop("official_year", None)
                st.session_state.pop("official_month", None)
                st.rerun()
        if st.session_state.user.get("rol") == "administrador":
            if st.button("Gestión de usuarios", use_container_width=True):
                st.session_state.page = "Gestión de usuarios"
                st.rerun()
        st.divider()
        if st.button("Cerrar sesión", use_container_width=True):
            st.session_state.clear()
            st.rerun()
    page = st.session_state.get("page", "Inicio")
    if page == "Inicio": landing()
    elif page == "Programas / Proyectos": programs()
    elif page == "Junta de Gobierno": board_government()
    elif page == "Gestión de usuarios": user_management()
    elif page == "Comités" and user_can(MODULE_COMMITTEES): committees()
    elif page == "Comités": st.error("No tienes permisos para acceder a Comités.")
    elif page == "Oficios Dirección General" and user_can(MODULE_OFFICIAL_LETTERS): official_letters()
    elif page == "Oficios Dirección General": st.error("No tienes permisos para acceder a Oficios Dirección General.")
    else: landing()
